# -*- coding: utf-8 -*-
"""Render each <section class="slide" id="sN"> in slides.html to PNG, then build PPTX.

Usage:
    python render_slides.py <slides_html_path> <output_pptx_path> [--total N]

Example:
    python render_slides.py .tmp/slides-html/slides.html output.pptx --total 10
"""
import asyncio
import argparse
from pathlib import Path
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("html", help="Path to slides HTML file")
    p.add_argument("output", help="Output PPTX file path")
    p.add_argument("--total", type=int, default=None, help="Number of slides (auto-detect if omitted)")
    return p.parse_args()

async def render(html_path: Path, out_dir: Path, total: int | None):
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(viewport={"width": 1920, "height": 1080}, device_scale_factor=1)
        page = await ctx.new_page()
        await page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
        await page.wait_for_timeout(1500)  # wait for fonts

        if total is None:
            # auto-detect: count #s1, #s2, ...
            total = 0
            while await page.query_selector(f"#s{total + 1}"):
                total += 1
            print(f"  auto-detected {total} slides")

        for i in range(1, total + 1):
            el = await page.query_selector(f"#s{i}")
            if not el:
                print(f"!! missing #s{i}")
                continue
            out = out_dir / f"slide-{i:02d}.png"
            await el.screenshot(path=str(out))
            print(f"  rendered {out.name}")
        await browser.close()
        return total

def build_pptx(out_dir: Path, total: int, pptx_out: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(1, total + 1):
        img = out_dir / f"slide-{i:02d}.png"
        if not img.exists():
            print(f"!! skip {img.name}")
            continue
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(pptx_out))
    print(f"\nSaved: {pptx_out}")

async def main():
    args = parse_args()
    html_path = Path(args.html)
    pptx_out = Path(args.output)
    out_dir = pptx_out.parent / "slides-png"
    out_dir.mkdir(exist_ok=True)

    total = await render(html_path, out_dir, args.total)
    build_pptx(out_dir, total, pptx_out)

if __name__ == "__main__":
    asyncio.run(main())
