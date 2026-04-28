"""
Build AI dev-platform comparison slides (McKinsey-style).
Compares: Claude Code Team / GitHub Copilot Enterprise / Kiro
Goal:    Which tool best realizes claude-code-components.pptx vision
Output:  .tmp/tool-comparison.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ==== Slide dimensions (16:9) ====
SLIDE_W = 9144000
SLIDE_H = 5143500

# ==== McKinsey-style palette ====
MCK_NAVY = RGBColor(26, 43, 76)
MCK_BLUE = RGBColor(74, 144, 226)
MCK_LIGHT_BLUE = RGBColor(208, 225, 245)
MCK_GRAY = RGBColor(138, 148, 166)
MCK_LIGHT_GRAY = RGBColor(230, 233, 239)
MCK_BG_LIGHT = RGBColor(246, 248, 251)
MCK_DARK = RGBColor(34, 34, 34)
MCK_RED = RGBColor(192, 48, 48)
MCK_RED_LIGHT = RGBColor(254, 243, 243)
MCK_GREEN = RGBColor(46, 125, 50)
MCK_GREEN_LIGHT = RGBColor(232, 245, 233)
MCK_ORANGE = RGBColor(214, 126, 30)
MCK_ORANGE_LIGHT = RGBColor(253, 243, 226)
MCK_PURPLE = RGBColor(103, 58, 183)
WHITE = RGBColor(255, 255, 255)

# ==== Layout ====
MARGIN_L = 457200
MARGIN_R = 457200
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
ACTION_TITLE_Y = 380000
ACTION_TITLE_H = 500000
RULE_Y = 950000
CONTENT_TOP = 1050000
CONTENT_BOTTOM = 4650000
FOOTER_Y = 4800000
TRACKER_Y = 260000


# ==== Helpers ====

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def add_rect(slide, left, top, width, height, fill=None, border=None, border_pt=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border is not None:
        shape.line.color.rgb = border
        if border_pt:
            shape.line.width = Pt(border_pt)
        else:
            shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=None, weight=0.75):
    if color is None:
        color = MCK_LIGHT_GRAY
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    return conn


def add_text(slide, left, top, width, height, text, size=10, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="Meiryo UI"):
    if color is None:
        color = MCK_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 30000
    tf.margin_right = 30000
    tf.margin_top = 20000
    tf.margin_bottom = 20000
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_multiline(slide, left, top, width, height, lines, size=10,
                  color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.25, font_name="Meiryo UI"):
    if color is None:
        color = MCK_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 30000
    tf.margin_right = 30000
    tf.margin_top = 20000
    tf.margin_bottom = 20000
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        if isinstance(line, tuple):
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            c = line[2] if len(line) > 2 and line[2] is not None else color
            s = line[3] if len(line) > 3 and line[3] is not None else size
        else:
            text = line
            bold = False
            c = color
            s = size
        if text == "":
            run = p.add_run()
            run.text = " "
            run.font.size = Pt(s)
            run.font.name = font_name
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(s)
        run.font.bold = bold
        run.font.color.rgb = c
        run.font.name = font_name
    return tb


def add_tracker(slide, label):
    add_text(slide, MARGIN_L, TRACKER_Y, CONTENT_W, 150000,
             label.upper(), size=8, bold=True, color=MCK_BLUE)


def add_action_title(slide, title, tracker=None):
    if tracker:
        add_tracker(slide, tracker)
    add_text(slide, MARGIN_L, ACTION_TITLE_Y, CONTENT_W, ACTION_TITLE_H,
             title, size=16, bold=True, color=MCK_NAVY)
    add_line(slide, MARGIN_L, RULE_Y, SLIDE_W - MARGIN_R, RULE_Y,
             color=MCK_NAVY, weight=1.5)


def add_footer(slide, page, total, source=None):
    if source:
        add_text(slide, MARGIN_L, FOOTER_Y, CONTENT_W - 500000, 200000,
                 f"Source: {source}", size=8, color=MCK_GRAY)
    add_text(slide, SLIDE_W - MARGIN_R - 800000, FOOTER_Y,
             800000, 200000,
             f"{page}", size=9, color=MCK_GRAY, align=PP_ALIGN.RIGHT)
    add_line(slide, MARGIN_L, FOOTER_Y - 30000,
             SLIDE_W - MARGIN_R, FOOTER_Y - 30000,
             color=MCK_LIGHT_GRAY, weight=0.5)


def col_header(slide, left, top, width, height, label, color=MCK_NAVY):
    add_rect(slide, left, top, width, height, fill=color)
    add_text(slide, left, top, width, height, label,
             size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def cell(slide, left, top, width, height, text, bold=False,
         color=MCK_DARK, fill=None, align=PP_ALIGN.LEFT, size=9):
    if fill is not None:
        add_rect(slide, left, top, width, height,
                 fill=fill, border=MCK_LIGHT_GRAY, border_pt=0.3)
    else:
        add_rect(slide, left, top, width, height,
                 border=MCK_LIGHT_GRAY, border_pt=0.3)
    add_text(slide, left + 50000, top, width - 100000, height,
             text, size=size, bold=bold, color=color,
             align=align, anchor=MSO_ANCHOR.MIDDLE)


def rating_cell(slide, left, top, width, height, rating):
    """rating: '◎' '○' '△' '×' — color-coded."""
    color_map = {
        '◎': (MCK_GREEN, MCK_GREEN_LIGHT),
        '○': (MCK_BLUE, MCK_LIGHT_BLUE),
        '△': (MCK_ORANGE, MCK_ORANGE_LIGHT),
        '×': (MCK_RED, MCK_RED_LIGHT),
    }
    fg, bg = color_map.get(rating, (MCK_DARK, MCK_BG_LIGHT))
    add_rect(slide, left, top, width, height,
             fill=bg, border=MCK_LIGHT_GRAY, border_pt=0.3)
    add_text(slide, left, top, width, height, rating,
             size=14, bold=True, color=fg,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================
# SLIDES
# ============================================

def s01_title(prs, page, total):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill=MCK_BLUE)
    add_text(slide, MARGIN_L, 1600000, CONTENT_W, 300000,
             "DX PORTAL TEAM  |  TOOLING DECISION", size=9, bold=True, color=MCK_BLUE)
    add_text(slide, MARGIN_L, 1950000, CONTENT_W, 700000,
             "集合知化を、どの道具で実現するか", size=30, bold=True, color=MCK_NAVY)
    add_text(slide, MARGIN_L, 2750000, CONTENT_W, 500000,
             "Claude Code Team / GitHub Copilot Enterprise / Kiro ― 3製品比較と推奨",
             size=13, color=MCK_DARK)
    add_text(slide, MARGIN_L, 3230000, CONTENT_W, 400000,
             "~ 5〜10名チームでの法人契約を想定 ~",
             size=11, color=MCK_GRAY)
    add_line(slide, MARGIN_L, 4550000, MARGIN_L + 1500000, 4550000,
             color=MCK_NAVY, weight=2)
    add_text(slide, MARGIN_L, 4600000, CONTENT_W, 250000,
             "意思決定資料 ― 予算承認ご依頼", size=10, color=MCK_GRAY)
    return slide


def s02_executive_summary(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "結論：Claude Code Team を推奨。L3（自動PR/自動修正）を Gitea 上で実現できる唯一の製品",
        tracker="Executive summary")

    # Three key reasons — horizontal cards
    card_top = 1150000
    card_h = 1350000
    card_w = (CONTENT_W - 400000) // 3
    gap = 200000

    cards = [
        ("1", "L3 を Gitea 上で実現できる",
         "CLI + MCP で Gitea に自動 PR 作成／ラベルで自動修正が可能。Copilot の Coding Agent は GitHub 専用で Gitea では動かない。",
         MCK_NAVY),
        ("2", "5要素すべてを満たす",
         "CLAUDE.md / Skills / Subagents / Hooks / MCP を単一製品で実装可能。他2製品は部分一致に留まる。",
         MCK_BLUE),
        ("3", "商用プランで学習利用なし",
         "Team プランは契約上モデル学習に使わない。SSO/監査は Enterprise で強化。",
         MCK_GREEN),
    ]
    for i, (num, head, body, color) in enumerate(cards):
        x = MARGIN_L + i * (card_w + gap)
        add_rect(slide, x, card_top, card_w, card_h,
                 fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
        add_rect(slide, x, card_top, 100000, card_h, fill=color)
        add_text(slide, x + 180000, card_top + 120000, card_w - 200000, 350000,
                 num, size=28, bold=True, color=color)
        add_text(slide, x + 180000, card_top + 530000, card_w - 300000, 300000,
                 head, size=12, bold=True, color=MCK_NAVY)
        add_text(slide, x + 180000, card_top + 830000, card_w - 300000, 480000,
                 body, size=9, color=MCK_DARK)

    # Bottom: ask
    ask_top = 2700000
    add_rect(slide, MARGIN_L, ask_top, CONTENT_W, 450000,
             fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, ask_top, CONTENT_W - 400000, 450000,
             "依頼事項：Claude Code Team 契約 10 席、初年度 約540千円 の承認をお願いしたい",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Rationale one-liner
    add_multiline(slide, MARGIN_L, 3250000, CONTENT_W, 1300000,
                  [
                      ("なぜ今か：", True, MCK_NAVY, 10),
                      ("・目標は L3（AIが自動で PR 作成／ラベルで自動修正）。L1（Chatでサマリー生成）止まりでは構想に到達しない。",
                       False, MCK_DARK, 10),
                      ("・Copilot の Coding Agent は GitHub 専用。Gitea 上で L3 を成立させるには Claude Code の CLI + MCP が必須。",
                       False, MCK_DARK, 10),
                      ("・Team プランの Projects 共有・使用量分析で、5〜10名規模での運用可視化と定着が可能。",
                       False, MCK_DARK, 10),
                      "",
                      ("リスク：", True, MCK_RED, 10),
                      ("・Team プランの SSO/SAML は Enterprise プランへ移行が必要。情シス要件を要確認（10名時点では先送り可）。",
                       False, MCK_DARK, 10),
                      ("・Agent Teams は実験機能。本番前提は避け、PoC で見極める。",
                       False, MCK_DARK, 10),
                  ], line_spacing=1.2)

    add_footer(slide, page, total)
    return slide


def s03_context(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "実現したい姿：目標は L3（AI が自動で PR を作成する）。5要素を『team-standards』に集約して秘書を育てる",
        tracker="Context")

    # Left: the 5 elements stack
    left_x = MARGIN_L
    left_w = 4400000
    top = CONTENT_TOP + 80000
    add_text(slide, left_x, top, left_w, 280000,
             "claude-code-components.pptx が示す構想", size=11, bold=True, color=MCK_NAVY)

    elements = [
        ("CLAUDE.md", MCK_NAVY,
         "秘書が毎回セッション冒頭で自動的に読み込む『常識・ルール集』",
         "例：『any/unknown禁止』『TDD先行』『金曜15分振り返り』を1行ずつ記載",
         "配置：team-standards/CLAUDE.md（200行以内／当番制で育てる）"),
        ("Skills", MCK_BLUE,
         "名前で呼び出せる手順書。2 回以上やる作業を markdown で定義",
         "例：/review-my-changes, /post-to-gitea, /sprint-plan",
         "配置：team-standards/skills/*.md（公式 + obra/superpowers を採用）"),
        ("Subagents", MCK_PURPLE,
         "特定分野に特化した秘書。ツール制限とモデル選択で『安全な専門家』を作る",
         "例：security-auditor（Write/Bash 禁止・Read/Grep のみ・Haiku で軽量）",
         "配置：team-standards/agents/*.md（VoltAgent から 6 つコピーで開始）"),
        ("Hooks", MCK_ORANGE,
         "ライフサイクル イベントで自動実行されるスクリプト。決定論的な制御装置",
         "例：PreToolUse で rm -rf 遮断／PostToolUse で lint 実行／UserPromptSubmit で CLAUDE.md 違反検出",
         "配置：.claude/settings.json + team-standards/hooks/"),
        ("MCP", MCK_GREEN,
         "外部ツールと繋げる標準プロトコル。Gitea / DB / ブラウザに秘書が直接触れる",
         "例：gitea-mcp で PR 作成・Issue 参照、postgres-mcp で DB 問合せ、playwright-mcp で E2E",
         "配置：.mcp.json + 個人 .env（公式 + 検証済みサーバーのみ採用）"),
    ]
    y = top + 340000
    row_h = 620000
    name_w = 1150000
    for name, color, what, example, place in elements:
        # colored bar + name label
        add_rect(slide, left_x, y, 120000, row_h - 40000, fill=color)
        add_text(slide, left_x + 180000, y, name_w - 200000, row_h - 40000,
                 name, size=11, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        # description: 3 lines
        add_multiline(slide, left_x + name_w, y + 20000,
                      left_w - name_w - 40000, row_h - 60000,
                      [
                          (what, True, MCK_DARK, 8),
                          ("例：" + example[2:] if example.startswith("例：") else example,
                           False, color, 8),
                          (place, False, MCK_GRAY, 7),
                      ], line_spacing=1.3, size=8)
        y += row_h

    # Right: the question
    right_x = MARGIN_L + left_w + 300000
    right_w = CONTENT_W - left_w - 300000

    add_rect(slide, right_x, top, right_w, 3100000,
             fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
    add_text(slide, right_x + 150000, top + 150000, right_w - 300000, 400000,
             "本資料が答える問い", size=11, bold=True, color=MCK_NAVY)

    add_multiline(slide,
                  right_x + 150000, top + 580000, right_w - 300000, 2400000,
                  [
                      ("Q1. どの製品で『5要素 × team-standards』を実装できるか", True, MCK_NAVY, 10),
                      ("    → スライド 9（実現度マトリクス）",
                       False, MCK_GRAY, 9),
                      "",
                      ("Q2. Gitea を残したまま使えるか", True, MCK_NAVY, 10),
                      ("    → スライド 8（連携比較）", False, MCK_GRAY, 9),
                      "",
                      ("Q3. セキュリティ要件を満たせるか", True, MCK_NAVY, 10),
                      ("    → スライド 10（統制比較）", False, MCK_GRAY, 9),
                      "",
                      ("Q4. コストは妥当か", True, MCK_NAVY, 10),
                      ("    → スライド 12（3年 TCO）", False, MCK_GRAY, 9),
                      "",
                      ("Q5. 何から始めるべきか", True, MCK_NAVY, 10),
                      ("    → スライド 15（ロードマップ）", False, MCK_GRAY, 9),
                  ], line_spacing=1.15, size=10)

    add_footer(slide, page, total,
               source="claude-code-components.pptx / architecture-overview.md")
    return slide


def s04_decision_frame(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "意思決定の3軸：『Gitea連携』『5要素の実現度』『セキュリティ統制』で優劣を確定させる",
        tracker="Decision frame")

    axes = [
        ("軸 1", "Gitea で L3 が成立するか",
         "自社 Gitea 上で『AI が自動で PR 作成／ラベルで自動修正』(L3) を実現できるか。L1 (Chatでサマリー生成) は全製品で可能、差は L3 で出る。",
         "Bot / Agent から Gitea API へ書き込む手段の有無で決まる",
         MCK_NAVY),
        ("軸 2", "5要素の実現度",
         "CLAUDE.md / Skills / Subagents / Hooks / MCP を製品機能として満たせるか。",
         "部分一致は『自作でカバー』になり運用が脆くなる",
         MCK_BLUE),
        ("軸 3", "セキュリティ統制",
         "学習利用・データ保持・SSO/SAML・監査ログの4項目で情シス要件を満たすか。",
         "商用プランかどうか、プラン階層で差がつく",
         MCK_PURPLE),
    ]

    top = CONTENT_TOP + 200000
    row_h = 1050000
    for i, (num, head, body, key, color) in enumerate(axes):
        y = top + i * row_h
        # Number block
        add_rect(slide, MARGIN_L, y, 900000, row_h - 100000, fill=color)
        add_text(slide, MARGIN_L, y, 900000, 400000,
                 num, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN_L, y + 380000, 900000, 550000,
                 head, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        bx = MARGIN_L + 1000000
        bw = CONTENT_W - 1000000
        add_rect(slide, bx, y, bw, row_h - 100000,
                 fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
        add_text(slide, bx + 150000, y + 80000, bw - 300000, 400000,
                 body, size=10, color=MCK_DARK)
        add_text(slide, bx + 150000, y + 500000, bw - 300000, 400000,
                 "→ " + key, size=10, bold=True, color=color)

    # Bottom note
    add_text(slide, MARGIN_L, 4400000, CONTENT_W, 250000,
             "コスト・学習コストはサブ軸。上記3軸で甲乙が付いた上で、最終確認として用いる。",
             size=9, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s05_claude_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Claude Code Team：CLI エージェント + チーム共有。5要素をネイティブに備える唯一の製品",
        tracker="Product 1 / 3")

    _product_overview(
        slide,
        product="Claude Code Team",
        vendor="Anthropic",
        price="$30 / 人・月（≒ 4,500円）",
        seats_10="450,000円 / 年（10名）",
        strengths=[
            "CLAUDE.md / Skills / Subagents / Hooks / MCP すべて公式サポート",
            "Gitea に MCP で直結可能（Git ホスト中立）",
            "Projects 共有・Usage Analytics で5〜10名運用に必要な可視化を具備",
            "商用プラン扱いで学習利用なし（契約で保証）",
        ],
        weaknesses=[
            "SSO/SAML は Enterprise プラン（$60+）への移行が必要",
            "Agent Teams は実験機能、本番前提は避ける",
            "国内法人サポートはまだ発展途上",
        ],
        verdict="第1選択肢 ― 5要素を満たす唯一の製品",
        verdict_color=MCK_GREEN)
    add_footer(slide, page, total, source="claude.com/pricing/team")
    return slide


def s06_copilot_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "GitHub Copilot Enterprise：IDE補完の最高峰、ただし GitHub 前提で Gitea とは噛み合わない",
        tracker="Product 2 / 3")

    _product_overview(
        slide,
        product="GitHub Copilot Enterprise",
        vendor="GitHub (Microsoft)",
        price="$39 / 人・月（≒ 4,000〜5,900円）",
        seats_10="約 470,000〜710,000円 / 年（10名・為替次第）",
        strengths=[
            "IDE補完は現時点で最も成熟・高速（L1 用途では最強）",
            "Chat で git diff を要約し PR 説明文を書くことは可能（L1）",
            "SSO/SAML・監査ログ・IP制限などエンタープライズ統制が標準装備",
            "国内 GitHub Japan 経由のサポート導線",
        ],
        weaknesses=[
            "Coding Agent（自律PR作成）は GitHub 専用 ― Gitea では L3 不可",
            "PR 自動サマリー UI / Code Review は GitHub のみ（L2 不可）",
            "Hooks 相当の仕組みが弱く、CLAUDE.md を『法律化』できない",
            "Subagents のようなツール制限付き専門エージェント機構なし",
        ],
        verdict="選外 ― L3 目標 × Gitea 継続 では成立しない（L1 のみなら Business で十分）",
        verdict_color=MCK_RED)
    add_footer(slide, page, total, source="github.com/pricing")
    return slide


def s07_kiro_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Kiro：Spec駆動開発に強み。ただし Gitea 対応とエコシステム成熟度に不安",
        tracker="Product 3 / 3")

    _product_overview(
        slide,
        product="Kiro",
        vendor="AWS",
        price="プレビュー期 無料〜低額（正式料金未確定）",
        seats_10="未確定 ― 本格採用時に再計算要",
        strengths=[
            "Spec→Design→Tasks→実装の型化が製品機能として組み込まれている",
            "MCP サポートにより gitea-mcp 経由で Gitea 連携（PR/Issue 自動化）が可能",
            "AWS Bedrock 経由でデータ境界を AWS 内に閉じやすい",
            "AWS 既存契約・IAM Identity Center に乗せやすい",
        ],
        weaknesses=[
            "公式 gitea-mcp の安定性・サポートは Claude Code より未成熟",
            "Hooks / Subagents 等の拡張点が Claude Code ほど整備されていない",
            "エコシステム（公式スキル・コミュニティ）がまだ薄い",
            "5〜10名規模のチーム運用事例が乏しい",
        ],
        verdict="次点 ― MCP で Gitea L3 は可。ただし Claude Code より工数大・成熟度低",
        verdict_color=MCK_ORANGE)
    add_footer(slide, page, total, source="aws.amazon.com/kiro")
    return slide


def _product_overview(slide, product, vendor, price, seats_10,
                       strengths, weaknesses, verdict, verdict_color):
    # Top band: product name and price
    top = CONTENT_TOP + 100000
    add_rect(slide, MARGIN_L, top, CONTENT_W, 600000,
             fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
    add_text(slide, MARGIN_L + 200000, top + 80000, 4000000, 200000,
             vendor.upper(), size=8, bold=True, color=MCK_GRAY)
    add_text(slide, MARGIN_L + 200000, top + 260000, 4000000, 350000,
             product, size=16, bold=True, color=MCK_NAVY)
    # Right price
    px = MARGIN_L + CONTENT_W - 3200000
    add_text(slide, px, top + 80000, 3000000, 200000,
             "LIST PRICE", size=8, bold=True, color=MCK_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, px, top + 260000, 3000000, 200000,
             price, size=11, bold=True, color=MCK_NAVY, align=PP_ALIGN.RIGHT)
    add_text(slide, px, top + 450000, 3000000, 180000,
             seats_10, size=9, color=MCK_DARK, align=PP_ALIGN.RIGHT)

    # Strengths / Weaknesses columns
    cy = top + 750000
    col_w = (CONTENT_W - 200000) // 2
    # Strengths
    add_rect(slide, MARGIN_L, cy, 100000, 350000, fill=MCK_GREEN)
    add_text(slide, MARGIN_L + 150000, cy, col_w - 150000, 350000,
             "強み", size=11, bold=True, color=MCK_GREEN, anchor=MSO_ANCHOR.MIDDLE)
    lines = [("・" + s, False, MCK_DARK, 9) for s in strengths]
    add_multiline(slide, MARGIN_L, cy + 400000, col_w, 1700000,
                  lines, line_spacing=1.3, size=9)

    # Weaknesses
    wx = MARGIN_L + col_w + 200000
    add_rect(slide, wx, cy, 100000, 350000, fill=MCK_RED)
    add_text(slide, wx + 150000, cy, col_w - 150000, 350000,
             "弱み・留意点", size=11, bold=True, color=MCK_RED, anchor=MSO_ANCHOR.MIDDLE)
    lines = [("・" + w, False, MCK_DARK, 9) for w in weaknesses]
    add_multiline(slide, wx, cy + 400000, col_w, 1700000,
                  lines, line_spacing=1.3, size=9)

    # Verdict
    vy = 4350000
    add_rect(slide, MARGIN_L, vy, CONTENT_W, 400000, fill=verdict_color)
    add_text(slide, MARGIN_L + 200000, vy, CONTENT_W - 400000, 400000,
             "判定：" + verdict, size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)


def s07b_claude_unique(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Claude Code ならではの機能（2026/4 時点）：L3 と『秘書を育てる』を支える固有の仕組み",
        tracker="Claude Code unique features")

    top = CONTENT_TOP + 150000
    col_w = (CONTENT_W - 300000) // 3
    row_h = 1500000
    gap_x = 150000
    gap_y = 150000

    features = [
        ("Subagents 並列実行",
         "最大 10 並列。frontmatter で tools / disallowedTools / model / permissionMode を明示し、『できる事・できない事』を制御できる専門秘書。",
         MCK_PURPLE),
        ("Hooks（決定論的制御）",
         "PreToolUse / PostToolUse / SessionStart 等のライフサイクルで任意スクリプトを発火。CLAUDE.md を『法律化』する唯一の仕組み。",
         MCK_ORANGE),
        ("MCP：クライアント兼サーバー",
         "Claude Code 自身が MCP クライアントかつサーバー。他製品は片側のみ。Gitea・DB・Playwright に直結しつつ、他ツールから呼ばれる側にもなれる。",
         MCK_GREEN),
        ("Git Worktree 隔離実行",
         "v2.0.60+。バックグラウンド Agent が独立した worktree で動くため、本流の作業と干渉しない。複数タスクの並走が安全。",
         MCK_BLUE),
        ("Checkpoints / Sandbox",
         "権限モード・チェックポイント・サンドボックス・信頼検証・managed settings による多層安全装置。危険操作の自動遮断に活用。",
         MCK_NAVY),
        ("マルチサーフェス + /loop",
         "CLI / デスクトップアプリ / VS Code・JetBrains / web / スマホ / 音声入力 / macOS Computer Use。/loop で定期タスクも自動実行。",
         MCK_PURPLE),
    ]

    for i, (name, desc, color) in enumerate(features):
        r, c = divmod(i, 3)
        x = MARGIN_L + c * (col_w + gap_x)
        y = top + r * (row_h + gap_y)
        # Card
        add_rect(slide, x, y, col_w, row_h,
                 fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
        # Accent bar
        add_rect(slide, x, y, col_w, 70000, fill=color)
        # Name
        add_text(slide, x + 150000, y + 130000, col_w - 300000, 400000,
                 name, size=12, bold=True, color=color)
        # Desc
        add_text(slide, x + 150000, y + 550000, col_w - 300000, row_h - 600000,
                 desc, size=9, color=MCK_DARK)

    # Bottom band — emphasis on exclusivity
    by = top + 2 * (row_h + gap_y) + 150000
    add_rect(slide, MARGIN_L, by, CONTENT_W, 350000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, by, CONTENT_W - 400000, 350000,
             "Subagents の並列実行・Hooks・MCP サーバー兼用・Worktree 隔離は、Copilot / Kiro には同等機構がない",
             size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page, total,
               source="code.claude.com/docs / release notes 2026-04")
    return slide


def s08_axis1_gitea(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "軸 1 ─ Gitea × L3 の実現コストで差がつく。Claude Code は MCP でネイティブ対応、Copilot は L3 不可",
        tracker="Decision axis 1")

    # Top: L1/L2/L3 ladder visualization
    ladder_y = CONTENT_TOP + 80000
    ladder_h = 900000
    steps = [
        ("L1", "Chat で説明文生成・手 PR",
         "人が差分を Chat に食わせ要約、手で PR 作成",
         MCK_LIGHT_GRAY, MCK_GRAY),
        ("L2", "PR 画面 UI 統合",
         "ボタン1つで要約生成（GitHub のみ）",
         MCK_LIGHT_GRAY, MCK_GRAY),
        ("L3", "Bot / Agent が自動 PR",
         "★ 今回の目標 ― ラベル貼付・Webhook で完全自動化",
         MCK_NAVY, WHITE),
    ]
    step_w = (CONTENT_W - 300000) // 3
    for i, (code, head, desc, bg, fg) in enumerate(steps):
        x = MARGIN_L + i * (step_w + 150000)
        add_rect(slide, x, ladder_y, step_w, ladder_h,
                 fill=bg, border=MCK_LIGHT_GRAY)
        add_text(slide, x + 150000, ladder_y + 100000, 600000, 350000,
                 code, size=22, bold=True, color=fg)
        add_text(slide, x + 750000, ladder_y + 150000, step_w - 900000, 300000,
                 head, size=11, bold=True, color=fg)
        add_text(slide, x + 150000, ladder_y + 500000, step_w - 300000, 380000,
                 desc, size=9, color=fg)
        if i < 2:
            # arrow between steps
            ax = x + step_w + 20000
            ay = ladder_y + ladder_h // 2
            add_line(slide, ax, ay, ax + 110000, ay,
                     color=MCK_NAVY, weight=1.5)

    # Matrix: one row per product, columns L1 / L3-auto-PR / L3-label-fix
    mx_top = ladder_y + ladder_h + 200000
    col0_w = 2400000
    col_w = (CONTENT_W - col0_w) // 3
    row_h = 600000

    col_header(slide, MARGIN_L, mx_top, col0_w, row_h, "製品 × Gitea 環境")
    col_header(slide, MARGIN_L + col0_w, mx_top, col_w, row_h, "L1 手PR", MCK_GRAY)
    col_header(slide, MARGIN_L + col0_w + col_w, mx_top, col_w, row_h,
               "L3 自動 PR 作成", MCK_NAVY)
    col_header(slide, MARGIN_L + col0_w + 2 * col_w, mx_top, col_w, row_h,
               "L3 ラベルで自動修正", MCK_NAVY)

    rows = [
        ("Claude Code Team", MCK_NAVY,
         "◎", "◎", "◎",
         "CLI + gitea-mcp で Gitea API に直結、Hooks で自動化"),
        ("Copilot Enterprise", MCK_GRAY,
         "◎", "×", "×",
         "Coding Agent は GitHub 専用／Gitea 上で L3 は成立しない"),
        ("Kiro", MCK_GRAY,
         "○", "○", "○",
         "MCP 経由で Gitea 接続可。公式サポート・安定性は Claude Code より低く初期構築工数が増える"),
    ]

    y = mx_top + row_h
    for name, color, r1, r2, r3, note in rows:
        # Product name + short note (merged cell)
        add_rect(slide, MARGIN_L, y, col0_w, row_h,
                 fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY, border_pt=0.3)
        add_text(slide, MARGIN_L + 120000, y + 50000, col0_w - 160000, 300000,
                 name, size=11, bold=True, color=color)
        add_text(slide, MARGIN_L + 120000, y + 330000, col0_w - 160000, row_h - 360000,
                 note, size=8, color=MCK_DARK)
        rating_cell(slide, MARGIN_L + col0_w, y, col_w, row_h, r1)
        rating_cell(slide, MARGIN_L + col0_w + col_w, y, col_w, row_h, r2)
        rating_cell(slide, MARGIN_L + col0_w + 2 * col_w, y, col_w, row_h, r3)
        y += row_h

    # Conclusion band
    cy = 4380000
    add_rect(slide, MARGIN_L, cy, CONTENT_W, 400000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, cy, CONTENT_W - 400000, 400000,
             "結論：Copilot は L3 × Gitea が不可。Kiro は MCP で○だが Claude Code より工数大・成熟度で劣る",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page, total)
    return slide


def s09_axis2_five_elements(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "軸 2 ─ 5要素の実現度：Claude Code Team のみ全要素をネイティブに満たす",
        tracker="Decision axis 2")

    # Matrix: rows = 5 elements, cols = Claude / Copilot / Kiro / "失うもの"
    table_top = CONTENT_TOP + 100000
    col0_w = 1200000
    lose_w = 2500000
    col_w = (CONTENT_W - col0_w - lose_w) // 3
    row_h = 560000

    # Headers
    col_header(slide, MARGIN_L, table_top, col0_w, row_h, "構成要素")
    col_header(slide, MARGIN_L + col0_w, table_top, col_w, row_h,
               "Claude Code Team", MCK_NAVY)
    col_header(slide, MARGIN_L + col0_w + col_w, table_top, col_w, row_h,
               "Copilot Enterprise", MCK_GRAY)
    col_header(slide, MARGIN_L + col0_w + 2 * col_w, table_top, col_w, row_h,
               "Kiro", MCK_GRAY)
    col_header(slide, MARGIN_L + col0_w + 3 * col_w, table_top, lose_w, row_h,
               "Copilot / Kiro で失うもの", MCK_RED)

    rows = [
        ("CLAUDE.md",
         "◎", "◎", "○",
         "共有の仕組みは代替可。要件としては差が小さい"),
        ("Skills",
         "◎", "△", "△",
         "『2回目から呼び出せる手順書』が持てず、運用が属人化する"),
        ("Subagents",
         "◎", "×", "△",
         "『security-auditor は Write 禁止』等のツール制限付き専門家が作れない"),
        ("Hooks",
         "◎", "×", "△",
         "CLAUDE.md の『法律化』が成立しない。危険操作遮断・自動lintが止まる"),
        ("MCP",
         "◎", "△", "○",
         "gitea-mcp / postgres-mcp / playwright-mcp の公式・安定性に差が出る"),
    ]

    y = table_top + row_h
    for name, r1, r2, r3, lose in rows:
        cell(slide, MARGIN_L, y, col0_w, row_h, name, bold=True,
             fill=MCK_BG_LIGHT, size=11)
        rating_cell(slide, MARGIN_L + col0_w, y, col_w, row_h, r1)
        rating_cell(slide, MARGIN_L + col0_w + col_w, y, col_w, row_h, r2)
        rating_cell(slide, MARGIN_L + col0_w + 2 * col_w, y, col_w, row_h, r3)
        cell(slide, MARGIN_L + col0_w + 3 * col_w, y, lose_w, row_h,
             lose, size=9, fill=MCK_RED_LIGHT, color=MCK_DARK)
        y += row_h

    # Takeaway
    ty = y + 150000
    add_rect(slide, MARGIN_L, ty, CONTENT_W, 400000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, ty, CONTENT_W - 400000, 400000,
             "結論：Subagents と Hooks の不在は致命的 ― CLAUDE.md の『法律化』が成立しなくなる",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page, total, source="各製品公式ドキュメント（2026-04時点）")
    return slide


def s10_axis3_security(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "軸 3 ─ セキュリティ統制：情シス要件の4項目で比較。いずれも商用プランで学習利用なし",
        tracker="Decision axis 3")

    table_top = CONTENT_TOP + 80000
    col0_w = 1600000
    need_w = 1500000
    col_w = (CONTENT_W - col0_w - need_w) // 3
    row_h = 600000

    col_header(slide, MARGIN_L, table_top, col0_w, row_h, "統制項目")
    col_header(slide, MARGIN_L + col0_w, table_top, need_w, row_h,
               "10名規模での必要性", MCK_ORANGE)
    col_header(slide, MARGIN_L + col0_w + need_w, table_top, col_w, row_h,
               "Claude Code Team", MCK_NAVY)
    col_header(slide, MARGIN_L + col0_w + need_w + col_w, table_top, col_w, row_h,
               "Copilot Enterprise", MCK_NAVY)
    col_header(slide, MARGIN_L + col0_w + need_w + 2 * col_w, table_top, col_w, row_h,
               "Kiro", MCK_NAVY)

    # rows: (name, need_level, need_note, r1, d1, r2, d2, r3, d3)
    rows = [
        ("学習利用",
         "必須", "情シスが最も気にする",
         "◎", "契約で利用なし",
         "◎", "契約で利用なし",
         "◎", "AWS 契約で利用なし"),
        ("データ保持",
         "必須", "DPA で確認可",
         "○", "一定期間あり／ZDR は Ent.",
         "○", "一定期間あり",
         "◎", "AWS内・Bedrock 基準"),
        ("SSO / SAML",
         "将来", "10名では未必須",
         "△", "Enterprise で対応",
         "◎", "標準装備",
         "◎", "IAM Identity Center"),
        ("監査ログ",
         "将来", "拡大時に必要",
         "△", "Enterprise で提供",
         "◎", "標準装備",
         "◎", "CloudTrail 統合"),
    ]

    y = table_top + row_h
    for name, need, need_note, r1, d1, r2, d2, r3, d3 in rows:
        cell(slide, MARGIN_L, y, col0_w, row_h, name, bold=True,
             fill=MCK_BG_LIGHT, size=11)
        # Need cell — color-coded
        need_color = MCK_RED if need == "必須" else MCK_ORANGE
        need_bg = MCK_RED_LIGHT if need == "必須" else MCK_ORANGE_LIGHT
        add_rect(slide, MARGIN_L + col0_w, y, need_w, row_h,
                 fill=need_bg, border=MCK_LIGHT_GRAY, border_pt=0.3)
        add_text(slide, MARGIN_L + col0_w + 100000, y + 50000,
                 need_w - 200000, 280000,
                 need, size=12, bold=True, color=need_color)
        add_text(slide, MARGIN_L + col0_w + 100000, y + 310000,
                 need_w - 200000, row_h - 340000,
                 need_note, size=8, color=MCK_DARK)

        base = MARGIN_L + col0_w + need_w
        rating_cell(slide, base, y, 380000, row_h, r1)
        cell(slide, base + 380000, y, col_w - 380000, row_h, d1, size=8)
        rating_cell(slide, base + col_w, y, 380000, row_h, r2)
        cell(slide, base + col_w + 380000, y, col_w - 380000, row_h, d2, size=8)
        rating_cell(slide, base + 2 * col_w, y, 380000, row_h, r3)
        cell(slide, base + 2 * col_w + 380000, y, col_w - 380000, row_h, d3, size=8)
        y += row_h

    # Conclusion
    cy = y + 150000
    add_rect(slide, MARGIN_L, cy, CONTENT_W, 400000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, cy, CONTENT_W - 400000, 400000,
             "結論：10名時点では『必須』2項目を全社満たす。『将来』項目は拡大時に Enterprise へ移行で対応可",
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page, total, source="各ベンダー Trust Center (2026-04時点)")
    return slide


def s11_overall_matrix(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "総合比較マトリクス：3軸 + サブ軸を合算すると Claude Code Team が明確に優位",
        tracker="Overall matrix")

    table_top = CONTENT_TOP + 50000
    weight_w = 700000
    col0_w = 2200000
    col_w = (CONTENT_W - col0_w - weight_w) // 3
    main_h = 470000
    sub_h = 340000

    col_header(slide, MARGIN_L, table_top, col0_w, main_h, "評価項目")
    col_header(slide, MARGIN_L + col0_w, table_top, weight_w, main_h, "重み", MCK_ORANGE)
    col_header(slide, MARGIN_L + col0_w + weight_w, table_top, col_w, main_h,
               "Claude Code Team", MCK_NAVY)
    col_header(slide, MARGIN_L + col0_w + weight_w + col_w, table_top, col_w, main_h,
               "Copilot Enterprise", MCK_GRAY)
    col_header(slide, MARGIN_L + col0_w + weight_w + 2 * col_w, table_top, col_w, main_h,
               "Kiro", MCK_GRAY)

    # rows: (name, weight, is_main, r1, r2, r3)
    rows = [
        ("Gitea で L3 が成立",                "×3", True,  "◎", "×", "○"),
        ("5要素の実現度",                     "×3", True,  "◎", "△", "△"),
        ("セキュリティ統制(10名規模)",         "×2", True,  "○", "◎", "◎"),
        ("コスト (10名/年)",                  "×1", False, "○", "○", "◎"),
        ("学習・運用コスト",                  "×1", False, "○", "◎", "△"),
        ("ベンダーロックイン",                "×1", False, "◎", "△", "△"),
        ("国内サポート",                      "×1", False, "△", "◎", "○"),
        ("将来拡張性 (秘書構想)",             "×1", False, "◎", "△", "○"),
    ]

    score_map = {"◎": 4, "○": 3, "△": 2, "×": 1}
    totals = [0, 0, 0]
    max_total = 0

    y = table_top + main_h
    for name, weight, is_main, r1, r2, r3 in rows:
        rh = main_h if is_main else sub_h
        bg = MCK_LIGHT_BLUE if is_main else MCK_BG_LIGHT
        label_size = 11 if is_main else 9
        # name cell
        cell(slide, MARGIN_L, y, col0_w, rh, name, bold=True,
             fill=bg, size=label_size)
        # weight cell
        w_fg = MCK_RED if is_main else MCK_GRAY
        cell(slide, MARGIN_L + col0_w, y, weight_w, rh, weight,
             bold=True, fill=bg, color=w_fg, size=12, align=PP_ALIGN.CENTER)
        base = MARGIN_L + col0_w + weight_w
        rating_cell(slide, base, y, col_w, rh, r1)
        rating_cell(slide, base + col_w, y, col_w, rh, r2)
        rating_cell(slide, base + 2 * col_w, y, col_w, rh, r3)

        w = int(weight[1:])
        totals[0] += score_map[r1] * w
        totals[1] += score_map[r2] * w
        totals[2] += score_map[r3] * w
        max_total += 4 * w
        y += rh

    # Weighted totals
    score_y = y + 100000
    cell(slide, MARGIN_L, score_y, col0_w + weight_w, 450000,
         f"加重合計（満点 {max_total}）",
         bold=True, fill=MCK_NAVY, color=WHITE, size=11)
    base = MARGIN_L + col0_w + weight_w
    labels = [
        (totals[0], MCK_GREEN, MCK_GREEN_LIGHT),
        (totals[1], MCK_ORANGE, MCK_ORANGE_LIGHT),
        (totals[2], MCK_ORANGE, MCK_ORANGE_LIGHT),
    ]
    for i, (val, fg, bg) in enumerate(labels):
        cell(slide, base + i * col_w, score_y, col_w, 450000,
             f"{val} / {max_total}",
             bold=True, fill=bg, color=fg, size=14, align=PP_ALIGN.CENTER)

    add_text(slide, MARGIN_L, score_y + 480000, CONTENT_W, 200000,
             "主軸（青背景）= 重み×3／×2、副軸（グレー）= 重み×1。◎=4 ○=3 △=2 ×=1 点",
             size=8, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s12_cost(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "コスト比較：10名・3年TCO でも Claude Code Team は許容範囲内。投資対効果で正当化可能",
        tracker="Cost / TCO")

    # Table: rows = cost items, cols = 3 products
    table_top = CONTENT_TOP + 100000
    col0_w = 2400000
    col_w = (CONTENT_W - col0_w) // 3
    row_h = 420000

    col_header(slide, MARGIN_L, table_top, col0_w, row_h, "費目")
    col_header(slide, MARGIN_L + col0_w, table_top, col_w, row_h, "Claude Code Team", MCK_NAVY)
    col_header(slide, MARGIN_L + col0_w + col_w, table_top, col_w, row_h, "Copilot Enterprise", MCK_NAVY)
    col_header(slide, MARGIN_L + col0_w + 2 * col_w, table_top, col_w, row_h, "Kiro", MCK_NAVY)

    rows = [
        ("ライセンス (円/人・月)",       "4,500",        "5,900 (≒$39)",   "未確定(暫定 3,000)"),
        ("年額 (10名・円)",              "540,000",      "708,000",        "360,000 (暫定)"),
        ("3年 TCO (10名・円)",           "1,620,000",    "2,124,000",      "1,080,000 (暫定)"),
        ("初期構築工数 (人日)",          "3 日",         "5 日",           "15〜20 日"),
        ("Gitea連携追加工数 (人日)",     "1 日",         "(不可)",         "3〜5 日（MCP経由）"),
    ]
    y = table_top + row_h
    for i, (name, v1, v2, v3) in enumerate(rows):
        bg = MCK_BG_LIGHT if i % 2 == 0 else None
        cell(slide, MARGIN_L, y, col0_w, row_h, name, bold=True, fill=bg, size=10)
        cell(slide, MARGIN_L + col0_w, y, col_w, row_h, v1, fill=bg, size=10, align=PP_ALIGN.CENTER)
        cell(slide, MARGIN_L + col0_w + col_w, y, col_w, row_h, v2, fill=bg, size=10, align=PP_ALIGN.CENTER)
        cell(slide, MARGIN_L + col0_w + 2 * col_w, y, col_w, row_h, v3, fill=bg, size=10, align=PP_ALIGN.CENTER)
        y += row_h

    # Verdict
    vy = y + 200000
    add_rect(slide, MARGIN_L, vy, CONTENT_W, 400000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, vy, CONTENT_W - 400000, 400000,
             "所見：Claude Code Team は Copilot Enterprise より年間 約17万円 安い。Kiro は工数込みなら同等以上に高くなる",
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page, total,
               source="各ベンダー価格ページ(2026-04) / $1=150円換算")
    return slide


def s13_risk(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "リスクと対策：Claude Code Team 採用時の3つの懸念はいずれも軽減策が確立している",
        tracker="Risk register")

    table_top = CONTENT_TOP + 150000
    col_w = [2400000, 2800000, 3000000]
    row_h = 700000

    # Headers
    x = MARGIN_L
    col_header(slide, x, table_top, col_w[0], 420000, "リスク")
    col_header(slide, x + col_w[0], table_top, col_w[1], 420000, "影響")
    col_header(slide, x + col_w[0] + col_w[1], table_top, col_w[2], 420000, "緩和策")

    risks = [
        ("SSO / SAML が Team プランに無い",
         "情シスが10名超の拡大時に統制NGを出す可能性",
         "10名以内で運用 → 拡大時に Enterprise($60+)へ移行。契約条項で前提を明記"),
        ("Agent Teams が実験機能",
         "仕様変更・提供停止の可能性あり",
         "PoC 運用に留め、本番フローからは切り離す。コアは Skills + Subagents で構成"),
        ("MCP サーバーの脆弱性",
         "悪意のあるサーバー経由で社内情報漏洩の可能性",
         "公式 + 検証済みのみ採用。.mcp.json をリポ管理＋Hooks で未許可サーバーを遮断"),
    ]

    y = table_top + 420000
    for i, (risk, impact, mitig) in enumerate(risks):
        bg = MCK_BG_LIGHT if i % 2 == 0 else None
        cell(slide, x, y, col_w[0], row_h, risk, bold=True, fill=bg, size=10)
        cell(slide, x + col_w[0], y, col_w[1], row_h, impact, fill=bg, size=9)
        cell(slide, x + col_w[0] + col_w[1], y, col_w[2], row_h, mitig, fill=bg, size=9, color=MCK_GREEN)
        y += row_h

    # Note
    ny = y + 150000
    add_text(slide, MARGIN_L, ny, CONTENT_W, 300000,
             "上記3件はいずれも運用ルールで制御可能。『致命的に採用を阻む』リスクは現時点で確認されていない。",
             size=10, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s14_recommended_architecture(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "推奨構成：Claude Code Team + team-standards リポ + Gitea MCP で5要素を統合運用",
        tracker="Recommended architecture")

    # Central stack diagram
    top = CONTENT_TOP + 150000

    # Layer 1: Users
    l1_y = top
    add_rect(slide, MARGIN_L, l1_y, CONTENT_W, 400000,
             fill=MCK_LIGHT_BLUE, border=MCK_BLUE)
    add_text(slide, MARGIN_L + 200000, l1_y, CONTENT_W - 400000, 400000,
             "チームメンバー（5〜10名）  ―  全員が Claude Code CLI + IDE で同じ秘書を呼び出す",
             size=10, bold=True, color=MCK_NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # Layer 2: Claude Code Team
    l2_y = l1_y + 500000
    add_rect(slide, MARGIN_L, l2_y, CONTENT_W, 500000,
             fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, l2_y, CONTENT_W - 400000, 500000,
             "Claude Code Team プラン  ―  Projects 共有 / Usage Analytics / チームメモリ",
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Layer 3: team-standards repo (5 elements)
    l3_y = l2_y + 600000
    col_w = (CONTENT_W - 400000) // 5
    blocks = [
        ("CLAUDE.md", MCK_NAVY, "ルール"),
        ("Skills", MCK_BLUE, "手順書"),
        ("Subagents", MCK_PURPLE, "専門秘書"),
        ("Hooks", MCK_ORANGE, "法律化"),
        ("MCP", MCK_GREEN, "外部連携"),
    ]
    add_text(slide, MARGIN_L, l3_y, CONTENT_W, 250000,
             "team-standards リポジトリ（単一のソースオブトゥルース）",
             size=9, bold=True, color=MCK_GRAY)
    bh = 700000
    for i, (name, color, desc) in enumerate(blocks):
        x = MARGIN_L + i * (col_w + 100000)
        add_rect(slide, x, l3_y + 280000, col_w, bh, fill=color)
        add_text(slide, x, l3_y + 280000, col_w, 350000,
                 name, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, l3_y + 610000, col_w, 350000,
                 desc, size=9, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Layer 4: external systems
    l4_y = l3_y + 280000 + bh + 150000
    ext_w = (CONTENT_W - 400000) // 3
    externals = [
        ("Gitea (自社運用)", "PR / Issue / Branch", MCK_GREEN),
        ("社内 DB / ドキュメント", "PostgreSQL / Markdown", MCK_BLUE),
        ("CI / テスト基盤", "Playwright / lint / test", MCK_PURPLE),
    ]
    for i, (name, desc, color) in enumerate(externals):
        x = MARGIN_L + i * (ext_w + 200000)
        add_rect(slide, x, l4_y, ext_w, 400000,
                 fill=WHITE, border=color, border_pt=1.5)
        add_text(slide, x, l4_y, ext_w, 180000,
                 name, size=10, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, l4_y + 180000, ext_w, 200000,
                 desc, size=8, color=MCK_GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Note
    add_text(slide, MARGIN_L, l4_y + 450000, CONTENT_W, 200000,
             "MCP 経由で Gitea・DB・CI に接続。Hooks で危険操作を遮断、Subagents でレビューを自動化。",
             size=9, color=MCK_GRAY, align=PP_ALIGN.CENTER)

    add_footer(slide, page, total)
    return slide


def s15_roadmap(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "導入ロードマップ：Day1 から3ヶ月で『秘書を育てる』基盤を立ち上げる",
        tracker="Rollout roadmap")

    # 4 phases
    top = CONTENT_TOP + 200000
    col_w = (CONTENT_W - 300000) // 4
    col_h = 2800000

    phases = [
        ("Day 1", "契約 + 最小セット",
         [
             "Team 10席契約",
             "team-standards リポ作成",
             "CLAUDE.md v0.1 (30行)",
             "Hooks: save-transcript",
             "Hooks: bash-guard",
         ], MCK_NAVY),
        ("Week 1", "MCP で Gitea 接続",
         [
             "gitea-mcp 導入",
             "PR / Issue から秘書が参照",
             "Day 1 Skills 3つ投入",
             "メンバー全員オンボード",
         ], MCK_BLUE),
        ("Month 1", "Subagents 投入",
         [
             "VoltAgent から6つ移植",
             "architect-reviewer / security-auditor 等",
             "自動 lint / review フロー構築",
             "金曜15分の振り返り運用",
         ], MCK_PURPLE),
        ("Month 3", "横展開 / 評価",
         [
             "Usage Analytics で定着度測定",
             "CLAUDE.md 100行超",
             "SSO/SAML 要否を情シスと合意",
             "Enterprise 移行 / 増員を判断",
         ], MCK_GREEN),
    ]

    for i, (phase, head, items, color) in enumerate(phases):
        x = MARGIN_L + i * (col_w + 100000)
        # Header block
        add_rect(slide, x, top, col_w, 500000, fill=color)
        add_text(slide, x, top, col_w, 220000,
                 phase, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, top + 220000, col_w, 260000,
                 head, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_rect(slide, x, top + 500000, col_w, col_h - 500000,
                 fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
        lines = [("・" + it, False, MCK_DARK, 9) for it in items]
        add_multiline(slide, x + 100000, top + 580000,
                      col_w - 200000, col_h - 600000,
                      lines, line_spacing=1.4, size=9)

    # Bottom arrow
    ay = top + col_h + 150000
    add_line(slide, MARGIN_L, ay, SLIDE_W - MARGIN_R, ay,
             color=MCK_NAVY, weight=2)
    add_text(slide, MARGIN_L, ay + 30000, CONTENT_W, 200000,
             "判断チェックポイント： Month 3 時点で『定着度』『5要素カバレッジ』『情シス要件』の3つをレビューし、継続 / 拡大 / 乗換を決定",
             size=9, bold=True, color=MCK_NAVY, align=PP_ALIGN.CENTER)

    add_footer(slide, page, total)
    return slide


def s16_next_actions(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "決裁ポイントと次のアクション：ご承認いただきたい3点とこちらで進める5点",
        tracker="Decision points")

    # Left: decision needs (ご承認)
    lx = MARGIN_L
    lw = (CONTENT_W - 200000) // 2
    top = CONTENT_TOP + 150000

    add_rect(slide, lx, top, 100000, 400000, fill=MCK_NAVY)
    add_text(slide, lx + 150000, top, lw - 150000, 400000,
             "ご承認いただきたい", size=12, bold=True, color=MCK_NAVY, anchor=MSO_ANCHOR.MIDDLE)

    approval_items = [
        ("1.", "契約締結の可否", "Claude Code Team 10席、初年度 約540千円"),
        ("2.", "対象チームの確定", "先行10名の選定／拡大条件の合意"),
        ("3.", "情シス合意の進め方", "DPA 入手と Trust Center 確認の担当明確化"),
    ]
    y = top + 500000
    for num, head, body in approval_items:
        add_rect(slide, lx + 100000, y, lw - 100000, 580000,
                 fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
        add_text(slide, lx + 180000, y + 50000, 400000, 250000,
                 num, size=16, bold=True, color=MCK_NAVY)
        add_text(slide, lx + 600000, y + 50000, lw - 700000, 250000,
                 head, size=11, bold=True, color=MCK_DARK)
        add_text(slide, lx + 600000, y + 310000, lw - 700000, 250000,
                 body, size=9, color=MCK_GRAY)
        y += 700000

    # Right: our side actions
    rx = MARGIN_L + lw + 200000
    rw = lw
    add_rect(slide, rx, top, 100000, 400000, fill=MCK_GREEN)
    add_text(slide, rx + 150000, top, rw - 150000, 400000,
             "こちらで進めます", size=12, bold=True, color=MCK_GREEN, anchor=MSO_ANCHOR.MIDDLE)

    actions = [
        ("契約後 3 日", "team-standards リポ初期化 + CLAUDE.md v0.1"),
        ("契約後 1 週", "Gitea MCP 接続、10名オンボード"),
        ("契約後 1 か月", "Subagents 6つ投入、金曜15分振り返り開始"),
        ("契約後 3 か月", "Usage Analytics レビュー、継続可否判断"),
        ("継続中", "隔週でリーダーへ進捗共有、半期で成果報告"),
    ]
    y = top + 500000
    for i, (when, what) in enumerate(actions):
        bg = MCK_BG_LIGHT if i % 2 == 0 else None
        if bg:
            add_rect(slide, rx + 100000, y, rw - 100000, 440000,
                     fill=bg, border=MCK_LIGHT_GRAY, border_pt=0.3)
        else:
            add_rect(slide, rx + 100000, y, rw - 100000, 440000,
                     border=MCK_LIGHT_GRAY, border_pt=0.3)
        add_text(slide, rx + 180000, y, 1500000, 440000,
                 when, size=10, bold=True, color=MCK_GREEN, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, rx + 1700000, y, rw - 1800000, 440000,
                 what, size=9, color=MCK_DARK, anchor=MSO_ANCHOR.MIDDLE)
        y += 440000

    # Closing band
    cy = 4400000
    add_rect(slide, MARGIN_L, cy, CONTENT_W, 280000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, cy, CONTENT_W - 400000, 280000,
             "ご判断のほど、よろしくお願いいたします。",
             size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page, total)
    return slide


# ==========================================
# MAIN
# ==========================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        s01_title,
        s02_executive_summary,
        s03_context,
        s04_decision_frame,
        s05_claude_overview,
        s06_copilot_overview,
        s07_kiro_overview,
        s07b_claude_unique,
        s08_axis1_gitea,
        s09_axis2_five_elements,
        s10_axis3_security,
        s11_overall_matrix,
        s12_cost,
        s13_risk,
        s14_recommended_architecture,
        s15_roadmap,
        s16_next_actions,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        fn(prs, i, total)

    out_path = Path(__file__).parent / "tool-comparison.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")
    print(f"Total: {total} slides")


if __name__ == "__main__":
    build()
