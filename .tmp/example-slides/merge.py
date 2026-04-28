"""
Merge example slides into the original presentation.
Insert positions (0-indexed, after insertion adjustments applied sequentially):
  - ex1 after slide 3 (困っていること detail) -> position 3
  - ex2 after slide 5 (解決策) -> position 6 (original 5 + 1 inserted)
  - ex3 after slide 6 (秘書がいると開発が変わる) -> position 8 (original 6 + 2 inserted)
  - ex4 after slide 8 (3層構造) -> position 11 (original 8 + 3 inserted)
  - ex5 after slide 11 (組織の働き方も変わる) -> position 15 (original 11 + 4 inserted)
"""
import copy
import sys
from pathlib import Path
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

def copy_slide(src_prs, src_idx, dst_prs, insert_idx):
    """Copy a slide from src_prs at src_idx into dst_prs at insert_idx."""
    src_slide = src_prs.slides[src_idx]

    # Get the slide layout from destination (use blank or first layout)
    slide_layout = dst_prs.slide_layouts[0]  # Only layout available

    # Add a new slide
    new_slide = dst_prs.slides.add_slide(slide_layout)

    # Clear the new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy all elements from source slide to new slide
    for elem in src_slide._element:
        new_slide._element.append(copy.deepcopy(elem))

    # Copy images and other relationships
    for rel in src_slide.part.rels.values():
        if "image" in rel.reltype:
            # Copy image
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    # Move slide to correct position
    # The new slide is added at the end, we need to move it
    slide_list = dst_prs.slides._sldIdLst
    slide_ids = list(slide_list)
    # The new slide is the last one
    new_slide_elem = slide_ids[-1]
    slide_list.remove(new_slide_elem)
    # Insert at the correct position
    if insert_idx >= len(slide_ids):
        slide_list.append(new_slide_elem)
    else:
        slide_list.insert(insert_idx, new_slide_elem)

    return new_slide

def main():
    base = Path(__file__).resolve().parent.parent.parent
    original_path = base / ".tmp" / "ai-collaborative-dev-platform.pptx"
    examples_path = base / ".tmp" / "example-slides" / "examples.pptx"
    output_path = base / ".tmp" / "ai-collaborative-dev-platform-v2.pptx"

    original = Presentation(str(original_path))
    examples = Presentation(str(examples_path))

    # Insert positions: (example_slide_index, insert_after_original_slide)
    # Original slides are 0-indexed: slide1=0, slide2=1, ...
    # After slide 3 (index 2) = insert at index 3
    # After slide 5 (index 4) = insert at index 5 + 1 already inserted = 6
    # After slide 6 (index 5) = insert at index 6 + 2 already inserted = 8
    # After slide 8 (index 7) = insert at index 8 + 3 already inserted = 11
    # After slide 11 (index 10) = insert at index 11 + 4 already inserted = 15

    insertions = [
        (0, 3),   # ex1 after "困っていること" (slide 3)
        (1, 6),   # ex2 after "解決策" (slide 5, +1)
        (2, 8),   # ex3 after "秘書がいると開発が変わる" (slide 6, +2)
        (3, 11),  # ex4 after "3層構造" (slide 8, +3)
        (4, 15),  # ex5 after "組織の働き方も変わる" (slide 11, +4)
    ]

    for ex_idx, insert_pos in insertions:
        print(f"Inserting example slide {ex_idx+1} at position {insert_pos}")
        copy_slide(examples, ex_idx, original, insert_pos)

    original.save(str(output_path))
    print(f"Saved: {output_path}")
    print(f"Total slides: {len(original.slides)}")

if __name__ == "__main__":
    main()
