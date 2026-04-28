"""
Build team explanation slides - McKinsey-style.
Output: .tmp/team-collective-intelligence-v2.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ==== Slide dimensions (16:9) ====
SLIDE_W = 9144000
SLIDE_H = 5143500

# ==== McKinsey-style palette ====
MCK_NAVY = RGBColor(26, 43, 76)          # #1a2b4c - primary
MCK_BLUE = RGBColor(74, 144, 226)         # #4a90e2 - accent
MCK_LIGHT_BLUE = RGBColor(208, 225, 245)  # #d0e1f5 - light accent
MCK_GRAY = RGBColor(138, 148, 166)        # #8a94a6 - secondary text
MCK_LIGHT_GRAY = RGBColor(230, 233, 239)  # #e6e9ef - dividers
MCK_BG_LIGHT = RGBColor(246, 248, 251)    # #f6f8fb - light bg fills
MCK_DARK = RGBColor(34, 34, 34)           # #222 - body text
MCK_RED = RGBColor(192, 48, 48)           # warning
MCK_GREEN = RGBColor(46, 125, 50)         # positive
MCK_ORANGE = RGBColor(214, 126, 30)       # attention
WHITE = RGBColor(255, 255, 255)

# ==== Layout constants ====
MARGIN_L = 457200   # 0.5"
MARGIN_R = 457200
MARGIN_T = 380000   # small top margin
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
ACTION_TITLE_Y = 380000
ACTION_TITLE_H = 500000
RULE_Y = 950000     # horizontal rule under title
CONTENT_TOP = 1050000
CONTENT_BOTTOM = 4650000   # leave ~500K for footer
FOOTER_Y = 4800000

TRACKER_LEFT = MARGIN_L     # small uppercase breadcrumb label
TRACKER_Y = 260000


# ==== Helpers ====

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def add_rect(slide, left, top, width, height, fill=None, border=None, border_pt=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border is not None:
        shape.line.color.rgb = border
        if border_pt:
            shape.line.width = Pt(border_pt)
        else:
            shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=None, weight=0.75):
    if color is None:
        color = MCK_LIGHT_GRAY
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    return conn


def add_text(slide, left, top, width, height, text, size=10, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="Meiryo UI"):
    if color is None:
        color = MCK_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 30000
    tf.margin_right = 30000
    tf.margin_top = 20000
    tf.margin_bottom = 20000
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_multiline(slide, left, top, width, height, lines, size=10,
                  color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.25, font_name="Meiryo UI"):
    """lines: list of str or tuples (text, bold, color_override, size_override)"""
    if color is None:
        color = MCK_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 30000
    tf.margin_right = 30000
    tf.margin_top = 20000
    tf.margin_bottom = 20000
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        if isinstance(line, tuple):
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            c = line[2] if len(line) > 2 and line[2] is not None else color
            s = line[3] if len(line) > 3 and line[3] is not None else size
        else:
            text = line
            bold = False
            c = color
            s = size
        if text == "":
            run = p.add_run()
            run.text = " "
            run.font.size = Pt(s)
            run.font.name = font_name
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(s)
        run.font.bold = bold
        run.font.color.rgb = c
        run.font.name = font_name
    return tb


def add_tracker(slide, label):
    """Small uppercase breadcrumb at top"""
    add_text(slide, TRACKER_LEFT, TRACKER_Y, CONTENT_W, 150000,
             label.upper(), size=8, bold=True, color=MCK_BLUE)


def add_action_title(slide, title, tracker=None):
    """McKinsey action title + horizontal rule"""
    if tracker:
        add_tracker(slide, tracker)
    add_text(slide, MARGIN_L, ACTION_TITLE_Y, CONTENT_W, ACTION_TITLE_H,
             title, size=16, bold=True, color=MCK_NAVY)
    # Horizontal rule
    add_line(slide, MARGIN_L, RULE_Y, SLIDE_W - MARGIN_R, RULE_Y,
             color=MCK_NAVY, weight=1.5)


def add_footer(slide, page, total, source=None):
    if source:
        add_text(slide, MARGIN_L, FOOTER_Y, CONTENT_W - 500000, 200000,
                 f"Source: {source}", size=8, color=MCK_GRAY)
    add_text(slide, SLIDE_W - MARGIN_R - 800000, FOOTER_Y,
             800000, 200000,
             f"{page}", size=9, color=MCK_GRAY, align=PP_ALIGN.RIGHT)
    # Thin line above footer
    add_line(slide, MARGIN_L, FOOTER_Y - 30000,
             SLIDE_W - MARGIN_R, FOOTER_Y - 30000,
             color=MCK_LIGHT_GRAY, weight=0.5)


def add_labeled_box(slide, left, top, width, height, label, body_lines,
                    label_color=None, border_color=None, label_bg=None):
    """Box with colored label bar and body text."""
    if label_color is None:
        label_color = MCK_NAVY
    if border_color is None:
        border_color = MCK_LIGHT_GRAY
    # Outer border
    add_rect(slide, left, top, width, height,
             fill=WHITE, border=border_color, border_pt=0.75)
    # Label bar
    label_h = 320000
    if label_bg is None:
        label_bg = label_color
    add_rect(slide, left, top, width, label_h, fill=label_bg)
    add_text(slide, left + 150000, top + 50000, width - 300000, 220000,
             label, size=10, bold=True, color=WHITE)
    # Body
    add_multiline(slide, left + 150000, top + label_h + 80000,
                  width - 300000, height - label_h - 160000,
                  body_lines, size=10, color=MCK_DARK, line_spacing=1.35)


def add_chat_pair(slide, left, top, width, user_lines, ai_lines, height=None):
    """Compact chat exchange. Returns actual height used."""
    user_h = 70000 + 180000 * len(user_lines)
    ai_h = 70000 + 180000 * len(ai_lines)
    gap = 100000

    # User
    add_rect(slide, left, top, width, 180000, fill=MCK_BG_LIGHT)
    add_text(slide, left + 60000, top, 600000, 180000,
             "👤", size=9, color=MCK_NAVY)
    add_text(slide, left + 260000, top, width - 260000, 180000,
             "You", size=8, bold=True, color=MCK_NAVY)
    add_multiline(slide, left + 60000, top + 180000, width - 120000, user_h,
                  user_lines, size=9, color=MCK_DARK, line_spacing=1.25)

    ai_top = top + user_h + 180000 + gap
    add_rect(slide, left, ai_top, width, 180000, fill=MCK_LIGHT_BLUE)
    add_text(slide, left + 60000, ai_top, 600000, 180000,
             "🤖", size=9, color=MCK_NAVY)
    add_text(slide, left + 260000, ai_top, width - 260000, 180000,
             "Claude", size=8, bold=True, color=MCK_NAVY)
    add_multiline(slide, left + 60000, ai_top + 180000, width - 120000, ai_h,
                  ai_lines, size=9, color=MCK_DARK, line_spacing=1.25)

    total_h = user_h + 180000 + gap + ai_h + 180000
    return total_h


def add_code_block(slide, left, top, width, height, code_lines, size=8):
    add_rect(slide, left, top, width, height,
             fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY, border_pt=0.5)
    lines = [(line, False, MCK_DARK, size) if isinstance(line, str) else line
             for line in code_lines]
    add_multiline(slide, left + 100000, top + 80000,
                  width - 200000, height - 160000,
                  lines, size=size, line_spacing=1.2,
                  font_name="Consolas")


def add_section_number(slide, left, top, size, number, color=None):
    """Large number as section marker."""
    if color is None:
        color = MCK_BLUE
    add_text(slide, left, top, size, size,
             str(number), size=int(size / 38000), bold=True,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================
# SLIDE BUILDERS
# ==========================================

def s01_title(prs, page, total):
    slide = new_slide(prs)
    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill=MCK_BLUE)

    # Main title block (left-aligned, McKinsey style)
    add_text(slide, MARGIN_L, 1700000, CONTENT_W, 300000,
             "DX PORTAL TEAM  |  PROPOSAL", size=9, bold=True, color=MCK_BLUE)
    add_text(slide, MARGIN_L, 2050000, CONTENT_W, 600000,
             "チームで集合知化を目指す", size=32, bold=True, color=MCK_NAVY)
    add_text(slide, MARGIN_L, 2750000, CONTENT_W, 500000,
             "Claude Code × Gitea で、個人の知恵をチームの力に変える",
             size=16, color=MCK_DARK)
    add_text(slide, MARGIN_L, 3250000, CONTENT_W, 400000,
             "~ 方向性の共有と、一緒に進めたいことの相談 ~",
             size=11, color=MCK_GRAY)
    # Bottom bar
    add_line(slide, MARGIN_L, 4600000, MARGIN_L + 1500000, 4600000,
             color=MCK_NAVY, weight=2)
    add_text(slide, MARGIN_L, 4650000, CONTENT_W, 250000,
             "DX ポータル開発チーム", size=10, color=MCK_GRAY)
    return slide


def s02_executive_summary(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "個人のノウハウがチームの力に変わる仕組みを、3名×3ヶ月のパイロットで立ち上げる",
        tracker="Executive summary")

    # 5 key messages as boxes
    box_top = CONTENT_TOP + 100000
    box_h = 3400000
    col_w = (CONTENT_W - 200000) // 5
    gap = 50000

    messages = [
        ("01", "現状", "優秀な人材のノウハウが\n個人の頭の中に留まり\n消えていく",
         "属人化・同じ失敗の反復・\n新人の立ち上がり遅延",
         MCK_RED),
        ("02", "目指す姿", "個人の知恵が消えず\nチームの力になる\n「集合知化」",
         "残す・揃える・掛け合わせる\n3つの方向で",
         MCK_NAVY),
        ("03", "手段", "既存の Gitea を\n作業場所と脳に変える\n（+ Claude Code）",
         "Issue → 壁打ち → PR →\nレビュー → 学習定着",
         MCK_BLUE),
        ("04", "仕組み", "team-standards に\nCLAUDE.md・Skills を\n集約し育てる",
         "共通層 + プロジェクト層で\n階層管理",
         MCK_BLUE),
        ("05", "進め方", "3名×3ヶ月で\nパイロット\n→ 成功したら7名に展開",
         "撤退条件・属人化対策を\n最初から組み込む",
         MCK_GREEN),
    ]

    for i, (num, label, main, detail, color) in enumerate(messages):
        x = MARGIN_L + i * (col_w + gap)
        # Number bar
        add_rect(slide, x, box_top, col_w, 70000, fill=color)
        # Box
        add_rect(slide, x, box_top + 70000, col_w, box_h - 70000,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        # Number
        add_text(slide, x + 100000, box_top + 150000, col_w - 200000, 300000,
                 num, size=11, bold=True, color=color)
        # Label
        add_text(slide, x + 100000, box_top + 430000, col_w - 200000, 280000,
                 label, size=9, bold=True, color=MCK_GRAY)
        # Main
        main_lines = [(l, True, MCK_NAVY, 10) for l in main.split("\n")]
        add_multiline(slide, x + 100000, box_top + 730000, col_w - 200000, 1200000,
                      main_lines, size=10, color=MCK_NAVY,
                      line_spacing=1.3)
        # Divider
        add_line(slide, x + 150000, box_top + 2000000,
                 x + col_w - 150000, box_top + 2000000,
                 color=MCK_LIGHT_GRAY)
        # Detail
        add_multiline(slide, x + 100000, box_top + 2100000, col_w - 200000, 1100000,
                      detail.split("\n"), size=8, color=MCK_GRAY,
                      line_spacing=1.35)

    add_footer(slide, page, total)
    return slide


def s03_current_state(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "優秀な人材がいるのに、ノウハウが個人に留まり消えている — もったいない",
        tracker="現状認識")

    # Top: one-line message box
    top = CONTENT_TOP
    add_rect(slide, MARGIN_L, top, CONTENT_W, 500000,
             fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
    add_multiline(slide, MARGIN_L + 200000, top + 80000,
                  CONTENT_W - 400000, 350000,
                  [("バイブコーディングはしている。各自それなりに使いこなしている。", False, MCK_DARK, 11),
                   ("でも、やり方がバラバラで、ノウハウが個人の頭の中だけ。", True, MCK_NAVY, 11)],
                  size=11, line_spacing=1.3)

    # Below: 2 examples side by side
    ex_top = top + 650000
    ex_h = 2650000
    col_w = (CONTENT_W - 200000) // 2

    # Example 1: Aさんが抜けたら
    left1 = MARGIN_L
    add_labeled_box(slide, left1, ex_top, col_w, ex_h,
                    "ケース ①  Aさんのノウハウは、Aさんと一緒に消える",
                    [],
                    label_color=MCK_RED, border_color=MCK_LIGHT_GRAY,
                    label_bg=MCK_RED)

    # Arrow flow inside box
    inner_top = ex_top + 400000
    # Before
    add_text(slide, left1 + 150000, inner_top, 1800000, 250000,
             "Aさんの頭の中", size=9, bold=True, color=MCK_GRAY)
    add_multiline(slide, left1 + 150000, inner_top + 280000, 1800000, 1500000,
                  [("● DBタイムアウト対処法", False, MCK_DARK, 8),
                   ("● 認証トークンの落とし穴", False, MCK_DARK, 8),
                   ("● 失敗した設計3つ", False, MCK_DARK, 8),
                   ("● 効いたプロンプト", False, MCK_DARK, 8),
                   ("● レビューでの指摘", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.4)
    # Arrow
    add_text(slide, left1 + 2050000, inner_top + 700000, 400000, 300000,
             "異動", size=9, bold=True, color=MCK_RED, align=PP_ALIGN.CENTER)
    add_text(slide, left1 + 2050000, inner_top + 950000, 400000, 300000,
             "→", size=16, bold=True, color=MCK_RED, align=PP_ALIGN.CENTER)
    # After
    add_text(slide, left1 + 2500000, inner_top, 1500000, 250000,
             "残ったもの", size=9, bold=True, color=MCK_GRAY)
    add_text(slide, left1 + 2500000, inner_top + 280000, 1500000, 700000,
             "???", size=36, bold=True, color=MCK_RED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, left1 + 2500000, inner_top + 1000000, 1500000, 600000,
             "コードは残るが\n「なぜこう作ったか」は\n誰も知らない",
             size=8, color=MCK_GRAY, align=PP_ALIGN.CENTER)

    # Example 2: 同じ失敗の反復
    left2 = MARGIN_L + col_w + 200000
    add_labeled_box(slide, left2, ex_top, col_w, ex_h,
                    "ケース ②  同じ DB タイムアウト問題で2人が3時間ハマる",
                    [],
                    label_color=MCK_RED, border_color=MCK_LIGHT_GRAY,
                    label_bg=MCK_RED)

    # Timeline inside
    ti_top = ex_top + 400000
    add_text(slide, left2 + 150000, ti_top, 500000, 250000,
             "2月", size=9, bold=True, color=MCK_GRAY)
    add_text(slide, left2 + 650000, ti_top, col_w - 800000, 250000,
             "Aさんが3時間ハマる", size=10, bold=True, color=MCK_DARK)
    add_multiline(slide, left2 + 650000, ti_top + 260000,
                  col_w - 800000, 600000,
                  [("タイムアウト値変更 → 効果なし", False, MCK_DARK, 8),
                   ("ドライバ調査 → 関係なし", False, MCK_DARK, 8),
                   ("プール設定で解決、記録せず", False, MCK_GREEN, 8)],
                  size=8, line_spacing=1.3)

    add_line(slide, left2 + 150000, ti_top + 950000,
             left2 + col_w - 150000, ti_top + 950000,
             color=MCK_LIGHT_GRAY)

    add_text(slide, left2 + 150000, ti_top + 1000000, 500000, 250000,
             "5月", size=9, bold=True, color=MCK_GRAY)
    add_text(slide, left2 + 650000, ti_top + 1000000, col_w - 800000, 250000,
             "Bさんが同じ問題で3時間ハマる", size=10, bold=True, color=MCK_RED)
    add_multiline(slide, left2 + 650000, ti_top + 1260000,
                  col_w - 800000, 700000,
                  [("同じ順序で同じ失敗を辿る", False, MCK_DARK, 8),
                   ("結局プール設定で解決", False, MCK_GREEN, 8),
                   ("Aさんも同じことを経験済みと後で知る", False, MCK_RED, 8)],
                  size=8, line_spacing=1.3)

    add_footer(slide, page, total)
    return slide


def s04_risks(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "このまま進むと、個人成長はあってもチームは賢くならない — 5つのリスク",
        tracker="放置した場合のリスク")

    top = CONTENT_TOP + 100000
    row_h = 650000

    risks = [
        ("1", "同じ失敗を繰り返す",
         "Aさんの3時間 → Bさんの3時間 → Cさんの3時間。\n年間で数十時間の損失",
         MCK_RED),
        ("2", "人が抜けた瞬間、知識が消える",
         "「あれどうなってるの？」が誰にも答えられない状態",
         MCK_RED),
        ("3", "新メンバーの立ち上がりに時間がかかる",
         "「先輩に聞いて覚える」しかない。先輩も新人も疲弊する",
         MCK_ORANGE),
        ("4", "AIを使っても、チーム全体は賢くならない",
         "個人は熟達する。でもチームとしての成長がない",
         MCK_ORANGE),
        ("5", "DX推進側なのに、自分たちの開発は旧態依然",
         "「DXを語る資格があるのか」と見られるリスク",
         MCK_ORANGE),
    ]
    for i, (num, title, desc, color) in enumerate(risks):
        y = top + i * row_h
        # Number
        add_rect(slide, MARGIN_L, y + 80000, 400000, 400000, fill=color)
        add_text(slide, MARGIN_L, y + 80000, 400000, 400000,
                 num, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(slide, MARGIN_L + 550000, y + 50000, CONTENT_W - 700000, 320000,
                 title, size=12, bold=True, color=color)
        # Description
        add_multiline(slide, MARGIN_L + 550000, y + 360000,
                      CONTENT_W - 700000, 280000,
                      desc.split("\n"), size=9, color=MCK_DARK, line_spacing=1.3)
    add_footer(slide, page, total)
    return slide


def s05_goal(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "目指すのは集合知化 — 個人の知恵が3つの軸でチームの力に変わる状態",
        tracker="目指す姿")

    # Definition
    top = CONTENT_TOP + 100000
    add_rect(slide, MARGIN_L, top, CONTENT_W, 500000,
             fill=MCK_NAVY)
    add_text(slide, MARGIN_L, top + 80000, CONTENT_W, 350000,
             "集合知化 (Collective Intelligence)",
             size=13, bold=True, color=MCK_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_L, top + 250000, CONTENT_W, 250000,
             "個人の知恵が、消えずに、チームの力になる",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 3 dimensions
    dim_top = top + 700000
    dim_h = 2600000
    col_w = (CONTENT_W - 400000) // 3

    dims = [
        ("残す", "時間軸で", "優秀な人の知恵が\n異動・退職後も残る",
         "Aさんが見つけた「DBタイムアウトはプール設定」が、\nAさんが抜けた後も全員の秘書が知っている",
         MCK_NAVY),
        ("揃える", "横軸で", "誰がやっても同じレベル\nでアウトプットが出せる",
         "新人でもベテランでも、秘書が同じ設計手順・\n同じ規約・同じチェックで作業する",
         MCK_BLUE),
        ("掛け合わせる", "化学反応で", "1+1が3になる",
         "Aさんの知恵 + Bさんの知恵 + Cさんの知恵が\n組み合わさって、一人では届かない品質に",
         MCK_GREEN),
    ]
    for i, (word, sub, main, example, color) in enumerate(dims):
        x = MARGIN_L + i * (col_w + 200000)
        # Card
        add_rect(slide, x, dim_top, col_w, dim_h,
                 fill=WHITE, border=color, border_pt=1.5)
        # Header bar
        add_rect(slide, x, dim_top, col_w, 70000, fill=color)
        # Sub
        add_text(slide, x + 100000, dim_top + 130000, col_w - 200000, 220000,
                 sub, size=8, color=MCK_GRAY)
        # Word
        add_text(slide, x + 100000, dim_top + 370000, col_w - 200000, 500000,
                 word, size=20, bold=True, color=color)
        # Main message
        main_lines = [(l, True, MCK_DARK, 11) for l in main.split("\n")]
        add_multiline(slide, x + 100000, dim_top + 900000, col_w - 200000, 700000,
                      main_lines, size=11, color=MCK_DARK,
                      line_spacing=1.3)
        # Line
        add_line(slide, x + 100000, dim_top + 1750000,
                 x + col_w - 100000, dim_top + 1750000,
                 color=MCK_LIGHT_GRAY)
        # Example
        add_text(slide, x + 100000, dim_top + 1820000, col_w - 200000, 200000,
                 "具体例", size=8, color=MCK_GRAY)
        add_multiline(slide, x + 100000, dim_top + 2020000, col_w - 200000, 500000,
                      example.split("\n"), size=8, color=MCK_DARK,
                      line_spacing=1.35)

    add_footer(slide, page, total)
    return slide


def s06_vegapunk_and_why_gitea(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "ベガパンクのサテライトと同じ仕組みを、既存のGiteaで実現できる",
        tracker="アプローチの原型")

    # Left: Vegapunk analogy
    top = CONTENT_TOP + 100000
    box_h = 3400000
    left_w = (CONTENT_W - 200000) * 45 // 100

    add_labeled_box(slide, MARGIN_L, top, left_w, box_h,
                    "たとえ話：ベガパンクのサテライト",
                    [("1人の天才が6体のサテライト（分身）に分かれ、", False, MCK_DARK, 10),
                     ("それぞれ独立して動く。", False, MCK_DARK, 10),
                     ("", False),
                     ("全員が『パンクレコーズ』という", False, MCK_DARK, 10),
                     ("共有の知識ベースでつながっている。", False, MCK_DARK, 10),
                     ("", False),
                     ("→ 1体が学んだことは、全員の知識になる", False, MCK_NAVY, 10),
                     ("", False),
                     ("━━━━━━━━━━━━━━━━━━━━", False, MCK_LIGHT_GRAY, 8),
                     ("", False),
                     ("私たちの場合", True, MCK_NAVY, 10),
                     ("● サテライト = 各自の AI 秘書", False, MCK_DARK, 10),
                     ("● パンクレコーズ = team-standards", False, MCK_DARK, 10),
                     ("● 共有知識 = CLAUDE.md / Skills", False, MCK_DARK, 10)],
                    label_color=MCK_NAVY, label_bg=MCK_NAVY)

    # Right: Why Gitea (3 reasons)
    right_x = MARGIN_L + left_w + 200000
    right_w = CONTENT_W - left_w - 200000
    add_text(slide, right_x, top, right_w, 300000,
             "なぜ GitHub ではなく Gitea か", size=11, bold=True, color=MCK_NAVY)

    reasons = [
        ("01", "すでに社内にある", "サーバー・アカウント完備。\n追加投資ゼロ", MCK_NAVY),
        ("02", "コードを社外に出せない", "セキュリティ・ポリシー上\nGitHub Cloud は使えない", MCK_RED),
        ("03", "コストを最小化できる", "GitHub Team の月額でも\n予算化が難しい", MCK_ORANGE),
    ]
    r_top = top + 350000
    r_h = 950000
    for i, (num, title, desc, color) in enumerate(reasons):
        y = r_top + i * (r_h + 50000)
        add_rect(slide, right_x, y, right_w, r_h,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        # Number stripe
        add_rect(slide, right_x, y, 50000, r_h, fill=color)
        add_text(slide, right_x + 150000, y + 100000, 400000, 300000,
                 num, size=10, bold=True, color=color)
        add_text(slide, right_x + 600000, y + 100000, right_w - 700000, 300000,
                 title, size=11, bold=True, color=MCK_DARK)
        add_multiline(slide, right_x + 600000, y + 430000,
                      right_w - 700000, 500000,
                      desc.split("\n"), size=9, color=MCK_GRAY, line_spacing=1.3)
    add_footer(slide, page, total)
    return slide


def s07_direction(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "方向性：Gitea を『完成コードの墓場』から『チームの脳と作業場所』に変える",
        tracker="方向性")

    # Top: one-liner
    top = CONTENT_TOP + 50000
    add_rect(slide, MARGIN_L, top, CONTENT_W, 550000,
             fill=MCK_NAVY)
    add_text(slide, MARGIN_L, top + 100000, CONTENT_W, 300000,
             "書く・考える・決める・学ぶ",
             size=12, bold=True, color=MCK_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_L, top + 280000, CONTENT_W, 250000,
             "その全部がチームの資産として溜まる仕組みを作る",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Diagram: Workstation + Workshop
    dia_top = top + 750000
    dia_h = 2600000

    # Left: Workstation
    left_w = (CONTENT_W - 1400000) // 2
    add_rect(slide, MARGIN_L, dia_top, left_w, dia_h,
             fill=MCK_BG_LIGHT, border=MCK_BLUE, border_pt=1.5)
    add_rect(slide, MARGIN_L, dia_top, left_w, 320000, fill=MCK_BLUE)
    add_text(slide, MARGIN_L + 150000, dia_top + 50000, left_w - 300000, 220000,
             "職人の作業台：手元のPC", size=10, bold=True, color=WHITE)

    add_text(slide, MARGIN_L + 150000, dia_top + 400000, left_w - 300000, 300000,
             "VS Code + Claude Code", size=13, bold=True, color=MCK_NAVY)
    add_text(slide, MARGIN_L + 150000, dia_top + 730000, left_w - 300000, 200000,
             "変わらない。むしろ強化する", size=9, color=MCK_GRAY)

    add_multiline(slide, MARGIN_L + 150000, dia_top + 1000000,
                  left_w - 300000, 1500000,
                  [("● コードを書く", False, MCK_DARK, 10),
                   ("● 秘書と壁打ち", False, MCK_DARK, 10),
                   ("● テスト・デバッグ", False, MCK_DARK, 10),
                   ("● 設計の試行錯誤", False, MCK_DARK, 10)],
                  size=10, line_spacing=1.5)

    # Middle: arrows
    arrow_x = MARGIN_L + left_w + 100000
    arrow_w = 1200000
    add_text(slide, arrow_x, dia_top + 800000, arrow_w, 300000,
             "push / PR", size=9, color=MCK_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x, dia_top + 1050000, arrow_w, 400000,
             "▶", size=22, bold=True, color=MCK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x, dia_top + 1500000, arrow_w, 400000,
             "◀", size=22, bold=True, color=MCK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x, dia_top + 1900000, arrow_w, 300000,
             "pull / review", size=9, color=MCK_GRAY, align=PP_ALIGN.CENTER)

    # Right: Workshop
    right_x = arrow_x + arrow_w + 100000
    right_w = CONTENT_W - left_w - 1400000
    add_rect(slide, right_x, dia_top, right_w, dia_h,
             fill=MCK_BG_LIGHT, border=MCK_NAVY, border_pt=1.5)
    add_rect(slide, right_x, dia_top, right_w, 320000, fill=MCK_NAVY)
    add_text(slide, right_x + 150000, dia_top + 50000, right_w - 300000, 220000,
             "共同作業場：Gitea", size=10, bold=True, color=WHITE)

    add_text(slide, right_x + 150000, dia_top + 400000, right_w - 300000, 300000,
             "チームの脳と記憶", size=13, bold=True, color=MCK_NAVY)
    add_text(slide, right_x + 150000, dia_top + 730000, right_w - 300000, 200000,
             "使い方を変えるだけ。追加投資ゼロ", size=9, color=MCK_GRAY)

    add_multiline(slide, right_x + 150000, dia_top + 1000000,
                  right_w - 300000, 1500000,
                  [("● Issue（何を・なぜ）", False, MCK_DARK, 10),
                   ("● 壁打ちログの要約", False, MCK_DARK, 10),
                   ("● PR / レビュー（判断）", False, MCK_DARK, 10),
                   ("● CLAUDE.md（チームの脳）", False, MCK_DARK, 10)],
                  size=10, line_spacing=1.5)

    add_footer(slide, page, total)
    return slide


def s08_day_in_life(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "1日の開発は、Gitea と手元ツールの行き来で進む — 8ステップ",
        tracker="日々の動き")

    top = CONTENT_TOP + 200000
    steps = [
        ("1", "朝", "Issue\n選定", "Gitea", MCK_NAVY),
        ("2", "開始", "壁打ち", "手元", MCK_BLUE),
        ("3", "整理", "要約を\nIssueへ", "Gitea", MCK_NAVY),
        ("4", "実装", "コード\n作成", "手元", MCK_BLUE),
        ("5", "途中", "Draft\nPR", "Gitea", MCK_NAVY),
        ("6", "確認", "セルフ\nレビュー", "手元", MCK_BLUE),
        ("7", "完了", "人間\nレビュー", "Gitea", MCK_NAVY),
        ("8", "学び", "CLAUDE.md\n更新", "Gitea", MCK_GREEN),
    ]
    n = len(steps)
    step_w = (CONTENT_W - (n - 1) * 80000) // n
    step_h = 1400000

    for i, (num, time, action, place, color) in enumerate(steps):
        x = MARGIN_L + i * (step_w + 80000)
        # Number circle at top
        circle_s = 280000
        cx = x + (step_w - circle_s) // 2
        add_rect(slide, cx, top, circle_s, circle_s, fill=color)
        add_text(slide, cx, top, circle_s, circle_s,
                 num, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Time
        add_text(slide, x, top + 320000, step_w, 200000,
                 time, size=9, color=MCK_GRAY, align=PP_ALIGN.CENTER)
        # Card
        card_top = top + 550000
        add_rect(slide, x, card_top, step_w, step_h - 550000,
                 fill=WHITE, border=color, border_pt=1)
        # Action
        act_lines = [(l, True, color, 10) for l in action.split("\n")]
        add_multiline(slide, x + 20000, card_top + 100000,
                      step_w - 40000, 600000,
                      act_lines, size=10, color=color,
                      align=PP_ALIGN.CENTER)
        # Place
        add_text(slide, x, card_top + 650000, step_w, 180000,
                 place, size=8, color=MCK_GRAY, align=PP_ALIGN.CENTER)

    # Legend
    leg_top = top + 1700000
    add_rect(slide, MARGIN_L, leg_top, CONTENT_W, 280000, fill=MCK_BG_LIGHT)
    add_multiline(slide, MARGIN_L + 200000, leg_top + 50000,
                  CONTENT_W - 400000, 200000,
                  [("■ 紺色＝Gitea（チームの場所）      ■ 青＝手元（VS Code + Claude Code）      ■ 緑＝学びの定着", False, MCK_GRAY, 9)],
                  size=9, line_spacing=1.2)

    # Key insight
    add_text(slide, MARGIN_L, leg_top + 450000, CONTENT_W, 350000,
             "書く場所は変わらない。変わるのは「流れ」と「記憶」の置き場所だけ",
             size=12, bold=True, color=MCK_NAVY, align=PP_ALIGN.CENTER)

    add_footer(slide, page, total)
    return slide


def s09_chats_planning(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "チャット例①：朝の壁打ちから、方針決定・Issue ポストまで",
        tracker="具体的な会話イメージ 1/2")

    # Three chat exchanges side by side
    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 300000) // 3
    col_h = 3500000

    # Column 1: Morning pickup
    x = MARGIN_L
    add_rect(slide, x, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 200000,
             "STEP 1  Issue を選んで壁打ち開始",
             size=9, bold=True, color=WHITE)
    add_chat_pair(slide, x, top + 320000, col_w,
                  ["Issue #42 『申請フォーム",
                   "リニューアル』を担当する。",
                   "全体の設計から相談したい。"],
                  ["Issue と CLAUDE.md を確認",
                   "しました。理解：",
                   "・入力15項目→8項目に削減",
                   "・紙申請廃止、オンライン化",
                   "・期限 4/15",
                   "",
                   "確認したいこと：",
                   "① 削減項目は決まってる？",
                   "② 既存データの移行は？",
                   "③ 承認フローは変わる？"])

    # Column 2: Design decision
    x = MARGIN_L + col_w + 150000
    add_rect(slide, x, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 200000,
             "STEP 2  方針を決める",
             size=9, bold=True, color=WHITE)
    add_chat_pair(slide, x, top + 320000, col_w,
                  ["① 未定（現場ヒアリング）",
                   "② 移行必要",
                   "③ 変わらない"],
                  ["2つの案があります：",
                   "",
                   "案A：段階的改修",
                   " → リスク小、中途半端期間",
                   "",
                   "案B：新画面を作って切替",
                   " → リスク大、クリーン移行",
                   "",
                   "過去事例を見ると『Aは2回",
                   "作業になる』記録があります。",
                   "今回の期間なら案B推奨。"])

    # Column 3: Issue post
    x = MARGIN_L + (col_w + 150000) * 2
    add_rect(slide, x, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 200000,
             "STEP 3  結論を Issue にポスト",
             size=9, bold=True, color=WHITE)
    add_chat_pair(slide, x, top + 320000, col_w,
                  ["案Bでいこう。今の議論を",
                   "Issue #42 にコメント投稿",
                   "しておいて。"],
                  ["4項目テンプレで投稿します：",
                   "",
                   "【前提】4/15期限で全面刷新",
                   "【選択肢】A段階改修/B新設",
                   "【結論】案B採用",
                   "【却下理由】Aは過去2回作業",
                   "",
                   "Issue #42 にコメント完了。",
                   "Draft PR の雛形も用意",
                   "しましょうか？"])

    add_footer(slide, page, total)
    return slide


def s10_chats_execution(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "チャット例②：実装・セルフレビュー・CLAUDE.md 更新まで",
        tracker="具体的な会話イメージ 2/2")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 300000) // 3

    # Column 1: Implementation
    x = MARGIN_L
    add_rect(slide, x, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 200000,
             "STEP 4  実装を始める",
             size=9, bold=True, color=WHITE)
    add_chat_pair(slide, x, top + 320000, col_w,
                  ["では実装へ。フォームの",
                   "画面から作りたい。"],
                  ["CLAUDE.md を再確認：",
                   "・Serendie を使う",
                   "・any 禁止",
                   "・テストファースト",
                   "",
                   "この順で進めます：",
                   "1. 型定義",
                   "2. バリデーションのテスト",
                   "3. フォーム実装",
                   "4. 確認画面",
                   "",
                   "型定義から始めます..."])

    # Column 2: Self-review
    x = MARGIN_L + col_w + 150000
    add_rect(slide, x, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 200000,
             "STEP 5  push 前にセルフレビュー",
             size=9, bold=True, color=WHITE)
    add_chat_pair(slide, x, top + 320000, col_w,
                  ["/review-my-changes"],
                  ["✓ OK",
                   "・Serendie 使用",
                   "・any なし",
                   "・テスト3件追加",
                   "",
                   "⚠ 気になる点",
                   "・validateZipCode にテスト無",
                   "・console.log が残っている",
                   "",
                   "❌ CLAUDE.md 違反",
                   "・コミット prefix 抜け",
                   "",
                   "修正しますか？"])

    # Column 3: CLAUDE.md update
    x = MARGIN_L + (col_w + 150000) * 2
    add_rect(slide, x, top, col_w, 280000, fill=MCK_GREEN)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 200000,
             "STEP 6  学びを CLAUDE.md に",
             size=9, bold=True, color=WHITE)
    add_chat_pair(slide, x, top + 320000, col_w,
                  ["郵便番号のバリデーションで",
                   "ハマった。ハイフン両対応が",
                   "大事と分かった。",
                   "CLAUDE.md に追記したい。"],
                  ["team-standards への追記",
                   "PR を作成します：",
                   "",
                   "## 郵便番号のバリデーション",
                   "- ハイフン あり/なし 両対応",
                   "- DB 保存は正規化後に統一",
                   "",
                   "この内容で PR を作成して",
                   "よろしいですか？"])

    add_footer(slide, page, total)
    return slide


def s11_team_standards(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "team-standards リポジトリに全員の共通資産を集約し、@import で各プロジェクトから読み込む",
        tracker="仕組みの中心")

    top = CONTENT_TOP + 100000

    # Left: team-standards structure
    left_w = (CONTENT_W - 200000) * 50 // 100
    add_rect(slide, MARGIN_L, top, left_w, 3400000,
             fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 150000, top + 120000, left_w - 300000, 280000,
             "📦  team-standards  /  Gitea 上の1リポジトリ",
             size=11, bold=True, color=WHITE)

    items = [
        ("CLAUDE.md", "常に読まれる基本ルール・規約", "any禁止、Serendie使用、規約"),
        ("skills/", "呼び出し式の手順書", "/implement-feature, /review"),
        ("agents/", "専門秘書（分野ごと）", "security / architecture / test"),
        ("hooks/", "自動化スクリプト", "pre-commit, session-end"),
        ("mcp/", "外部ツール連携設定", "Gitea MCP, Slack MCP"),
    ]
    for i, (name, desc, example) in enumerate(items):
        y = top + 550000 + i * 520000
        # Name
        add_text(slide, MARGIN_L + 200000, y, left_w - 400000, 240000,
                 name, size=10, bold=True, color=RGBColor(255, 220, 100))
        # Desc
        add_text(slide, MARGIN_L + 1800000, y, left_w - 2000000, 240000,
                 desc, size=9, color=WHITE)
        # Example
        add_text(slide, MARGIN_L + 1800000, y + 220000, left_w - 2000000, 220000,
                 f"例: {example}", size=8, color=MCK_LIGHT_BLUE)

    # Right: hierarchy
    right_x = MARGIN_L + left_w + 200000
    right_w = CONTENT_W - left_w - 200000

    add_text(slide, right_x, top, right_w, 250000,
             "2階層で管理", size=11, bold=True, color=MCK_NAVY)

    # Common box
    cb_top = top + 350000
    add_rect(slide, right_x, cb_top, right_w, 1100000,
             fill=MCK_LIGHT_BLUE, border=MCK_NAVY, border_pt=1.5)
    add_text(slide, right_x + 100000, cb_top + 80000, right_w - 200000, 240000,
             "共通（team-standards）", size=10, bold=True, color=MCK_NAVY)
    add_multiline(slide, right_x + 100000, cb_top + 350000,
                  right_w - 200000, 700000,
                  [("● Serendie を使う", False, MCK_DARK, 9),
                   ("● TypeScript で any 禁止", False, MCK_DARK, 9),
                   ("● 壁打ち→設計→実装のフロー", False, MCK_DARK, 9)],
                  size=9, line_spacing=1.3)

    # Arrow
    add_text(slide, right_x, cb_top + 1150000, right_w, 250000,
             "↑ @import で読み込む ↑", size=9, color=MCK_GRAY,
             align=PP_ALIGN.CENTER)

    # Project box
    pb_top = cb_top + 1450000
    add_rect(slide, right_x, pb_top, right_w, 1550000,
             fill=WHITE, border=MCK_BLUE, border_pt=1.5)
    add_text(slide, right_x + 100000, pb_top + 80000, right_w - 200000, 240000,
             "プロジェクト個別（各リポジトリ）", size=10, bold=True, color=MCK_BLUE)
    add_multiline(slide, right_x + 100000, pb_top + 350000,
                  right_w - 200000, 1100000,
                  [("● プロジェクトの目的・背景", False, MCK_DARK, 9),
                   ("● 使っている技術スタック", False, MCK_DARK, 9),
                   ("● ディレクトリ構造", False, MCK_DARK, 9),
                   ("● プロジェクト固有の注意点", False, MCK_DARK, 9),
                   ("● デプロイ手順", False, MCK_DARK, 9)],
                  size=9, line_spacing=1.3)

    add_footer(slide, page, total)
    return slide


def s12_components_detail(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "CLAUDE.md は常識・Skills は手順・Subagents は専門秘書 — 役割で使い分ける",
        tracker="構成要素の具体例")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 300000) // 3
    col_h = 3500000

    # Col 1: CLAUDE.md
    x = MARGIN_L
    add_rect(slide, x, top, col_w, 300000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 220000,
             "CLAUDE.md  （常識）", size=10, bold=True, color=WHITE)
    add_text(slide, x + 100000, top + 330000, col_w - 200000, 200000,
             "常に読まれる。短く（200行以内）", size=8, color=MCK_GRAY)

    code1 = [
        "# team-standards",
        "",
        "## 開発の流れ",
        "1. 15分以上→Issue",
        "2. 壁打ち→4項目テンプレ",
        "3. Draft PR",
        "4. セルフ/人間レビュー",
        "5. 学びは追記PR",
        "",
        "## 規約",
        "- any 禁止",
        "- Serendie 使用",
        "- コミット prefix 必須",
        "",
        "## 失敗あるある",
        "- DBタイムアウト→プール",
        "- 郵便番号→ハイフン対応",
    ]
    add_code_block(slide, x, top + 600000, col_w, col_h - 600000,
                   code1, size=8)

    # Col 2: Skills
    x = MARGIN_L + col_w + 150000
    add_rect(slide, x, top, col_w, 300000, fill=MCK_BLUE)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 220000,
             "Skills  （手順）", size=10, bold=True, color=WHITE)
    add_text(slide, x + 100000, top + 330000, col_w - 200000, 200000,
             "呼び出し式。/name で起動", size=8, color=MCK_GRAY)

    code2 = [
        "---",
        "name: implement-feature",
        "description: 機能追加の手順",
        "---",
        "",
        "# 機能追加の標準手順",
        "",
        "1. 要件を Issue で確認",
        "2. 過去事例を検索",
        "3. 4項目テンプレで方針整理",
        "4. Draft PR 作成",
        "5. テストを先に書く",
        "6. 実装する",
        "7. /review-my-changes",
        "8. Ready for review",
    ]
    add_code_block(slide, x, top + 600000, col_w, col_h - 600000,
                   code2, size=8)

    # Col 3: Subagents
    x = MARGIN_L + (col_w + 150000) * 2
    add_rect(slide, x, top, col_w, 300000, fill=MCK_GREEN)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 220000,
             "Subagents  （専門秘書）", size=10, bold=True, color=WHITE)
    add_text(slide, x + 100000, top + 330000, col_w - 200000, 200000,
             "特定分野の専門家を呼び出す", size=8, color=MCK_GRAY)

    # 3 agent cards
    agents = [
        ("security-reviewer", "脆弱性・認証漏れ", "SQLインジェクション\n権限チェック漏れ", MCK_RED),
        ("architecture-reviewer", "設計の妥当性", "責務分離\n拡張性", MCK_NAVY),
        ("test-coverage-checker", "テスト不足", "エッジケース\n異常系", MCK_GREEN),
    ]
    a_top = top + 600000
    a_h = 900000
    for i, (name, role, ex, color) in enumerate(agents):
        y = a_top + i * (a_h + 30000)
        add_rect(slide, x, y, col_w, a_h,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_rect(slide, x, y, 40000, a_h, fill=color)
        add_text(slide, x + 120000, y + 80000, col_w - 200000, 240000,
                 name, size=9, bold=True, color=color)
        add_text(slide, x + 120000, y + 320000, col_w - 200000, 220000,
                 role, size=8, color=MCK_DARK)
        add_multiline(slide, x + 120000, y + 550000, col_w - 200000, 320000,
                      ex.split("\n"), size=7, color=MCK_GRAY, line_spacing=1.2)

    add_footer(slide, page, total)
    return slide


def s13_hooks_mcp_and_rules(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Hooks で自動化・MCP で外部連携・Issue と壁打ちログは型を持って習慣化する",
        tracker="運用の仕組み")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # Left: Hooks + MCP
    # Hooks card
    hx = MARGIN_L
    hy = top
    hh = (col_h - 200000) // 2
    add_rect(slide, hx, hy, col_w, hh,
             fill=WHITE, border=MCK_ORANGE, border_pt=1.5)
    add_rect(slide, hx, hy, col_w, 50000, fill=MCK_ORANGE)
    add_text(slide, hx + 100000, hy + 100000, col_w - 200000, 280000,
             "⚡ Hooks  自動化スクリプト", size=11, bold=True, color=MCK_ORANGE)
    hooks_items = [
        ("コミット前", "lint と型チェックを自動実行"),
        ("セッション終了", "壁打ちログを自動保存"),
        ("PR作成時", "自動セルフレビュー起動"),
    ]
    for i, (w, d) in enumerate(hooks_items):
        y = hy + 500000 + i * 300000
        add_text(slide, hx + 150000, y, 1400000, 220000,
                 f"● {w}", size=9, bold=True, color=MCK_DARK)
        add_text(slide, hx + 1600000, y, col_w - 1800000, 220000,
                 d, size=9, color=MCK_GRAY)
    add_text(slide, hx + 150000, hy + hh - 300000, col_w - 300000, 200000,
             "本体は team-standards/hooks/ で共有", size=8, color=MCK_GRAY)

    # MCP card
    my = hy + hh + 100000
    add_rect(slide, hx, my, col_w, hh,
             fill=WHITE, border=MCK_BLUE, border_pt=1.5)
    add_rect(slide, hx, my, col_w, 50000, fill=MCK_BLUE)
    add_text(slide, hx + 100000, my + 100000, col_w - 200000, 280000,
             "🔌 MCP  外部ツール連携", size=11, bold=True, color=MCK_BLUE)
    mcp_items = [
        ("Gitea MCP", "Issue/PR/コードを秘書が読み書き"),
        ("Slack MCP", "チャンネル投稿を秘書が代行"),
        ("DB MCP", "DBスキーマを秘書が直接参照"),
    ]
    for i, (n, d) in enumerate(mcp_items):
        y = my + 500000 + i * 300000
        add_text(slide, hx + 150000, y, 1400000, 220000,
                 f"● {n}", size=9, bold=True, color=MCK_DARK)
        add_text(slide, hx + 1600000, y, col_w - 1800000, 220000,
                 d, size=9, color=MCK_GRAY)
    add_text(slide, hx + 150000, my + hh - 300000, col_w - 300000, 200000,
             "Gitea 公式 MCP が存在する", size=8, color=MCK_GRAY)

    # Right: Rules
    rx = MARGIN_L + col_w + 200000

    # Issue rule
    add_rect(slide, rx, top, col_w, hh,
             fill=WHITE, border=MCK_NAVY, border_pt=1.5)
    add_rect(slide, rx, top, col_w, 50000, fill=MCK_NAVY)
    add_text(slide, rx + 100000, top + 100000, col_w - 200000, 280000,
             "📋 Issue 起票ルール", size=11, bold=True, color=MCK_NAVY)
    add_rect(slide, rx + 150000, top + 500000, col_w - 300000, 350000,
             fill=MCK_BG_LIGHT)
    add_text(slide, rx + 150000, top + 510000, col_w - 300000, 330000,
             "15分以上かかる作業は Issue 必須",
             size=11, bold=True, color=MCK_NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_multiline(slide, rx + 150000, top + 920000, col_w - 300000, 800000,
                  [("立てる：新規機能・大きな改修・バグ調査", False, MCK_DARK, 9),
                   ("立てない：タイポ、1行修正、変数リネーム", False, MCK_GRAY, 9),
                   ("", False),
                   ("立ち上げ期は『迷ったら立てる』でOK", False, MCK_NAVY, 9)],
                  size=9, line_spacing=1.3)

    # Kabeuchi template
    ky = top + hh + 100000
    add_rect(slide, rx, ky, col_w, hh,
             fill=WHITE, border=MCK_NAVY, border_pt=1.5)
    add_rect(slide, rx, ky, col_w, 50000, fill=MCK_NAVY)
    add_text(slide, rx + 100000, ky + 100000, col_w - 200000, 280000,
             "📝 壁打ちログ 4項目テンプレ", size=11, bold=True, color=MCK_NAVY)
    add_multiline(slide, rx + 150000, ky + 500000, col_w - 300000, 1200000,
                  [("【前提】状況・背景", True, MCK_NAVY, 9),
                   ("【検討した選択肢】2-3個", True, MCK_NAVY, 9),
                   ("【結論】採用した案", True, MCK_NAVY, 9),
                   ("【却下理由】選ばなかった理由", True, MCK_NAVY, 9),
                   ("", False),
                   ("→ ログ全体ではなく、結論と経緯だけ", False, MCK_GRAY, 8),
                   ("→ 5分で書ける。3日で続かない落とし穴を回避", False, MCK_GRAY, 8)],
                  size=9, line_spacing=1.4)

    add_footer(slide, page, total)
    return slide


def s14_review_parallel_conflict(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "レビューは秘書と人間で2段階・並列開発は Git の基本・衝突の8割は秘書が解決",
        tracker="レビューと並行作業")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 300000) // 3
    col_h = 3500000

    # Col 1: 2-stage review
    x = MARGIN_L
    add_rect(slide, x, top, col_w, 300000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 220000,
             "PR レビュー  2段階", size=10, bold=True, color=WHITE)

    # Stage 1
    s1_top = top + 400000
    add_rect(slide, x, s1_top, col_w, 1400000,
             fill=WHITE, border=MCK_GREEN)
    add_text(slide, x + 100000, s1_top + 80000, col_w - 200000, 250000,
             "1. 秘書レビュー", size=10, bold=True, color=MCK_GREEN)
    add_multiline(slide, x + 100000, s1_top + 360000, col_w - 200000, 1000000,
                  [("✓ CLAUDE.md 準拠", False, MCK_DARK, 8),
                   ("✓ テストの存在", False, MCK_DARK, 8),
                   ("✓ 過去の失敗事例チェック", False, MCK_DARK, 8),
                   ("✓ 明らかなバグ・タイポ", False, MCK_DARK, 8),
                   ("✓ コミットメッセージ規約", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.35)

    # Stage 2
    s2_top = s1_top + 1500000
    add_rect(slide, x, s2_top, col_w, 1400000,
             fill=WHITE, border=MCK_NAVY)
    add_text(slide, x + 100000, s2_top + 80000, col_w - 200000, 250000,
             "2. 人間レビュー", size=10, bold=True, color=MCK_NAVY)
    add_multiline(slide, x + 100000, s2_top + 360000, col_w - 200000, 1000000,
                  [("✓ アプローチの妥当性", False, MCK_DARK, 8),
                   ("✓ 業務ロジックの意図", False, MCK_DARK, 8),
                   ("✓ もっと良いやり方", False, MCK_DARK, 8),
                   ("✓ 影響範囲の見落とし", False, MCK_DARK, 8),
                   ("✓ 命名・可読性", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.35)

    # Col 2: Parallel development
    x = MARGIN_L + col_w + 150000
    add_rect(slide, x, top, col_w, 300000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 220000,
             "並列開発  問題なし", size=10, bold=True, color=WHITE)

    # A
    a_top = top + 400000
    add_rect(slide, x, a_top, col_w, 800000,
             fill=MCK_LIGHT_BLUE, border=MCK_BLUE)
    add_text(slide, x + 100000, a_top + 80000, col_w - 200000, 250000,
             "👤 Aさん", size=10, bold=True, color=MCK_NAVY)
    add_multiline(slide, x + 100000, a_top + 330000, col_w - 200000, 500000,
                  [("feature/issue-42", False, MCK_DARK, 8),
                   ("独立 Claude セッション", False, MCK_GRAY, 8),
                   ("PR #50", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.3)

    # B
    b_top = a_top + 900000
    add_rect(slide, x, b_top, col_w, 800000,
             fill=MCK_BG_LIGHT, border=MCK_BLUE)
    add_text(slide, x + 100000, b_top + 80000, col_w - 200000, 250000,
             "👤 Bさん", size=10, bold=True, color=MCK_NAVY)
    add_multiline(slide, x + 100000, b_top + 330000, col_w - 200000, 500000,
                  [("feature/issue-43", False, MCK_DARK, 8),
                   ("独立 Claude セッション", False, MCK_GRAY, 8),
                   ("PR #51", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.3)

    # Both → Gitea
    g_top = b_top + 900000
    add_rect(slide, x, g_top, col_w, 500000, fill=MCK_NAVY)
    add_text(slide, x, g_top + 50000, col_w, 250000,
             "🏭 Gitea", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, x, g_top + 280000, col_w, 200000,
             "PR #50 と #51 が並列で存在", size=8, color=MCK_LIGHT_BLUE,
             align=PP_ALIGN.CENTER)

    # Col 3: Merge conflicts
    x = MARGIN_L + (col_w + 150000) * 2
    add_rect(slide, x, top, col_w, 300000, fill=MCK_NAVY)
    add_text(slide, x + 100000, top + 50000, col_w - 200000, 220000,
             "マージ衝突  8割は解決できる", size=10, bold=True, color=WHITE)

    # Good
    g_top = top + 400000
    add_rect(slide, x, g_top, col_w, 1400000,
             fill=WHITE, border=MCK_GREEN)
    add_text(slide, x + 100000, g_top + 80000, col_w - 200000, 250000,
             "✓ 秘書が解決できる", size=10, bold=True, color=MCK_GREEN)
    add_multiline(slide, x + 100000, g_top + 360000, col_w - 200000, 1000000,
                  [("● 両方の意図を汲んだ統合", False, MCK_DARK, 8),
                   ("● 行追加系の単純衝突", False, MCK_DARK, 8),
                   ("● 指示しての解決", False, MCK_DARK, 8),
                   ("", False),
                   ("例：AはAPI追加、Bはリファクタ", False, MCK_GRAY, 7),
                   ("→ 両方取り込める", False, MCK_GRAY, 7)],
                  size=8, line_spacing=1.35)

    # Bad
    b_top = g_top + 1500000
    add_rect(slide, x, b_top, col_w, 1400000,
             fill=WHITE, border=MCK_RED)
    add_text(slide, x + 100000, b_top + 80000, col_w - 200000, 250000,
             "⚠ 人間の判断が必要", size=10, bold=True, color=MCK_RED)
    add_multiline(slide, x + 100000, b_top + 360000, col_w - 200000, 1000000,
                  [("● 設計判断の矛盾", False, MCK_DARK, 8),
                   ("● セマンティック衝突", False, MCK_DARK, 8),
                   ("● CLAUDE.md 自体の衝突", False, MCK_DARK, 8),
                   ("", False),
                   ("→ チームミーティングで合意", False, MCK_GRAY, 7)],
                  size=8, line_spacing=1.35)

    add_footer(slide, page, total)
    return slide


def s15_claude_md_growth(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "CLAUDE.md は当番制で育て、迷ったら置き場所の判断フローで判定する",
        tracker="CLAUDE.md の育て方")

    top = CONTENT_TOP + 100000

    # Top message (compressed)
    add_rect(slide, MARGIN_L, top, CONTENT_W, 320000,
             fill=MCK_BG_LIGHT)
    add_text(slide, MARGIN_L, top + 60000, CONTENT_W, 200000,
             "「気づいた人が書く」だけでは続かない。仕組みと判断基準の両方が必要",
             size=11, bold=True, color=MCK_NAVY, align=PP_ALIGN.CENTER)

    # 2 columns
    col_top = top + 420000
    col_w = (CONTENT_W - 200000) // 2

    # ==== Left column: 育てる3つの仕組み ====
    lx = MARGIN_L
    add_rect(slide, lx, col_top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, lx + 100000, col_top + 50000, col_w - 200000, 220000,
             "育てる3つの仕組み   HOW", size=10, bold=True, color=WHITE)

    mechanisms = [
        ("1", "週1の振り返り定例", "15分／週", "ハマりどころ共有 → 当番が PR 化", MCK_NAVY),
        ("2", "その場で小さな PR", "随時", "気づいた人が 3 行追記、チームがレビュー", MCK_BLUE),
        ("3", "月1の刈り込み", "月1回", "古いルール削除 → 肥大化防止", MCK_ORANGE),
    ]
    m_top = col_top + 380000
    m_h = 760000
    for i, (num, title, when, desc, color) in enumerate(mechanisms):
        y = m_top + i * (m_h + 40000)
        add_rect(slide, lx, y, col_w, m_h,
                 fill=WHITE, border=color, border_pt=1)
        # Number bar
        add_rect(slide, lx, y, 320000, m_h, fill=color)
        add_text(slide, lx, y, 320000, m_h, num,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(slide, lx + 380000, y + 80000, col_w - 550000, 260000,
                 title, size=11, bold=True, color=color)
        # When
        add_text(slide, lx + 380000, y + 320000, col_w - 550000, 200000,
                 when, size=8, color=MCK_GRAY)
        # Desc
        add_text(slide, lx + 380000, y + 510000, col_w - 550000, 240000,
                 desc, size=9, color=MCK_DARK)

    # ==== Right column: 迷ったら？ ====
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, col_top, col_w, 280000, fill=MCK_BLUE)
    add_text(slide, rx + 100000, col_top + 50000, col_w - 200000, 220000,
             "迷ったら？   置き場所の判断   WHAT", size=10, bold=True, color=WHITE)

    # 4 destinations table
    destinations = [
        ("自分の作業スタイル・好み", "~/.claude/CLAUDE.md", "個人メモ", MCK_GRAY),
        ("このプロジェクト固有の話", "プロジェクト CLAUDE.md", "プロジェクト層", MCK_BLUE),
        ("複数プロジェクトに効く", "team-standards CLAUDE.md", "共通層", MCK_NAVY),
        ("2回以上やる手順（10行超）", "team-standards/skills/", "Skill", MCK_GREEN),
    ]
    d_top = col_top + 380000
    header_h = 280000
    row_h = 380000
    # Header
    add_rect(slide, rx, d_top, col_w, header_h, fill=MCK_BG_LIGHT)
    add_text(slide, rx + 100000, d_top + 50000, 2400000, 200000,
             "性質", size=9, bold=True, color=MCK_GRAY)
    add_text(slide, rx + 2550000, d_top + 50000, col_w - 2650000, 200000,
             "行き先", size=9, bold=True, color=MCK_GRAY)
    # Rows
    for i, (prop, dest, label, color) in enumerate(destinations):
        y = d_top + header_h + i * row_h
        if i % 2 == 0:
            add_rect(slide, rx, y, col_w, row_h, fill=WHITE)
        else:
            add_rect(slide, rx, y, col_w, row_h, fill=MCK_BG_LIGHT)
        add_text(slide, rx + 100000, y + 60000, 2400000, 260000,
                 prop, size=9, color=MCK_DARK)
        add_text(slide, rx + 2550000, y + 40000, col_w - 2650000, 220000,
                 dest, size=9, bold=True, color=color)
        add_text(slide, rx + 2550000, y + 200000, col_w - 2650000, 180000,
                 label, size=7, color=MCK_GRAY)

    # Bottom guidance
    g_top = d_top + header_h + 4 * row_h + 40000
    add_rect(slide, rx, g_top, col_w, 440000,
             fill=MCK_LIGHT_BLUE)
    add_multiline(slide, rx + 150000, g_top + 50000, col_w - 300000, 360000,
                  [("● 迷ったらプロジェクト側。後から共通に昇格", False, MCK_DARK, 8),
                   ("● それでも迷ったら「⚠ 判断依頼」ラベルで小さな PR", False, MCK_DARK, 8),
                   ("● 週1定例で合意形成 → 却下されても損失ゼロ", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.3)

    add_footer(slide, page, total)
    return slide


def s15b_skills_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "必要な機能の80%は既存の公式/メジャー Skill でカバーできる — 自作はほぼ不要",
        tracker="必要な Skills / 全体像")

    top = CONTENT_TOP + 100000

    # Top message
    add_rect(slide, MARGIN_L, top, CONTENT_W, 400000,
             fill=MCK_BG_LIGHT)
    add_multiline(slide, MARGIN_L + 300000, top + 60000, CONTENT_W - 600000, 300000,
                  [("調査対象: 13 シナリオ", False, MCK_GRAY, 10),
                   ("→ 7 つはそのまま / 3 つは拡張 / 3 つは自作（合計 13 中 10 は既存活用）", True, MCK_NAVY, 11)],
                  size=11, align=PP_ALIGN.CENTER, line_spacing=1.3)

    # 3 layers as horizontal stack
    layer_top = top + 550000
    layer_h = 2900000
    col_w = (CONTENT_W - 300000) // 3

    layers = [
        ("レイヤー 1", "そのまま使える", "10 / 13 機能",
         "公式プラグイン\n+ 超メジャーな Skill\n+ Gitea 公式 MCP",
         "superpowers, code-review,\nclaude-md-management,\nfeature-dev, gitea-mcp,\nVoltAgent subagents",
         MCK_NAVY, "80%"),
        ("レイヤー 2", "既存を拡張", "3 / 13 機能",
         "ベースは既存 Skill\n数時間〜半日の追加作業",
         "壁打ちログ 4項目抽出\nセルフレビュー拡張\nCLAUDE.md 月1刈り込み",
         MCK_BLUE, "15%"),
        ("レイヤー 3", "完全に自作", "3 / 13 機能",
         "既存のものがない\n低〜中難易度で自作",
         "/post-to-gitea-issue\nteam-standards 初期設定\nGitea 自動レビュー CI",
         MCK_ORANGE, "5%"),
    ]

    for i, (num, title, count, approach, examples, color, pct) in enumerate(layers):
        x = MARGIN_L + i * (col_w + 150000)
        # Card
        add_rect(slide, x, layer_top, col_w, layer_h,
                 fill=WHITE, border=color, border_pt=1.5)
        # Top band
        add_rect(slide, x, layer_top, col_w, 80000, fill=color)
        # Number
        add_text(slide, x + 100000, layer_top + 150000, col_w - 200000, 260000,
                 num, size=10, bold=True, color=MCK_GRAY)
        # Title
        add_text(slide, x + 100000, layer_top + 420000, col_w - 200000, 400000,
                 title, size=16, bold=True, color=color)
        # Pct
        add_text(slide, x + 100000, layer_top + 830000, col_w - 200000, 400000,
                 pct, size=32, bold=True, color=color)
        # Count
        add_text(slide, x + 100000, layer_top + 1250000, col_w - 200000, 250000,
                 count, size=10, color=MCK_GRAY)
        # Divider
        add_line(slide, x + 100000, layer_top + 1520000,
                 x + col_w - 100000, layer_top + 1520000,
                 color=MCK_LIGHT_GRAY)
        # Approach
        add_multiline(slide, x + 100000, layer_top + 1560000, col_w - 200000, 500000,
                      approach.split("\n"), size=9, color=MCK_DARK, line_spacing=1.35)
        # Examples
        add_text(slide, x + 100000, layer_top + 2100000, col_w - 200000, 200000,
                 "例", size=8, bold=True, color=MCK_GRAY)
        add_multiline(slide, x + 100000, layer_top + 2300000, col_w - 200000, 550000,
                      examples.split("\n"), size=8, color=MCK_DARK, line_spacing=1.35)

    add_footer(slide, page, total)
    return slide


def s15c_skills_layer1(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "レイヤー1：Anthropic 公式プラグインとコミュニティ決定版 Skill をまず入れる",
        tracker="必要な Skills / レイヤー1 そのまま使える")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # ==== Left: Official ====
    lx = MARGIN_L
    add_rect(slide, lx, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, lx + 100000, top + 50000, col_w - 200000, 220000,
             "Anthropic 公式プラグイン", size=10, bold=True, color=WHITE)

    official = [
        ("superpowers", "統合ワークフロー集  (壁打ち・TDD・並列レビュー・デバッグ)", "★必須"),
        ("claude-md-management", "CLAUDE.md 品質監査  (/revise-claude-md)", "★必須"),
        ("code-review", "5並列エージェント PR レビュー  (信頼度フィルタ)", "★必須"),
        ("feature-dev", "7段階機能開発フロー", "推奨"),
        ("pr-review-toolkit", "6種類の専門レビュー  (コメント/テスト/型/品質)", "推奨"),
        ("commit-commands", "Git 自動化  (/commit, /commit-push-pr)", "推奨"),
        ("hookify", "カスタム Hook 作成", "必要時"),
        ("skill-creator", "自作 Skill のテンプレ生成", "必要時"),
    ]
    o_top = top + 380000
    for i, (name, desc, tag) in enumerate(official):
        y = o_top + i * 340000
        if i % 2 == 0:
            add_rect(slide, lx, y, col_w, 320000, fill=MCK_BG_LIGHT)
        add_text(slide, lx + 100000, y + 30000, 1700000, 260000,
                 name, size=9, bold=True, color=MCK_NAVY)
        add_text(slide, lx + 1850000, y + 30000, col_w - 2500000, 260000,
                 desc, size=8, color=MCK_DARK)
        tag_color = MCK_RED if "必須" in tag else (MCK_BLUE if "推奨" in tag else MCK_GRAY)
        add_text(slide, lx + col_w - 550000, y + 30000, 450000, 260000,
                 tag, size=8, bold=True, color=tag_color,
                 align=PP_ALIGN.RIGHT)

    # ==== Right: Community ====
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_BLUE)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "コミュニティ決定版 Skill", size=10, bold=True, color=WHITE)

    # Featured boxes
    community = [
        ("obra/superpowers", "144k stars", "統合スキル集の決定版", "★必須"),
        ("VoltAgent/awesome-claude-code-subagents", "130+ subagents", "security/architect/test-coverage 等", "★必須"),
        ("gitea/gitea-mcp", "Gitea 公式", "Issue/PR/コード操作の MCP", "★必須"),
        ("hesreallyhim/awesome-claude-code", "キュレーション", "リポジトリ集の決定版", "参考"),
        ("claude-config-doctor", "CLAUDE.md ツール", "複数ファイルの conflict 検出", "参考"),
    ]
    c_top = top + 380000
    for i, (name, meta, desc, tag) in enumerate(community):
        y = c_top + i * 560000
        add_rect(slide, rx, y, col_w, 520000,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_text(slide, rx + 100000, y + 60000, col_w - 700000, 260000,
                 name, size=9, bold=True, color=MCK_NAVY)
        tag_color = MCK_RED if "必須" in tag else MCK_GRAY
        add_text(slide, rx + col_w - 550000, y + 60000, 450000, 260000,
                 tag, size=8, bold=True, color=tag_color,
                 align=PP_ALIGN.RIGHT)
        add_text(slide, rx + 100000, y + 290000, col_w - 200000, 220000,
                 f"{meta}  ·  {desc}", size=8, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s15d_skills_layer23(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "レイヤー2と3：既存を少し拡張すれば3つ、ゼロから作るのは3つだけ",
        tracker="必要な Skills / レイヤー2+3")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # ==== Left: Layer 2 - Customize ====
    lx = MARGIN_L
    add_rect(slide, lx, top, col_w, 280000, fill=MCK_BLUE)
    add_text(slide, lx + 100000, top + 50000, col_w - 200000, 220000,
             "レイヤー 2   既存を拡張する", size=10, bold=True, color=WHITE)

    layer2 = [
        ("壁打ちログ → 4項目抽出",
         "ベース: superpowers:brainstorming\n       + Meeting Notes Architect",
         "「前提・選択肢・結論・却下理由」\nの抽出ロジックを追加",
         "4-8 時間"),
        ("セルフレビュー",
         "ベース: Git Diff Reviewer",
         "CLAUDE.md 準拠チェック・\nテスト存在確認を追加",
         "4-6 時間"),
        ("CLAUDE.md 月1刈り込み",
         "ベース: claude-md-management\n       + Auto Dream",
         "200行上限監視と月1トリガー、\n削除判定ロジック",
         "6-10 時間"),
    ]
    l2_top = top + 380000
    l2_h = 1000000
    for i, (title, base, what, time) in enumerate(layer2):
        y = l2_top + i * (l2_h + 30000)
        add_rect(slide, lx, y, col_w, l2_h,
                 fill=WHITE, border=MCK_BLUE, border_pt=1)
        add_rect(slide, lx, y, 40000, l2_h, fill=MCK_BLUE)
        add_text(slide, lx + 100000, y + 60000, col_w - 200000, 240000,
                 title, size=10, bold=True, color=MCK_NAVY)
        add_multiline(slide, lx + 100000, y + 300000, col_w - 200000, 400000,
                      base.split("\n"), size=7, color=MCK_GRAY, line_spacing=1.3)
        add_multiline(slide, lx + 100000, y + 650000, col_w - 700000, 300000,
                      what.split("\n"), size=8, color=MCK_DARK, line_spacing=1.3)
        add_text(slide, lx + col_w - 600000, y + 700000, 550000, 220000,
                 time, size=8, bold=True, color=MCK_ORANGE,
                 align=PP_ALIGN.RIGHT)

    # ==== Right: Layer 3 - Build ====
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_ORANGE)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "レイヤー 3   ゼロから自作", size=10, bold=True, color=WHITE)

    layer3 = [
        ("/post-to-gitea-issue",
         "壁打ちログ要約を Issue にコメント投稿",
         "Gitea MCP 経由で\ncreate_issue_comment を呼ぶ",
         "2-4 時間",
         "参考: markwylde/claude-code-gitea-action"),
        ("team-standards 初期設定",
         "共有リポジトリの雛形を作る",
         "CLAUDE.md 初版\nSkills/agents/hooks/mcp の構造",
         "半日",
         "参考: anthropics/claude-plugins-official"),
        ("Gitea 自動レビュー CI",
         "PR ラベルで Claude が自動レビュー",
         "Gitea Actions で CI 構築\n※Phase 3 以降",
         "1-2 日",
         "参考: markwylde/claude-code-gitea-action"),
    ]
    l3_top = top + 380000
    l3_h = 1000000
    for i, (title, purpose, how, time, ref) in enumerate(layer3):
        y = l3_top + i * (l3_h + 30000)
        add_rect(slide, rx, y, col_w, l3_h,
                 fill=WHITE, border=MCK_ORANGE, border_pt=1)
        add_rect(slide, rx, y, 40000, l3_h, fill=MCK_ORANGE)
        add_text(slide, rx + 100000, y + 60000, col_w - 700000, 240000,
                 title, size=10, bold=True, color=MCK_NAVY)
        add_text(slide, rx + col_w - 600000, y + 60000, 550000, 240000,
                 time, size=8, bold=True, color=MCK_ORANGE,
                 align=PP_ALIGN.RIGHT)
        add_text(slide, rx + 100000, y + 300000, col_w - 200000, 220000,
                 purpose, size=8, color=MCK_DARK)
        add_multiline(slide, rx + 100000, y + 520000, col_w - 200000, 300000,
                      how.split("\n"), size=8, color=MCK_GRAY, line_spacing=1.3)
        add_text(slide, rx + 100000, y + 820000, col_w - 200000, 180000,
                 ref, size=7, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s15e_install_order(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "導入順序：Day 1 に必須セット → Week 1 で強化 → Month 1 で仕上げ",
        tracker="必要な Skills / 導入順序")

    top = CONTENT_TOP + 100000
    phase_w = (CONTENT_W - 200000) // 3
    phase_h = 3500000

    phases = [
        ("PHASE 1", "Day 1", "必須セット", "2時間以内",
         MCK_NAVY,
         [
             ("1.", "Gitea MCP インストール", "各自 Token 発行"),
             ("", "", "claude mcp add ..."),
             ("", "", ""),
             ("2.", "公式プラグイン", "/plugin install"),
             ("", "", "  superpowers"),
             ("", "", "  claude-md-management"),
             ("", "", ""),
             ("3.", "team-standards 作成", "Gitea にリポジトリ"),
             ("", "", "CLAUDE.md 初版を"),
             ("", "", "3名で合意"),
         ]),
        ("PHASE 2", "Week 1", "強化", "数日〜1週間",
         MCK_BLUE,
         [
             ("4.", "追加プラグイン", "/plugin install"),
             ("", "", "  code-review"),
             ("", "", "  feature-dev"),
             ("", "", ""),
             ("5.", "VoltAgent subagents", "必要なものを"),
             ("", "", "team-standards/agents/"),
             ("", "", "にコピー"),
             ("", "", ""),
             ("6.", "/post-to-gitea-issue", "自作 slash command"),
             ("", "", "2-4 時間"),
         ]),
        ("PHASE 3", "Month 1", "仕上げ", "余裕があれば",
         MCK_GREEN,
         [
             ("7.", "セルフレビュー拡張", "Git Diff Reviewer"),
             ("", "", "+ CLAUDE.md 準拠"),
             ("", "", ""),
             ("8.", "CLAUDE.md 刈り込み", "月1トリガー"),
             ("", "", "Hook で自動化"),
             ("", "", ""),
             ("9.", "Gitea CI 連携", "※任意"),
             ("", "", "claude-code-gitea-"),
             ("", "", "action を参考に"),
             ("", "", ""),
         ]),
    ]

    for i, (phase, when, label, duration, color, items) in enumerate(phases):
        x = MARGIN_L + i * (phase_w + 100000)
        # Header
        add_rect(slide, x, top, phase_w, 600000, fill=color)
        add_text(slide, x + 100000, top + 50000, phase_w - 200000, 200000,
                 phase, size=8, bold=True, color=RGBColor(200, 215, 230))
        add_text(slide, x + 100000, top + 230000, phase_w - 200000, 250000,
                 label, size=14, bold=True, color=WHITE)
        add_text(slide, x + 100000, top + 450000, phase_w - 200000, 180000,
                 f"{when}  ·  {duration}", size=8, color=RGBColor(200, 215, 230))
        # Body
        body_top = top + 600000
        body_h = phase_h - 600000
        add_rect(slide, x, body_top, phase_w, body_h,
                 fill=WHITE, border=color, border_pt=1)
        # Items
        for j, (num, title, desc) in enumerate(items):
            y = body_top + 60000 + j * 230000
            if num:
                add_text(slide, x + 80000, y, 300000, 200000,
                         num, size=9, bold=True, color=color)
            if title:
                add_text(slide, x + 380000, y, phase_w - 460000, 200000,
                         title, size=9, bold=True, color=MCK_DARK)
            if desc:
                add_text(slide, x + 380000, y, phase_w - 460000, 200000,
                         desc, size=8, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s16_pilot_plan(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "3名×3ヶ月のパイロットで立ち上げ、成功したら残り7名に展開する",
        tracker="進め方の提案")

    top = CONTENT_TOP + 100000

    # Why not all
    add_rect(slide, MARGIN_L, top, CONTENT_W, 550000,
             fill=MCK_NAVY)
    add_text(slide, MARGIN_L, top + 80000, CONTENT_W, 280000,
             "いきなり全員ではなく、3名で3ヶ月試す",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_L, top + 320000, CONTENT_W, 200000,
             "失敗コスト 1/3  |  ワークフローを磨ける  |  上司を説得しやすい  |  成功事例を作れる",
             size=9, color=MCK_LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # 3-month timeline
    tl_top = top + 700000
    tl_h = 2500000
    phase_w = (CONTENT_W - 200000) // 3

    phases = [
        ("1ヶ月目", "準備と慣らし", MCK_NAVY, [
            "team-standards リポジトリ作成",
            "CLAUDE.md 初版を3名で書く",
            "小さな Issue で全フロー試す",
            "4項目テンプレを使い始める",
        ]),
        ("2ヶ月目", "本格運用", MCK_BLUE, [
            "3名の実作業をフル運用",
            "Skills・Subagents を追加",
            "週1振り返りで CLAUDE.md 更新",
            "当番制を徹底する",
        ]),
        ("3ヶ月目", "チェックと展開判断", MCK_GREEN, [
            "3ヶ月目チェック項目を確認",
            "推進役の2週間休暇テスト",
            "成果と課題を整理",
            "残り7名への展開を判断",
        ]),
    ]
    for i, (month, phase, color, tasks) in enumerate(phases):
        x = MARGIN_L + i * (phase_w + 100000)
        # Header
        add_rect(slide, x, tl_top, phase_w, 450000, fill=color)
        add_text(slide, x + 100000, tl_top + 50000, phase_w - 200000, 200000,
                 month, size=9, bold=True, color=WHITE)
        add_text(slide, x + 100000, tl_top + 220000, phase_w - 200000, 220000,
                 phase, size=12, bold=True, color=WHITE)
        # Body
        add_rect(slide, x, tl_top + 450000, phase_w, tl_h - 450000,
                 fill=WHITE, border=color, border_pt=1)
        tasks_tuples = [("● " + t, False, MCK_DARK, 9) for t in tasks]
        add_multiline(slide, x + 150000, tl_top + 600000,
                      phase_w - 300000, tl_h - 700000,
                      tasks_tuples, size=9, line_spacing=1.5)

    add_footer(slide, page, total)
    return slide


def s17_checkpoints_trap(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "撤退条件と推進役依存の罠対策を、最初から設計に組み込む",
        tracker="リスク対策と撤退基準")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # Left: Checkpoints
    add_rect(slide, MARGIN_L, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 100000, top + 50000, col_w - 200000, 220000,
             "チェックポイントと撤退条件", size=10, bold=True, color=WHITE)

    # 3-month check
    c3_top = top + 380000
    add_rect(slide, MARGIN_L, c3_top, col_w, 1400000,
             fill=WHITE, border=MCK_BLUE)
    add_text(slide, MARGIN_L + 100000, c3_top + 80000, col_w - 200000, 250000,
             "3ヶ月目チェック", size=10, bold=True, color=MCK_BLUE)
    add_multiline(slide, MARGIN_L + 100000, c3_top + 340000,
                  col_w - 200000, 1000000,
                  [("□ 推進役以外の CLAUDE.md 更新 PR が5件以上", False, MCK_DARK, 8),
                   ("□ 壁打ちログ4項目が週2件以上蓄積", False, MCK_DARK, 8),
                   ("□ CODEOWNERS 2人とも機能", False, MCK_DARK, 8),
                   ("□ 参加3名が『慣れた』と実感", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.4)

    # 6-month check
    c6_top = c3_top + 1500000
    add_rect(slide, MARGIN_L, c6_top, col_w, 1400000,
             fill=WHITE, border=MCK_ORANGE)
    add_text(slide, MARGIN_L + 100000, c6_top + 80000, col_w - 200000, 250000,
             "6ヶ月目チェック", size=10, bold=True, color=MCK_ORANGE)
    add_multiline(slide, MARGIN_L + 100000, c6_top + 340000,
                  col_w - 200000, 1000000,
                  [("□ 2週間休暇テストで回るか検証", False, MCK_DARK, 8),
                   ("□ 集合知化の実感（アンケート）", False, MCK_DARK, 8),
                   ("□ 定量指標（PR数・レビュー時間）", False, MCK_DARK, 8),
                   ("□ 残り7名への展開判断", False, MCK_DARK, 8)],
                  size=8, line_spacing=1.4)

    # Right: Trap & countermeasures
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_RED)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "推進役依存の罠", size=10, bold=True, color=WHITE)

    # Trap description
    t_top = top + 380000
    add_rect(slide, rx, t_top, col_w, 1000000,
             fill=RGBColor(254, 243, 243), border=MCK_RED)
    add_text(slide, rx + 100000, t_top + 80000, col_w - 200000, 250000,
             "⚠ 最大の失敗シナリオ", size=9, bold=True, color=MCK_RED)
    add_multiline(slide, rx + 100000, t_top + 340000,
                  col_w - 200000, 600000,
                  [("推進役Aさんだけが CLAUDE.md を更新、", False, MCK_DARK, 8),
                   ("他は『Aさんの秘書に聞く消費者』に。", False, MCK_DARK, 8),
                   ("Aさんの異動で全員立ち尽くす。", True, MCK_RED, 8)],
                  size=8, line_spacing=1.35)

    # Countermeasures
    cm_top = t_top + 1100000
    add_text(slide, rx, cm_top, col_w, 280000,
             "4つの対策", size=10, bold=True, color=MCK_NAVY)
    measures = [
        ("輪番制を徹底", "更新・議事・レビューを週替わり"),
        ("CODEOWNERS 2人体制", "team-standards のレビュアーを必ず2人"),
        ("属人化 KPI", "推進役以外のPR比率を定期確認"),
        ("2週間休暇テスト", "3ヶ月目に意図的に休ませて検証"),
    ]
    for i, (title, desc) in enumerate(measures):
        y = cm_top + 320000 + i * 360000
        add_text(slide, rx + 50000, y, 280000, 220000,
                 f"{i + 1}.", size=9, bold=True, color=MCK_NAVY)
        add_text(slide, rx + 380000, y, col_w - 500000, 220000,
                 title, size=9, bold=True, color=MCK_DARK)
        add_text(slide, rx + 380000, y + 180000, col_w - 500000, 200000,
                 desc, size=8, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s18_first_week_and_questions(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "合意できたら来週から試運転を始める — 並行してチームの率直な意見を聞きたい",
        tracker="次の一歩と Q&A")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2

    # Left: First week
    add_rect(slide, MARGIN_L, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 100000, top + 50000, col_w - 200000, 220000,
             "最初の1週間でやること", size=10, bold=True, color=WHITE)

    week_items = [
        ("Day 1-2", "準備", [
            "team-standards リポジトリ作成",
            "CLAUDE.md 初版を3名で書く",
            "~/dev/team-standards に clone",
        ], MCK_NAVY),
        ("Day 3-4", "試運転", [
            "小さな Issue で全フロー試す",
            "壁打ち → 4項目テンプレ → PR",
            "つまずき記録",
        ], MCK_BLUE),
        ("Day 5", "振り返り", [
            "3名で15分振り返り",
            "良い点・悪い点整理",
            "CLAUDE.md を早速更新",
        ], MCK_GREEN),
    ]
    it_top = top + 380000
    it_h = 1050000
    for i, (days, phase, tasks, color) in enumerate(week_items):
        y = it_top + i * (it_h + 30000)
        add_rect(slide, MARGIN_L, y, col_w, it_h,
                 fill=WHITE, border=color)
        # Label box
        add_rect(slide, MARGIN_L, y, 1100000, it_h, fill=color)
        add_text(slide, MARGIN_L + 50000, y + 200000, 1000000, 240000,
                 days, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, MARGIN_L + 50000, y + 500000, 1000000, 240000,
                 phase, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Tasks
        tasks_tuples = [("● " + t, False, MCK_DARK, 9) for t in tasks]
        add_multiline(slide, MARGIN_L + 1200000, y + 150000,
                      col_w - 1350000, it_h - 300000,
                      tasks_tuples, size=9, line_spacing=1.35)

    # Right: Questions to team
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_BLUE)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "みんなに聞きたいこと", size=10, bold=True, color=WHITE)

    questions = [
        ("この方向性、ピンと来る？", "『集合知化』という旗に心が動く？"),
        ("最初の3名、誰がやる？", "興味がある人、手を挙げてくれる人"),
        ("不安・懸念は？", "『これは続かない』『ここが面倒』率直に"),
        ("もっと良い方法ある？", "別のアプローチ、気になる案があれば"),
        ("上司への相談、どう進める？", "タイミング・伝え方・必要な準備"),
    ]
    q_top = top + 380000
    for i, (q, desc) in enumerate(questions):
        y = q_top + i * 600000
        add_rect(slide, rx, y, col_w, 550000,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_text(slide, rx + 50000, y + 80000, 350000, 300000,
                 "?", size=20, bold=True, color=MCK_BLUE)
        add_text(slide, rx + 450000, y + 100000, col_w - 550000, 260000,
                 q, size=10, bold=True, color=MCK_DARK)
        add_text(slide, rx + 450000, y + 320000, col_w - 550000, 220000,
                 desc, size=8, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s19_end(prs, page, total):
    slide = new_slide(prs)
    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill=MCK_BLUE)

    # Center content
    add_text(slide, MARGIN_L, 1800000, CONTENT_W, 300000,
             "DX PORTAL TEAM", size=9, bold=True, color=MCK_BLUE,
             align=PP_ALIGN.CENTER)
    add_line(slide, (SLIDE_W - 1000000) // 2, 2200000,
             (SLIDE_W + 1000000) // 2, 2200000,
             color=MCK_NAVY, weight=2)
    add_text(slide, MARGIN_L, 2300000, CONTENT_W, 600000,
             "ありがとうございました", size=28, bold=True, color=MCK_NAVY,
             align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_L, 2950000, CONTENT_W, 400000,
             "意見・質問、お待ちしています",
             size=14, color=MCK_GRAY, align=PP_ALIGN.CENTER)
    return slide


# ==========================================
# MAIN
# ==========================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        s01_title,
        s02_executive_summary,
        s03_current_state,
        s04_risks,
        s05_goal,
        s06_vegapunk_and_why_gitea,
        s07_direction,
        s08_day_in_life,
        s09_chats_planning,
        s10_chats_execution,
        s11_team_standards,
        s12_components_detail,
        s13_hooks_mcp_and_rules,
        s14_review_parallel_conflict,
        s15_claude_md_growth,
        s15b_skills_overview,
        s15c_skills_layer1,
        s15d_skills_layer23,
        s15e_install_order,
        s16_pilot_plan,
        s17_checkpoints_trap,
        s18_first_week_and_questions,
        s19_end,
    ]

    total = len(builders)
    for i, build in enumerate(builders, 1):
        build(prs, i, total)
        print(f"  Built {i:2d}/{total}: {build.__name__}")

    out_path = Path(__file__).parent / "team-collective-intelligence-v2.pptx"
    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")
    print(f"Total: {total} slides")


if __name__ == "__main__":
    main()
