"""
Build Claude Code components deep-dive slides (McKinsey-style).
Covers: CLAUDE.md / Skills / Subagents / Hooks / MCP
Output: .tmp/claude-code-components.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ==== Slide dimensions (16:9) ====
SLIDE_W = 9144000
SLIDE_H = 5143500

# ==== McKinsey-style palette ====
MCK_NAVY = RGBColor(26, 43, 76)
MCK_BLUE = RGBColor(74, 144, 226)
MCK_LIGHT_BLUE = RGBColor(208, 225, 245)
MCK_GRAY = RGBColor(138, 148, 166)
MCK_LIGHT_GRAY = RGBColor(230, 233, 239)
MCK_BG_LIGHT = RGBColor(246, 248, 251)
MCK_DARK = RGBColor(34, 34, 34)
MCK_RED = RGBColor(192, 48, 48)
MCK_RED_LIGHT = RGBColor(254, 243, 243)
MCK_GREEN = RGBColor(46, 125, 50)
MCK_ORANGE = RGBColor(214, 126, 30)
MCK_PURPLE = RGBColor(103, 58, 183)
WHITE = RGBColor(255, 255, 255)

# ==== Layout ====
MARGIN_L = 457200
MARGIN_R = 457200
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
ACTION_TITLE_Y = 380000
ACTION_TITLE_H = 500000
RULE_Y = 950000
CONTENT_TOP = 1050000
CONTENT_BOTTOM = 4650000
FOOTER_Y = 4800000
TRACKER_Y = 260000


# ==== Helpers ====

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def add_rect(slide, left, top, width, height, fill=None, border=None, border_pt=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border is not None:
        shape.line.color.rgb = border
        if border_pt:
            shape.line.width = Pt(border_pt)
        else:
            shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=None, weight=0.75):
    if color is None:
        color = MCK_LIGHT_GRAY
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    return conn


def add_text(slide, left, top, width, height, text, size=10, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="Meiryo UI"):
    if color is None:
        color = MCK_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 30000
    tf.margin_right = 30000
    tf.margin_top = 20000
    tf.margin_bottom = 20000
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_multiline(slide, left, top, width, height, lines, size=10,
                  color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.25, font_name="Meiryo UI"):
    if color is None:
        color = MCK_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 30000
    tf.margin_right = 30000
    tf.margin_top = 20000
    tf.margin_bottom = 20000
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        if isinstance(line, tuple):
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            c = line[2] if len(line) > 2 and line[2] is not None else color
            s = line[3] if len(line) > 3 and line[3] is not None else size
        else:
            text = line
            bold = False
            c = color
            s = size
        if text == "":
            run = p.add_run()
            run.text = " "
            run.font.size = Pt(s)
            run.font.name = font_name
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(s)
        run.font.bold = bold
        run.font.color.rgb = c
        run.font.name = font_name
    return tb


def add_tracker(slide, label):
    add_text(slide, MARGIN_L, TRACKER_Y, CONTENT_W, 150000,
             label.upper(), size=8, bold=True, color=MCK_BLUE)


def add_action_title(slide, title, tracker=None):
    if tracker:
        add_tracker(slide, tracker)
    add_text(slide, MARGIN_L, ACTION_TITLE_Y, CONTENT_W, ACTION_TITLE_H,
             title, size=16, bold=True, color=MCK_NAVY)
    add_line(slide, MARGIN_L, RULE_Y, SLIDE_W - MARGIN_R, RULE_Y,
             color=MCK_NAVY, weight=1.5)


def add_footer(slide, page, total, source=None):
    if source:
        add_text(slide, MARGIN_L, FOOTER_Y, CONTENT_W - 500000, 200000,
                 f"Source: {source}", size=8, color=MCK_GRAY)
    add_text(slide, SLIDE_W - MARGIN_R - 800000, FOOTER_Y,
             800000, 200000,
             f"{page}", size=9, color=MCK_GRAY, align=PP_ALIGN.RIGHT)
    add_line(slide, MARGIN_L, FOOTER_Y - 30000,
             SLIDE_W - MARGIN_R, FOOTER_Y - 30000,
             color=MCK_LIGHT_GRAY, weight=0.5)


def add_code_block(slide, left, top, width, height, code_lines, size=8):
    add_rect(slide, left, top, width, height,
             fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY, border_pt=0.5)
    lines = [(line, False, MCK_DARK, size) if isinstance(line, str) else line
             for line in code_lines]
    add_multiline(slide, left + 100000, top + 80000,
                  width - 200000, height - 160000,
                  lines, size=size, line_spacing=1.2,
                  font_name="Consolas")


# ============================================
# SLIDES
# ============================================

def s01_title(prs, page, total):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill=MCK_BLUE)

    add_text(slide, MARGIN_L, 1700000, CONTENT_W, 300000,
             "DX PORTAL TEAM  |  TECHNICAL DEEP DIVE", size=9, bold=True, color=MCK_BLUE)
    add_text(slide, MARGIN_L, 2050000, CONTENT_W, 600000,
             "集合知化の道具立て", size=32, bold=True, color=MCK_NAVY)
    add_text(slide, MARGIN_L, 2750000, CONTENT_W, 500000,
             "Claude Code の5つの構成要素 ― CLAUDE.md / Skills / Subagents / Hooks / MCP",
             size=14, color=MCK_DARK)
    add_text(slide, MARGIN_L, 3250000, CONTENT_W, 400000,
             "~ 何を・何のために・どう使うか ~",
             size=11, color=MCK_GRAY)
    add_line(slide, MARGIN_L, 4600000, MARGIN_L + 1500000, 4600000,
             color=MCK_NAVY, weight=2)
    add_text(slide, MARGIN_L, 4650000, CONTENT_W, 250000,
             "集合知化ワークフローのテクニカルリファレンス", size=10, color=MCK_GRAY)
    return slide


def s02_executive_summary(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "5つの構成要素はそれぞれ役割が異なる ― 組み合わせて集合知化を実現する",
        tracker="Executive summary")

    box_top = CONTENT_TOP + 100000
    box_h = 3400000
    col_w = (CONTENT_W - 200000) // 5
    gap = 50000

    elements = [
        ("CLAUDE.md", "常識", "常に読まれる\n基本ルール",
         "● 規約・背景知識\n● 200行以内\n● PRで育てる",
         MCK_NAVY),
        ("Skills", "手順", "呼び出し式の\n作業手順書",
         "● /implement-feature\n● 長い手順はここ\n● 2回以上やるなら化",
         MCK_BLUE),
        ("Subagents", "専門家", "特定領域の\n専門秘書",
         "● security-auditor\n● 独立 context\n● ツール制限可",
         MCK_PURPLE),
        ("Hooks", "自動化", "イベントで\n自動実行",
         "● lint/format自動\n● 危険コマンド遮断\n● 21 イベント",
         MCK_ORANGE),
        ("MCP", "外部連携", "外部ツールを\n秘書の道具に",
         "● Gitea/DB/Playwright\n● ⚠️ 脆弱性注意\n● 閉域 OK 優先",
         MCK_GREEN),
    ]

    for i, (name, role, main, detail, color) in enumerate(elements):
        x = MARGIN_L + i * (col_w + gap)
        add_rect(slide, x, box_top, col_w, 70000, fill=color)
        add_rect(slide, x, box_top + 70000, col_w, box_h - 70000,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_text(slide, x + 100000, box_top + 150000, col_w - 200000, 320000,
                 name, size=12, bold=True, color=color)
        add_text(slide, x + 100000, box_top + 500000, col_w - 200000, 250000,
                 f"役割: {role}", size=9, color=MCK_GRAY)
        main_lines = [(l, True, MCK_NAVY, 10) for l in main.split("\n")]
        add_multiline(slide, x + 100000, box_top + 800000, col_w - 200000, 1100000,
                      main_lines, size=10, line_spacing=1.3)
        add_line(slide, x + 150000, box_top + 2000000,
                 x + col_w - 150000, box_top + 2000000,
                 color=MCK_LIGHT_GRAY)
        add_multiline(slide, x + 100000, box_top + 2100000, col_w - 200000, 1200000,
                      detail.split("\n"), size=8, color=MCK_DARK, line_spacing=1.4)

    add_footer(slide, page, total)
    return slide


def s03_relationship(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "5要素の関係：CLAUDE.md が『脳』、他4つはそこから派生する専門機能",
        tracker="5要素の関係")

    top = CONTENT_TOP + 100000

    # Center: CLAUDE.md box
    center_w = 2800000
    center_h = 700000
    center_x = (SLIDE_W - center_w) // 2
    center_y = top + 100000
    add_rect(slide, center_x, center_y, center_w, center_h,
             fill=MCK_NAVY)
    add_text(slide, center_x, center_y + 80000, center_w, 300000,
             "CLAUDE.md", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, center_x, center_y + 400000, center_w, 250000,
             "チームの常識・背景知識", size=10, color=MCK_LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # 4 derived components around it
    derivs = [
        ("Skills", "「常に読む」には\n長すぎる手順を切り出し",
         "CLAUDE.md → 詳細は Skills へ",
         MCK_BLUE),
        ("Subagents", "専門分野だけに\n特化した秘書",
         "ツール制限つきで独立動作",
         MCK_PURPLE),
        ("Hooks", "ルール違反を\n自動で検知・阻止",
         "CLAUDE.md は約束、Hooks は法律",
         MCK_ORANGE),
        ("MCP", "秘書の手を\n外部ツールまで伸ばす",
         "Gitea・DB・Playwright 等",
         MCK_GREEN),
    ]

    # Position: 4 boxes at corners
    box_w = 2000000
    box_h = 1200000
    positions = [
        (MARGIN_L + 300000, top + 1400000),                      # top-left
        (SLIDE_W - MARGIN_R - box_w - 300000, top + 1400000),    # top-right
        (MARGIN_L + 300000, top + 2800000),                      # bottom-left
        (SLIDE_W - MARGIN_R - box_w - 300000, top + 2800000),    # bottom-right
    ]

    for i, (name, desc, rel, color) in enumerate(derivs):
        x, y = positions[i]
        add_rect(slide, x, y, box_w, box_h,
                 fill=WHITE, border=color, border_pt=2)
        add_rect(slide, x, y, box_w, 60000, fill=color)
        add_text(slide, x + 100000, y + 130000, box_w - 200000, 300000,
                 name, size=13, bold=True, color=color)
        add_multiline(slide, x + 100000, y + 440000, box_w - 200000, 500000,
                      desc.split("\n"), size=9, color=MCK_DARK, line_spacing=1.3)
        add_text(slide, x + 100000, y + 950000, box_w - 200000, 200000,
                 rel, size=8, color=MCK_GRAY)

    # Connecting lines
    cx_center = center_x + center_w // 2
    cy_center = center_y + center_h // 2
    for i, (x, y) in enumerate(positions):
        bx = x + box_w // 2
        by = y + box_h // 2
        add_line(slide, cx_center, cy_center, bx, by,
                 color=MCK_LIGHT_GRAY, weight=0.75)

    # Bottom summary
    add_text(slide, MARGIN_L, 4400000, CONTENT_W, 250000,
             "つまり CLAUDE.md を中心に、Skills で広げ、Subagents で深め、Hooks で固め、MCP で外とつなぐ",
             size=11, bold=True, color=MCK_NAVY, align=PP_ALIGN.CENTER)

    add_footer(slide, page, total)
    return slide


def s04_claudemd_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "CLAUDE.md：秘書が常に読む『書かれた常識』― 200行以内で育てる",
        tracker="構成要素 1 / CLAUDE.md")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # Left: What and How
    add_rect(slide, MARGIN_L, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 100000, top + 50000, col_w - 200000, 220000,
             "CLAUDE.md とは", size=10, bold=True, color=WHITE)

    add_multiline(slide, MARGIN_L + 100000, top + 380000, col_w - 200000, 2800000,
                  [("● チームの秘書（Claude Code）が", False, MCK_DARK, 10),
                   ("  セッション開始時に自動で読み込むファイル", False, MCK_DARK, 10),
                   ("", False),
                   ("● ここに書いた内容は全員の秘書の知識に", False, MCK_DARK, 10),
                   ("", False),
                   ("● Git でバージョン管理、PR で育てる", False, MCK_DARK, 10),
                   ("", False),
                   ("● 200 行以内が目安（肥大化で効力低下）", False, MCK_DARK, 10),
                   ("", False),
                   ("── 書くこと ──", True, MCK_NAVY, 10),
                   ("● 規約・背景知識・開発の流れ", False, MCK_DARK, 10),
                   ("● 過去の失敗あるある", False, MCK_DARK, 10),
                   ("● よく使うコマンド・パス", False, MCK_DARK, 10),
                   ("", False),
                   ("── 書かないこと ──", True, MCK_NAVY, 10),
                   ("● 自明なこと（一般的ベストプラクティス）", False, MCK_GRAY, 10),
                   ("● 長い手順 → Skills に切り出す", False, MCK_GRAY, 10),
                   ("● 個人のメモ → ~/.claude/CLAUDE.md へ", False, MCK_GRAY, 10)],
                  size=10, line_spacing=1.35)

    # Right: Sample code
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_BLUE)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "初版サンプル（約30行の骨格）", size=10, bold=True, color=WHITE)

    code_lines = [
        "# team-standards CLAUDE.md",
        "",
        "## チーム文脈",
        "- チーム: DX ポータル開発",
        "- 技術: React + TS + Python",
        "",
        "## 開発の流れ",
        "1. 15 分以上は Issue を立てる",
        "2. 壁打ち → 4 項目テンプレで Issue",
        "3. Draft PR → セルフレビュー",
        "4. 人間レビュー → マージ",
        "5. 学びがあれば追記 PR",
        "",
        "## 秘書への基本指示",
        "- 過去の失敗事例を先に検索",
        "- 迷ったら選択肢を提示",
        "- テストは実装前に書く",
        "- any / unknown は使わない",
        "",
        "## コーディング規約",
        "<!-- 気づいたら追記。最初は空 -->",
        "",
        "## 過去の失敗あるある",
        "<!-- ハマったらその場で追記 -->",
    ]
    add_code_block(slide, rx, top + 380000, col_w, col_h - 400000,
                   code_lines, size=9)

    add_footer(slide, page, total)
    return slide


def s05_claudemd_growth(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "CLAUDE.md は当番制で育て、迷ったら置き場所の判断フローで判定する",
        tracker="構成要素 1 / CLAUDE.md の育て方")

    top = CONTENT_TOP + 100000

    # Top message
    add_rect(slide, MARGIN_L, top, CONTENT_W, 320000, fill=MCK_BG_LIGHT)
    add_text(slide, MARGIN_L, top + 60000, CONTENT_W, 200000,
             "「気づいた人が書く」だけでは続かない。仕組みと判断基準の両方が必要",
             size=11, bold=True, color=MCK_NAVY, align=PP_ALIGN.CENTER)

    col_top = top + 420000
    col_w = (CONTENT_W - 200000) // 2

    # Left: 育てる3つの仕組み
    lx = MARGIN_L
    add_rect(slide, lx, col_top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, lx + 100000, col_top + 50000, col_w - 200000, 220000,
             "育てる3つの仕組み  HOW", size=10, bold=True, color=WHITE)

    mechanisms = [
        ("1", "週1の振り返り定例", "15分／週", "ハマりどころ → 当番が PR 化", MCK_NAVY),
        ("2", "その場で小さな PR", "随時", "気づいた人が3行追記", MCK_BLUE),
        ("3", "月1の刈り込み", "月1回", "古いルール削除で肥大化防止", MCK_ORANGE),
    ]
    m_top = col_top + 380000
    m_h = 760000
    for i, (num, title, when, desc, color) in enumerate(mechanisms):
        y = m_top + i * (m_h + 40000)
        add_rect(slide, lx, y, col_w, m_h,
                 fill=WHITE, border=color, border_pt=1)
        add_rect(slide, lx, y, 320000, m_h, fill=color)
        add_text(slide, lx, y, 320000, m_h, num,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, lx + 380000, y + 80000, col_w - 550000, 260000,
                 title, size=11, bold=True, color=color)
        add_text(slide, lx + 380000, y + 320000, col_w - 550000, 200000,
                 when, size=8, color=MCK_GRAY)
        add_text(slide, lx + 380000, y + 510000, col_w - 550000, 240000,
                 desc, size=9, color=MCK_DARK)

    # Right: 迷ったら判断フロー
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, col_top, col_w, 280000, fill=MCK_BLUE)
    add_text(slide, rx + 100000, col_top + 50000, col_w - 200000, 220000,
             "迷ったら？  置き場所の判断  WHAT", size=10, bold=True, color=WHITE)

    destinations = [
        ("自分の作業スタイル・好み", "~/.claude/CLAUDE.md", "個人メモ", MCK_GRAY),
        ("このプロジェクト固有の話", "プロジェクト CLAUDE.md", "プロジェクト層", MCK_BLUE),
        ("複数プロジェクトに効く", "team-standards CLAUDE.md", "共通層", MCK_NAVY),
        ("2回以上やる手順(10行超)", "team-standards/skills/", "Skill", MCK_GREEN),
    ]
    d_top = col_top + 380000
    header_h = 280000
    row_h = 380000
    add_rect(slide, rx, d_top, col_w, header_h, fill=MCK_BG_LIGHT)
    add_text(slide, rx + 100000, d_top + 50000, 2400000, 200000,
             "性質", size=9, bold=True, color=MCK_GRAY)
    add_text(slide, rx + 2550000, d_top + 50000, col_w - 2650000, 200000,
             "行き先", size=9, bold=True, color=MCK_GRAY)
    for i, (prop, dest, label, color) in enumerate(destinations):
        y = d_top + header_h + i * row_h
        if i % 2 == 0:
            add_rect(slide, rx, y, col_w, row_h, fill=WHITE)
        else:
            add_rect(slide, rx, y, col_w, row_h, fill=MCK_BG_LIGHT)
        add_text(slide, rx + 100000, y + 60000, 2400000, 260000,
                 prop, size=9, color=MCK_DARK)
        add_text(slide, rx + 2550000, y + 40000, col_w - 2650000, 220000,
                 dest, size=9, bold=True, color=color)
        add_text(slide, rx + 2550000, y + 200000, col_w - 2650000, 180000,
                 label, size=7, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s06_skills_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Skills：呼び出し式の手順書 ― 2 回以上やる作業は Skill 化する",
        tracker="構成要素 2 / Skills")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # Left: What
    lx = MARGIN_L
    add_rect(slide, lx, top, col_w, 280000, fill=MCK_BLUE)
    add_text(slide, lx + 100000, top + 50000, col_w - 200000, 220000,
             "Skills とは  用途と使い分け", size=10, bold=True, color=WHITE)

    add_multiline(slide, lx + 100000, top + 380000, col_w - 200000, 3000000,
                  [("● 呼び出したときだけ読み込まれる", False, MCK_DARK, 10),
                   ("  （CLAUDE.md と違って常駐しない）", False, MCK_GRAY, 9),
                   ("", False),
                   ("● 長い手順・チェックリスト・", False, MCK_DARK, 10),
                   ("  ワークフローを格納する", False, MCK_DARK, 10),
                   ("", False),
                   ("● /implement-feature 等で呼び出し", False, MCK_DARK, 10),
                   ("", False),
                   ("── 使い分けの原則 ──", True, MCK_NAVY, 10),
                   ("CLAUDE.md: 常識・ルール（常駐）", False, MCK_DARK, 9),
                   ("Skills: 手順・チェックリスト（呼出）", False, MCK_DARK, 9),
                   ("", False),
                   ("── 作る目安 ──", True, MCK_NAVY, 10),
                   ("● 2 回以上やる作業 → Skill 化", False, MCK_DARK, 9),
                   ("● 10 行超の手順 → Skill 化", False, MCK_DARK, 9),
                   ("● CLAUDE.md が肥大化 → Skill 切出", False, MCK_DARK, 9)],
                  size=10, line_spacing=1.35)

    # Right: Must install
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "必須インストール  公式＋コミュニティ", size=10, bold=True, color=WHITE)

    items = [
        ("superpowers", "公式", "統合ワークフロー集（壁打ち/TDD/並列）", "★必須"),
        ("claude-md-management", "公式", "CLAUDE.md 品質監査", "★必須"),
        ("code-review", "公式", "5並列エージェント PR レビュー", "★必須"),
        ("feature-dev", "公式", "7 段階機能開発フロー", "推奨"),
        ("pr-review-toolkit", "公式", "6 種類の専門レビュー", "推奨"),
        ("commit-commands", "公式", "/commit, /commit-push-pr", "推奨"),
        ("obra/superpowers", "コミュニティ", "144k stars 決定版 Skill 集", "★必須"),
    ]
    i_top = top + 380000
    row_h = 380000
    for i, (name, kind, desc, tag) in enumerate(items):
        y = i_top + i * row_h
        if i % 2 == 0:
            add_rect(slide, rx, y, col_w, row_h, fill=MCK_BG_LIGHT)
        add_text(slide, rx + 100000, y + 30000, 2000000, 220000,
                 name, size=9, bold=True, color=MCK_NAVY)
        add_text(slide, rx + 2200000, y + 30000, 700000, 220000,
                 kind, size=7, color=MCK_GRAY)
        add_text(slide, rx + 100000, y + 240000, col_w - 700000, 140000,
                 desc, size=8, color=MCK_DARK)
        tag_color = MCK_RED if "必須" in tag else (MCK_BLUE if "推奨" in tag else MCK_GRAY)
        add_text(slide, rx + col_w - 550000, y + 30000, 450000, 220000,
                 tag, size=8, bold=True, color=tag_color, align=PP_ALIGN.RIGHT)

    add_footer(slide, page, total)
    return slide


def s07_subagents_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Subagents：特定分野に特化した専門秘書 ― VoltAgent から6つコピーするだけ",
        tracker="構成要素 3 / Subagents")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2

    # Left: What
    lx = MARGIN_L
    add_rect(slide, lx, top, col_w, 280000, fill=MCK_PURPLE)
    add_text(slide, lx + 100000, top + 50000, col_w - 200000, 220000,
             "Subagents とは  役割と使い方", size=10, bold=True, color=WHITE)

    add_multiline(slide, lx + 100000, top + 380000, col_w - 200000, 3000000,
                  [("● 特定役割に特化した『分身』", False, MCK_DARK, 10),
                   ("  （security 専門、設計 専門、等）", False, MCK_GRAY, 9),
                   ("", False),
                   ("● 独立したコンテキストウィンドウ", False, MCK_DARK, 10),
                   ("", False),
                   ("● ツール制限を設定できる", False, MCK_DARK, 10),
                   ("  （例: security は Write 禁止）", False, MCK_GRAY, 9),
                   ("", False),
                   ("● @agent-name で呼び出し or 自動委譲", False, MCK_DARK, 10),
                   ("", False),
                   ("● モデルも個別に選べる", False, MCK_DARK, 10),
                   ("  Haiku: 探索（-80%コスト）", False, MCK_GRAY, 9),
                   ("  Sonnet: 品質重視", False, MCK_GRAY, 9),
                   ("  Opus: 複雑推論", False, MCK_GRAY, 9),
                   ("", False),
                   ("── 並列実行 ──", True, MCK_NAVY, 10),
                   ("superpowers:dispatching-parallel-agents", False, MCK_DARK, 9),
                   ("で複数同時起動可能", False, MCK_DARK, 9)],
                  size=10, line_spacing=1.35)

    # Right: 6 subagents to copy
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "VoltAgent からコピーする6つ", size=10, bold=True, color=WHITE)

    subagents = [
        ("code-reviewer", "一般的コード品質"),
        ("architect-reviewer", "設計の妥当性・拡張性"),
        ("security-auditor", "脆弱性・認証漏れ"),
        ("test-coverage-checker", "テスト不足検出"),
        ("debugger", "根本原因分析"),
        ("dependency-manager", "依存関係監査"),
    ]
    s_top = top + 380000
    row_h = 440000
    for i, (name, desc) in enumerate(subagents):
        y = s_top + i * row_h
        add_rect(slide, rx, y, col_w, row_h,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_rect(slide, rx, y, 40000, row_h, fill=MCK_PURPLE)
        add_text(slide, rx + 120000, y + 60000, col_w - 200000, 260000,
                 name, size=10, bold=True, color=MCK_NAVY)
        add_text(slide, rx + 120000, y + 280000, col_w - 200000, 180000,
                 desc, size=8, color=MCK_DARK)

    add_footer(slide, page, total)
    return slide


def s08_subagents_tools(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Subagent の設定例：ツール制限とモデル選択で『安全な専門家』を作る",
        tracker="構成要素 3 / Subagents の設定")

    top = CONTENT_TOP + 100000

    # Top: explanation
    add_rect(slide, MARGIN_L, top, CONTENT_W, 400000, fill=MCK_BG_LIGHT)
    add_multiline(slide, MARGIN_L + 300000, top + 60000, CONTENT_W - 600000, 300000,
                  [("Subagent は frontmatter で『何ができるか・できないか』を明示する。", False, MCK_DARK, 10),
                   ("セキュリティ系は Write 禁止、テスト系は Bash 限定 ― など役割別に制限する。", False, MCK_NAVY, 10)],
                  size=10, line_spacing=1.3)

    # Code example for security-auditor
    ex_top = top + 500000
    ex_h = 2100000
    col_w = (CONTENT_W - 200000) // 2

    # Left: security-auditor
    add_rect(slide, MARGIN_L, ex_top, col_w, 250000, fill=MCK_RED)
    add_text(slide, MARGIN_L + 100000, ex_top + 40000, col_w - 200000, 200000,
             "例1: security-auditor.md", size=9, bold=True, color=WHITE)

    sec_code = [
        "---",
        "name: security-auditor",
        "description: 脆弱性・認証漏れ専門",
        "tools: Read, Grep, WebSearch",
        "disallowedTools: Write, Bash",
        "model: sonnet",
        "---",
        "",
        "あなたはセキュリティ監査の専門家。",
        "以下の観点でコードをレビュー：",
        "",
        "- 認証・認可の漏れ",
        "- SQL/Command Injection",
        "- 機密情報の露出",
        "- OWASP Top 10 への該当",
        "",
        "指摘には CVSS スコア目安を",
        "付与する。修正は提案のみ、",
        "実装はメインエージェントに任せる。",
    ]
    add_code_block(slide, MARGIN_L, ex_top + 280000, col_w, ex_h - 280000,
                   sec_code, size=8)

    # Right: code-reviewer
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, ex_top, col_w, 250000, fill=MCK_BLUE)
    add_text(slide, rx + 100000, ex_top + 40000, col_w - 200000, 200000,
             "例2: code-reviewer.md", size=9, bold=True, color=WHITE)

    rev_code = [
        "---",
        "name: code-reviewer",
        "description: 一般的なコード品質",
        "tools: Read, Grep, Glob",
        "model: haiku",
        "---",
        "",
        "あなたはコードレビュアー。",
        "CLAUDE.md に準拠しているかを",
        "最優先にチェック。",
        "",
        "- 命名規約",
        "- テストの有無",
        "- エラーハンドリング",
        "- 可読性",
        "",
        "Haiku モデルで軽量に回す。",
        "深い指摘は他の subagent に委譲。",
    ]
    add_code_block(slide, rx, ex_top + 280000, col_w, ex_h - 280000,
                   rev_code, size=8)

    # Bottom: Key principles
    pp_top = ex_top + ex_h + 100000
    add_rect(slide, MARGIN_L, pp_top, CONTENT_W, 700000,
             fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
    add_multiline(slide, MARGIN_L + 200000, pp_top + 60000, CONTENT_W - 400000, 600000,
                  [("● セキュリティ系: Write 禁止 — レビューだけして、修正は人間／メインに委ねる", False, MCK_DARK, 9),
                   ("● 一次レビュー: Haiku で軽量に — 80% コスト削減、速い", False, MCK_DARK, 9),
                   ("● 深いレビュー: Sonnet で品質重視", False, MCK_DARK, 9),
                   ("● 配置: team-standards/agents/ に commit — チーム全員に配布", False, MCK_DARK, 9)],
                  size=9, line_spacing=1.3)

    add_footer(slide, page, total)
    return slide


def s09_hooks_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Hooks：特定のイベントで自動実行 ― CLAUDE.md を『法律』に変える仕組み",
        tracker="構成要素 4 / Hooks")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # Left: What
    lx = MARGIN_L
    add_rect(slide, lx, top, col_w, 280000, fill=MCK_ORANGE)
    add_text(slide, lx + 100000, top + 50000, col_w - 200000, 220000,
             "Hooks とは", size=10, bold=True, color=WHITE)

    add_multiline(slide, lx + 100000, top + 380000, col_w - 200000, 3000000,
                  [("● 21 種類の Event で自動実行される", False, MCK_DARK, 10),
                   ("  スクリプト", False, MCK_DARK, 10),
                   ("", False),
                   ("● CLAUDE.md は『約束』、", False, MCK_DARK, 10),
                   ("  Hooks は『法律』", False, MCK_NAVY, 10),
                   ("  → 違反を機械的に検知・阻止", False, MCK_GRAY, 9),
                   ("", False),
                   ("── 4つのハンドラータイプ ──", True, MCK_NAVY, 10),
                   ("1. command  Shell スクリプト", False, MCK_DARK, 9),
                   ("2. http     外部 API へ POST", False, MCK_DARK, 9),
                   ("3. prompt   LLM で単一ターン評価", False, MCK_DARK, 9),
                   ("4. agent    フル subagent 統合", False, MCK_DARK, 9),
                   ("", False),
                   ("── ⚠️ 必ず Node.js で書く ──", True, MCK_RED, 10),
                   ("Bash だと Windows/WSL で壊れる", False, MCK_DARK, 9),
                   ("→ .mjs でクロスプラットフォーム", False, MCK_GRAY, 9)],
                  size=10, line_spacing=1.35)

    # Right: Main event types
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "主要 Event Type と用途", size=10, bold=True, color=WHITE)

    events = [
        ("SessionStart", "team-standards 最新化"),
        ("UserPromptSubmit", "CLAUDE.md 違反検出"),
        ("PreToolUse", "rm -rf 等の危険コマンド遮断"),
        ("PostToolUse", "auto-format, lint 実行"),
        ("Stop", "PR 作成前の最終チェック"),
        ("SessionEnd", "壁打ちログの自動保存"),
    ]
    e_top = top + 380000
    row_h = 450000
    for i, (name, purpose) in enumerate(events):
        y = e_top + i * row_h
        if i % 2 == 0:
            add_rect(slide, rx, y, col_w, row_h, fill=MCK_BG_LIGHT)
        add_text(slide, rx + 100000, y + 60000, 2400000, 260000,
                 name, size=10, bold=True, color=MCK_ORANGE)
        add_text(slide, rx + 100000, y + 290000, col_w - 200000, 180000,
                 purpose, size=9, color=MCK_DARK)

    add_footer(slide, page, total)
    return slide


def s10_hooks_day1(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "Day 1 に入れる3つの Hook：壁打ちログ保存・自動lint・危険コマンド遮断",
        tracker="構成要素 4 / Hooks の実装")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 300000) // 3

    hooks_list = [
        ("1. save-transcript", "SessionEnd",
         "壁打ちログを .claude/logs/\ntranscripts/ に自動保存",
         [
             '"SessionEnd": [{',
             '  "hooks": [{',
             '    "type": "command",',
             '    "command": "node',
             '      .claude/hooks/',
             '      save-transcript.mjs"',
             '  }]',
             '}]',
             "",
             "→ transcript_path を",
             "  読んで markdown で保存",
         ],
         MCK_NAVY),
        ("2. lint-on-edit", "PostToolUse",
         "Write/Edit 後に\nprettier/eslint 自動実行",
         [
             '"PostToolUse": [{',
             '  "matcher": "Write|Edit",',
             '  "hooks": [{',
             '    "type": "command",',
             '    "command":',
             '      "npx prettier',
             '       --write ...",',
             '    "async": true,',
             '    "timeout": 60',
             '  }]',
             '}]',
         ],
         MCK_BLUE),
        ("3. bash-guard", "PreToolUse",
         "rm -rf 等の\n破壊的コマンドをブロック",
         [
             '"PreToolUse": [{',
             '  "matcher": "Bash",',
             '  "if": "Bash(rm -rf *)",',
             '  "hooks": [{',
             '    "type": "command",',
             '    "command": "node',
             '      .claude/hooks/',
             '      bash-guard.mjs"',
             '  }]',
             '}]',
             "",
             "→ exit 2 でブロック",
         ],
         MCK_RED),
    ]

    for i, (title, event, desc, code, color) in enumerate(hooks_list):
        x = MARGIN_L + i * (col_w + 150000)
        # Header
        add_rect(slide, x, top, col_w, 320000, fill=color)
        add_text(slide, x + 100000, top + 50000, col_w - 200000, 250000,
                 title, size=11, bold=True, color=WHITE)
        # Event tag
        add_rect(slide, x, top + 330000, col_w, 200000, fill=MCK_BG_LIGHT)
        add_text(slide, x + 100000, top + 340000, col_w - 200000, 180000,
                 f"Event: {event}", size=8, bold=True, color=color)
        # Description
        add_multiline(slide, x + 100000, top + 560000, col_w - 200000, 400000,
                      desc.split("\n"), size=9, color=MCK_DARK, line_spacing=1.3)
        # Code block
        code_top = top + 1000000
        add_code_block(slide, x, code_top, col_w, 2500000, code, size=8)

    # Bottom note
    note_top = top + 3600000
    add_rect(slide, MARGIN_L, note_top, CONTENT_W, 300000, fill=MCK_BG_LIGHT)
    add_text(slide, MARGIN_L + 200000, note_top + 50000, CONTENT_W - 400000, 200000,
             "本体スクリプトは team-standards/hooks/ に置く。設定は .claude/settings.json で git commit",
             size=9, color=MCK_DARK, align=PP_ALIGN.CENTER)

    add_footer(slide, page, total)
    return slide


def s10b_settings_deny(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "settings.json の禁止事項：CLAUDE.md は約束、settings.json は鍵、Hooks は番人",
        tracker="構成要素 4.5 / permissions.deny")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2

    # Left: 3 layers + deny list
    lx = MARGIN_L
    add_rect(slide, lx, top, col_w, 280000, fill=MCK_RED)
    add_text(slide, lx + 100000, top + 50000, col_w - 200000, 220000,
             "3つの防御層と使い分け", size=10, bold=True, color=WHITE)

    layers = [
        ("CLAUDE.md", "約束", "秘書への「お願い」", "守らないことがある", MCK_NAVY),
        ("settings.json", "鍵", "ツール・コマンドの禁止", "物理的にできない", MCK_RED),
        ("Hooks", "番人", "イベント時の自動検証", "exit 2 でブロック", MCK_ORANGE),
    ]
    lay_top = top + 380000
    row_h = 480000
    for i, (name, metaphor, desc, strength, color) in enumerate(layers):
        y = lay_top + i * row_h
        add_rect(slide, lx, y, col_w, row_h - 30000,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_rect(slide, lx, y, 40000, row_h - 30000, fill=color)
        add_text(slide, lx + 100000, y + 30000, 1200000, 200000,
                 f"{name} = {metaphor}", size=10, bold=True, color=color)
        add_text(slide, lx + 100000, y + 230000, col_w - 200000, 200000,
                 f"{desc}（{strength}）", size=8, color=MCK_DARK)

    # Principle
    p_top = lay_top + 3 * row_h + 50000
    add_rect(slide, lx, p_top, col_w, 350000, fill=MCK_BG_LIGHT)
    add_multiline(slide, lx + 150000, p_top + 50000, col_w - 300000, 250000,
                  [("「絶対ダメ」→ settings.json", True, MCK_RED, 9),
                   ("「状況による」→ Hooks", True, MCK_ORANGE, 9),
                   ("「心がけ」→ CLAUDE.md", True, MCK_NAVY, 9)],
                  size=9, line_spacing=1.3)

    # Right: deny list
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "チームで入れるべき禁止事項", size=10, bold=True, color=WHITE)

    denies = [
        ("rm -rf *", "ファイル全消失"),
        ("git push --force", "リモート履歴上書き"),
        ("git reset --hard", "未コミット変更消失"),
        ("git checkout -- .", "作業中の変更消失"),
        ("git clean -fd", "未追跡ファイル消失"),
        ("DROP / DELETE FROM", "DB 破壊的操作"),
        ("chmod 777", "セキュリティ無効化"),
        ("curl | bash", "外部スクリプト盲目実行"),
    ]
    d_top = top + 380000
    row_h = 380000
    for i, (cmd, reason) in enumerate(denies):
        y = d_top + i * row_h
        if i % 2 == 0:
            add_rect(slide, rx, y, col_w, row_h, fill=MCK_BG_LIGHT)
        add_text(slide, rx + 100000, y + 50000, 2200000, 180000,
                 cmd, size=9, bold=True, color=MCK_RED,
                 font_name="Consolas")
        add_text(slide, rx + 100000, y + 210000, col_w - 200000, 160000,
                 reason, size=8, color=MCK_DARK)

    add_footer(slide, page, total)
    return slide


def s11_mcp_overview(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "MCP：外部ツール連携の仕組み ― 秘書が Gitea や DB に直接触れるようになる",
        tracker="構成要素 5 / MCP")

    top = CONTENT_TOP + 100000
    col_w = (CONTENT_W - 200000) // 2
    col_h = 3500000

    # Left: What
    lx = MARGIN_L
    add_rect(slide, lx, top, col_w, 280000, fill=MCK_GREEN)
    add_text(slide, lx + 100000, top + 50000, col_w - 200000, 220000,
             "MCP とは  Model Context Protocol", size=10, bold=True, color=WHITE)

    add_multiline(slide, lx + 100000, top + 380000, col_w - 200000, 3000000,
                  [("● 秘書（Claude Code）が", False, MCK_DARK, 10),
                   ("  外部ツールと話すための共通規格", False, MCK_DARK, 10),
                   ("", False),
                   ("● これで秘書が『手元のコード以外』", False, MCK_DARK, 10),
                   ("  を読み書きできるようになる", False, MCK_DARK, 10),
                   ("", False),
                   ("── 使い道の例 ──", True, MCK_NAVY, 10),
                   ("● Gitea の Issue / PR 操作", False, MCK_DARK, 9),
                   ("● PostgreSQL に直接クエリ", False, MCK_DARK, 9),
                   ("● ファイルシステム拡張", False, MCK_DARK, 9),
                   ("● Playwright でブラウザ自動化", False, MCK_DARK, 9),
                   ("● 社内 Wiki / Slack 連携", False, MCK_DARK, 9),
                   ("", False),
                   ("── 設定方法 ──", True, MCK_NAVY, 10),
                   (".mcp.json （プロジェクトスコープ）", False, MCK_DARK, 9),
                   ("または claude mcp add コマンド", False, MCK_DARK, 9),
                   ("", False),
                   ("トークンは .env で管理", False, MCK_RED, 9),
                   ("(絶対に commit しない)", False, MCK_RED, 9)],
                  size=10, line_spacing=1.35)

    # Right: What to use (closed-network safe)
    rx = MARGIN_L + col_w + 200000
    add_rect(slide, rx, top, col_w, 280000, fill=MCK_NAVY)
    add_text(slide, rx + 100000, top + 50000, col_w - 200000, 220000,
             "閉域環境 OK な MCP", size=10, bold=True, color=WHITE)

    mcps = [
        ("gitea/gitea-mcp", "Gitea 公式", "Issue/PR/コード操作", "★必須"),
        ("server-postgres", "公式", "読み取り専用 DB アクセス", "任意"),
        ("server-filesystem", "公式", "ファイル操作の拡張", "任意"),
        ("server-git", "公式", "Git 操作", "任意"),
        ("server-puppeteer", "公式", "Playwright/E2E テスト", "任意"),
        ("server-time", "公式", "時刻・タイムゾーン", "任意"),
    ]
    m_top = top + 380000
    row_h = 460000
    for i, (name, kind, desc, tag) in enumerate(mcps):
        y = m_top + i * row_h
        if i % 2 == 0:
            add_rect(slide, rx, y, col_w, row_h, fill=MCK_BG_LIGHT)
        add_text(slide, rx + 100000, y + 40000, 2200000, 220000,
                 name, size=9, bold=True, color=MCK_NAVY)
        add_text(slide, rx + 2350000, y + 40000, 900000, 220000,
                 kind, size=7, color=MCK_GRAY)
        add_text(slide, rx + 100000, y + 250000, col_w - 700000, 180000,
                 desc, size=8, color=MCK_DARK)
        tag_color = MCK_RED if "必須" in tag else MCK_GRAY
        add_text(slide, rx + col_w - 550000, y + 40000, 450000, 220000,
                 tag, size=8, bold=True, color=tag_color, align=PP_ALIGN.RIGHT)

    add_footer(slide, page, total)
    return slide


def s12_mcp_security(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "⚠️ MCP には重大な脆弱性が多数 ― 公式以外は導入前にチェック必須",
        tracker="構成要素 5 / MCP セキュリティ")

    top = CONTENT_TOP + 100000

    # Warning header
    add_rect(slide, MARGIN_L, top, CONTENT_W, 450000,
             fill=MCK_RED_LIGHT, border=MCK_RED, border_pt=1.5)
    add_text(slide, MARGIN_L + 200000, top + 60000, CONTENT_W - 400000, 260000,
             "⚠  2025-2026 年に発見された人気 MCP の CVE",
             size=12, bold=True, color=MCK_RED)
    add_text(slide, MARGIN_L + 200000, top + 280000, CONTENT_W - 400000, 160000,
             "500k+ ダウンロード済みの MCP でも深刻な脆弱性が判明している",
             size=9, color=MCK_DARK)

    # CVE table
    table_top = top + 550000
    row_h = 320000
    # Header
    add_rect(slide, MARGIN_L, table_top, CONTENT_W, row_h, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 150000, table_top + 80000, 2200000, 200000,
             "MCP 名", size=9, bold=True, color=WHITE)
    add_text(slide, MARGIN_L + 2400000, table_top + 80000, 2000000, 200000,
             "脆弱性タイプ", size=9, bold=True, color=WHITE)
    add_text(slide, MARGIN_L + 4500000, table_top + 80000, CONTENT_W - 4600000, 200000,
             "影響", size=9, bold=True, color=WHITE)

    cves = [
        ("mcp-remote (npm)", "CVE-2025-6514", "558k ダウンロード済み"),
        ("mcp-server-git", "CVE-2025-68145-6 RCE", "リモートコード実行"),
        ("SQLite 参考実装", "SQL Injection", "5000+ fork 被害"),
        ("Microsoft MarkItDown", "SSRF", "EC2 メタデータ露出"),
        ("gemini-mcp-tool", "Shell Injection CVSS 9.8", "未パッチ"),
    ]
    for i, (name, vuln, impact) in enumerate(cves):
        y = table_top + (i + 1) * row_h
        if i % 2 == 0:
            add_rect(slide, MARGIN_L, y, CONTENT_W, row_h, fill=WHITE)
        else:
            add_rect(slide, MARGIN_L, y, CONTENT_W, row_h, fill=MCK_BG_LIGHT)
        add_text(slide, MARGIN_L + 150000, y + 60000, 2200000, 200000,
                 name, size=9, bold=True, color=MCK_DARK)
        add_text(slide, MARGIN_L + 2400000, y + 60000, 2000000, 200000,
                 vuln, size=9, color=MCK_RED)
        add_text(slide, MARGIN_L + 4500000, y + 60000, CONTENT_W - 4600000, 200000,
                 impact, size=9, color=MCK_DARK)

    # Bottom: checklist preview
    check_top = table_top + 7 * row_h
    add_rect(slide, MARGIN_L, check_top, CONTENT_W, 550000,
             fill=MCK_BG_LIGHT, border=MCK_LIGHT_GRAY)
    add_text(slide, MARGIN_L + 200000, check_top + 50000, CONTENT_W - 400000, 240000,
             "導入前チェック（全項目必須）",
             size=11, bold=True, color=MCK_NAVY)
    add_text(slide, MARGIN_L + 200000, check_top + 280000, CONTENT_W - 400000, 220000,
             "□ 配布元  □ CVE確認  □ 権限範囲  □ トークン管理  □ 監査ログ  □ コードレビュー  □ テスト運用",
             size=9, color=MCK_DARK)

    add_footer(slide, page, total, source="OWASP MCP Top 10")
    return slide


def s13_team_standards_structure(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "team-standards リポジトリ構造：5要素すべてを1箇所に集約する",
        tracker="統合 / ディレクトリ構造")

    top = CONTENT_TOP + 50000

    # Left: tree
    left_w = (CONTENT_W - 200000) * 55 // 100
    add_rect(slide, MARGIN_L, top, left_w, 280000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L + 100000, top + 50000, left_w - 200000, 220000,
             "team-standards/  ディレクトリツリー", size=10, bold=True, color=WHITE)

    tree_lines = [
        "team-standards/",
        "│",
        "├── README.md",
        "├── CLAUDE.md                    ← 常に読まれる規約 (200行以内)",
        "│",
        "├── skills/                       ← 呼び出し式の手順書",
        "│   ├── implement-feature.md",
        "│   ├── post-to-gitea-issue.md   (自作)",
        "│   ├── review-my-changes.md     (拡張)",
        "│   └── brainstorm-log-to-4items.md (自作)",
        "│",
        "├── agents/                       ← 専門秘書 (VoltAgent コピー)",
        "│   ├── code-reviewer.md",
        "│   ├── architect-reviewer.md",
        "│   ├── security-auditor.md",
        "│   ├── test-coverage-checker.md",
        "│   ├── debugger.md",
        "│   └── dependency-manager.md",
        "│",
        "├── hooks/                        ← 自動化 (.mjs / Node.js)",
        "│   ├── save-transcript.mjs",
        "│   ├── lint-on-edit.mjs",
        "│   ├── bash-guard.mjs",
        "│   └── team-standards-pull.mjs",
        "│",
        "├── mcp/",
        "│   └── .mcp.json                 ← MCP 接続設定テンプレ",
        "│",
        "├── .claude/",
        "│   └── settings.json             ← Hook 有効化 (git commit)",
        "│",
        "├── .env.example                  ← トークンのテンプレ (空値)",
        "└── .gitignore                    ← .env, secrets 除外",
    ]
    add_code_block(slide, MARGIN_L, top + 380000, left_w, 3400000,
                   tree_lines, size=8)

    # Right: mapping table
    rx = MARGIN_L + left_w + 200000
    right_w = CONTENT_W - left_w - 200000
    add_rect(slide, rx, top, right_w, 280000, fill=MCK_BLUE)
    add_text(slide, rx + 100000, top + 50000, right_w - 200000, 220000,
             "構成要素との対応", size=10, bold=True, color=WHITE)

    mappings = [
        ("CLAUDE.md", "CLAUDE.md", MCK_NAVY),
        ("Skills", "skills/", MCK_BLUE),
        ("Subagents", "agents/", MCK_PURPLE),
        ("Hooks", "hooks/\n+ .claude/settings.json", MCK_ORANGE),
        ("MCP", "mcp/.mcp.json\n+ .env (個人)", MCK_GREEN),
    ]
    map_top = top + 380000
    row_h = 620000
    for i, (elem, where, color) in enumerate(mappings):
        y = map_top + i * row_h
        add_rect(slide, rx, y, right_w, row_h - 30000,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_rect(slide, rx, y, 50000, row_h - 30000, fill=color)
        add_text(slide, rx + 150000, y + 80000, right_w - 250000, 260000,
                 elem, size=11, bold=True, color=color)
        add_multiline(slide, rx + 150000, y + 300000, right_w - 250000, 250000,
                      where.split("\n"), size=9, color=MCK_DARK, line_spacing=1.3)

    add_footer(slide, page, total)
    return slide


def s14_install_order(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "導入順序：Day 1 で基盤 → Week 1 で強化 → Month 1 で仕上げ",
        tracker="導入計画")

    top = CONTENT_TOP + 100000
    phase_w = (CONTENT_W - 200000) // 3
    phase_h = 3500000

    phases = [
        ("PHASE 1", "Day 1", "基盤", "2時間以内", MCK_NAVY,
         [
             ("1.", "Gitea MCP", "公式を各自インストール"),
             ("", "", "gitea-mcp + Token"),
             ("", "", ""),
             ("2.", "公式プラグイン", "superpowers"),
             ("", "", "claude-md-management"),
             ("", "", ""),
             ("3.", "team-standards", "リポジトリ作成"),
             ("", "", "CLAUDE.md 初版"),
             ("", "", "(30 行の骨格)"),
             ("", "", ""),
             ("4.", "最小 Hooks", "save-transcript"),
             ("", "", "bash-guard"),
         ]),
        ("PHASE 2", "Week 1", "強化", "数日〜1週間", MCK_BLUE,
         [
             ("5.", "Subagents", "VoltAgent から6つ"),
             ("", "", "コピー & ツール制限"),
             ("", "", ""),
             ("6.", "追加プラグイン", "code-review"),
             ("", "", "feature-dev"),
             ("", "", ""),
             ("7.", "lint-on-edit", "Hook 追加"),
             ("", "", "PostToolUse で自動化"),
             ("", "", ""),
             ("8.", "/post-to-gitea", "slash command 自作"),
             ("", "", "2-4 時間"),
         ]),
        ("PHASE 3", "Month 1", "仕上げ", "余裕があれば", MCK_GREEN,
         [
             ("9.", "セルフレビュー", "Git Diff Reviewer"),
             ("", "", "+ CLAUDE.md 準拠"),
             ("", "", ""),
             ("10.", "月1刈り込み", "CLAUDE.md の"),
             ("", "", "Pruning Hook"),
             ("", "", ""),
             ("11.", "追加 MCP", "PostgreSQL など"),
             ("", "", "(セキュリティ確認後)"),
             ("", "", ""),
             ("12.", "CI連携", "Gitea Actions"),
             ("", "", "(任意・Phase 3末)"),
         ]),
    ]

    for i, (phase, when, label, duration, color, items) in enumerate(phases):
        x = MARGIN_L + i * (phase_w + 100000)
        # Header
        add_rect(slide, x, top, phase_w, 600000, fill=color)
        add_text(slide, x + 100000, top + 50000, phase_w - 200000, 200000,
                 phase, size=8, bold=True, color=RGBColor(200, 215, 230))
        add_text(slide, x + 100000, top + 230000, phase_w - 200000, 250000,
                 label, size=14, bold=True, color=WHITE)
        add_text(slide, x + 100000, top + 450000, phase_w - 200000, 180000,
                 f"{when}  ·  {duration}", size=8, color=RGBColor(200, 215, 230))
        # Body
        body_top = top + 600000
        body_h = phase_h - 600000
        add_rect(slide, x, body_top, phase_w, body_h,
                 fill=WHITE, border=color, border_pt=1)
        for j, (num, title, desc) in enumerate(items):
            y = body_top + 60000 + j * 220000
            if num:
                add_text(slide, x + 80000, y, 300000, 200000,
                         num, size=9, bold=True, color=color)
            if title:
                add_text(slide, x + 380000, y, phase_w - 460000, 200000,
                         title, size=9, bold=True, color=MCK_DARK)
            if desc and not title:
                add_text(slide, x + 380000, y, phase_w - 460000, 200000,
                         desc, size=8, color=MCK_GRAY)

    add_footer(slide, page, total)
    return slide


def s15_summary(prs, page, total):
    slide = new_slide(prs)
    add_action_title(
        slide,
        "まとめ：5要素を理解し、順番に導入すれば集合知化の基盤が完成する",
        tracker="まとめ")

    top = CONTENT_TOP + 100000

    # 5 key takeaways
    items = [
        ("CLAUDE.md", "チームの『書かれた常識』。200行以内で育て、迷ったら置き場所の判断フローで判定", MCK_NAVY),
        ("Skills", "呼び出し式の手順書。公式 superpowers と claude-md-management で80%カバー", MCK_BLUE),
        ("Subagents", "VoltAgent から6つコピーするだけ。security は Write 禁止などツール制限必須", MCK_PURPLE),
        ("Hooks", "CLAUDE.md を『法律』にする仕組み。Day 1 は save-transcript と bash-guard の2つで十分", MCK_ORANGE),
        ("MCP", "外部連携の要。公式優先、閉域 OK を選ぶ。導入前にセキュリティチェックリスト必須", MCK_GREEN),
    ]

    for i, (label, msg, color) in enumerate(items):
        y = top + i * 600000
        add_rect(slide, MARGIN_L, y, 1500000, 480000, fill=color)
        add_text(slide, MARGIN_L, y + 80000, 1500000, 350000,
                 label, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, MARGIN_L + 1550000, y, CONTENT_W - 1550000, 480000,
                 fill=WHITE, border=MCK_LIGHT_GRAY)
        add_text(slide, MARGIN_L + 1650000, y + 100000, CONTENT_W - 1800000, 300000,
                 msg, size=10, color=MCK_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom message
    bot_y = top + 5 * 600000 + 100000
    add_rect(slide, MARGIN_L, bot_y, CONTENT_W, 350000, fill=MCK_NAVY)
    add_text(slide, MARGIN_L, bot_y + 50000, CONTENT_W, 260000,
             "自作はほぼ不要。既存を組み合わせて、小さく始めて育てる",
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(slide, page, total)
    return slide


def s16_end(prs, page, total):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill=MCK_BLUE)

    add_text(slide, MARGIN_L, 1800000, CONTENT_W, 300000,
             "DX PORTAL TEAM  |  TECHNICAL REFERENCE", size=9, bold=True, color=MCK_BLUE,
             align=PP_ALIGN.CENTER)
    add_line(slide, (SLIDE_W - 1000000) // 2, 2200000,
             (SLIDE_W + 1000000) // 2, 2200000,
             color=MCK_NAVY, weight=2)
    add_text(slide, MARGIN_L, 2300000, CONTENT_W, 600000,
             "ありがとうございました", size=28, bold=True, color=MCK_NAVY,
             align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN_L, 2950000, CONTENT_W, 400000,
             "詳細は skills-inventory.md を参照",
             size=12, color=MCK_GRAY, align=PP_ALIGN.CENTER)
    return slide


# ==========================================
# MAIN
# ==========================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        s01_title,
        s02_executive_summary,
        s03_relationship,
        s04_claudemd_overview,
        s05_claudemd_growth,
        s06_skills_overview,
        s07_subagents_overview,
        s08_subagents_tools,
        s09_hooks_overview,
        s10_hooks_day1,
        s10b_settings_deny,
        s11_mcp_overview,
        s12_mcp_security,
        s13_team_standards_structure,
        s14_install_order,
        s15_summary,
        s16_end,
    ]

    total = len(builders)
    for i, build in enumerate(builders, 1):
        build(prs, i, total)
        print(f"  Built {i:2d}/{total}: {build.__name__}")

    out_path = Path(__file__).parent / "claude-code-components.pptx"
    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")
    print(f"Total: {total} slides")


if __name__ == "__main__":
    main()
