"""
Build team explanation slides for collective intelligence initiative.
Output: .tmp/team-collective-intelligence.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ==== Constants ====
SLIDE_W = 9144000   # 10 inches
SLIDE_H = 5143500   # 5.625 inches (16:9)

# Colors
NAVY = RGBColor(0, 51, 102)
NAVY_LIGHT = RGBColor(0, 82, 153)
DARK_GRAY = RGBColor(51, 51, 51)
MID_GRAY = RGBColor(102, 102, 102)
LIGHT_GRAY = RGBColor(230, 230, 230)
VERY_LIGHT_GRAY = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(204, 0, 0)
RED_LIGHT = RGBColor(253, 237, 237)
GREEN = RGBColor(0, 102, 51)
GREEN_LIGHT = RGBColor(232, 245, 233)
ORANGE = RGBColor(230, 115, 0)
ORANGE_LIGHT = RGBColor(255, 243, 224)
BLUE = RGBColor(25, 103, 210)
BLUE_LIGHT = RGBColor(232, 240, 254)
PURPLE = RGBColor(103, 58, 183)
BORDER_GRAY = RGBColor(210, 210, 210)

CHAT_USER_BG = RGBColor(225, 237, 255)
CHAT_USER_BORDER = RGBColor(180, 200, 230)
CHAT_AI_BG = RGBColor(245, 245, 245)
CHAT_AI_BORDER = RGBColor(210, 210, 210)

# Layout
TITLE_BAR_H = 620000
CONTENT_LEFT = 500000
CONTENT_WIDTH = SLIDE_W - 1000000
CONTENT_TOP = 780000  # After title bar + margin


# ==== Helpers ====

def add_rect(slide, left, top, width, height, fill=None, border=None, line_pt=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border is not None:
        shape.line.color.rgb = border
        if line_pt:
            shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill=None, border=None, line_pt=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.15
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border is not None:
        shape.line.color.rgb = border
        if line_pt:
            shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=11, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    if color is None:
        color = DARK_GRAY
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 50000
    tf.margin_right = 50000
    tf.margin_top = 30000
    tf.margin_bottom = 30000
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_multiline(slide, left, top, width, height, lines, size=11,
                  color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.15):
    """lines: list of strings OR tuples (text, bold, color_override, size_override)"""
    if color is None:
        color = DARK_GRAY
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 50000
    tf.margin_right = 50000
    tf.margin_top = 30000
    tf.margin_bottom = 30000
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(line, tuple):
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            line_color = line[2] if len(line) > 2 and line[2] is not None else color
            line_size = line[3] if len(line) > 3 and line[3] is not None else size
        else:
            text = line
            bold = False
            line_color = color
            line_size = size
        if text == "":
            run = p.add_run()
            run.text = " "
            run.font.size = Pt(line_size)
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(line_size)
        run.font.bold = bold
        run.font.color.rgb = line_color
    return tb


def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, TITLE_BAR_H, fill=NAVY)
    if subtitle:
        add_text(slide, 500000, 120000, 8000000, 240000,
                 subtitle, size=10, color=RGBColor(180, 200, 220))
        add_text(slide, 500000, 280000, 8000000, 320000,
                 title, size=19, bold=True, color=WHITE)
    else:
        add_text(slide, 500000, 150000, 8000000, 380000,
                 title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    # Ensure white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def add_chat_user(slide, top, text_lines, name="あなた"):
    """Add a user chat bubble (right-aligned). Returns bottom y."""
    width = 6800000
    left = SLIDE_W - width - 500000
    # Label
    add_text(slide, left, top, width, 240000, f"👤 {name}",
             size=9, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    bubble_top = top + 250000
    # Calculate height
    n_lines = len(text_lines)
    max_chars = max(len(l) for l in text_lines) if text_lines else 10
    wrapped_lines = sum(max(1, (len(l) // 42) + 1) for l in text_lines)
    bubble_h = 180000 + 240000 * wrapped_lines
    # Bubble
    bubble = add_round_rect(slide, left, bubble_top, width, bubble_h,
                             fill=CHAT_USER_BG, border=CHAT_USER_BORDER, line_pt=1)
    bubble.adjustments[0] = 0.15
    # Text inside
    add_multiline(slide, left + 150000, bubble_top + 80000,
                  width - 300000, bubble_h - 160000,
                  text_lines, size=10, color=DARK_GRAY)
    return bubble_top + bubble_h + 100000


def add_chat_ai(slide, top, text_lines, name="秘書 (Claude Code)"):
    """Add an AI chat bubble (left-aligned). Returns bottom y."""
    width = 6800000
    left = 500000
    add_text(slide, left, top, width, 240000, f"🤖 {name}",
             size=9, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    bubble_top = top + 250000
    n_lines = len(text_lines)
    wrapped_lines = sum(max(1, (len(l) // 42) + 1) for l in text_lines)
    bubble_h = 180000 + 240000 * wrapped_lines
    bubble = add_round_rect(slide, left, bubble_top, width, bubble_h,
                             fill=CHAT_AI_BG, border=CHAT_AI_BORDER, line_pt=1)
    bubble.adjustments[0] = 0.15
    add_multiline(slide, left + 150000, bubble_top + 80000,
                  width - 300000, bubble_h - 160000,
                  text_lines, size=10, color=DARK_GRAY)
    return bubble_top + bubble_h + 100000


def add_code_block(slide, left, top, width, height, code_lines, size=9, lang_color=None):
    """Add a code block with monospace appearance."""
    add_rect(slide, left, top, width, height, fill=VERY_LIGHT_GRAY, border=BORDER_GRAY, line_pt=0.75)
    lines_formatted = []
    for line in code_lines:
        if isinstance(line, tuple):
            lines_formatted.append(line)
        else:
            lines_formatted.append((line, False, DARK_GRAY, size))
    add_multiline(slide, left + 100000, top + 80000, width - 200000, height - 160000,
                  lines_formatted, size=size, line_spacing=1.1)


def add_bullet_list(slide, left, top, width, height, items, size=12,
                     color=None, bullet_color=None, line_spacing=1.3):
    """Add a bulleted list."""
    if color is None:
        color = DARK_GRAY
    if bullet_color is None:
        bullet_color = NAVY
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 50000
    tf.margin_top = 30000
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # Bullet
        run_b = p.add_run()
        run_b.text = "● "
        run_b.font.size = Pt(size)
        run_b.font.color.rgb = bullet_color
        run_b.font.bold = True
        # Text
        if isinstance(item, tuple):
            text, bold = item[0], item[1] if len(item) > 1 else False
        else:
            text, bold = item, False
        run_t = p.add_run()
        run_t.text = text
        run_t.font.size = Pt(size)
        run_t.font.color.rgb = color
        run_t.font.bold = bold
    return tb


def add_footer(slide, page_num, total):
    add_text(slide, SLIDE_W - 1500000, SLIDE_H - 350000,
             1400000, 250000,
             f"{page_num} / {total}", size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ==== Slide builders ====

def slide_title(prs, page, total):
    slide = new_slide(prs)
    # Colored band
    add_rect(slide, 0, 1500000, SLIDE_W, 2100000, fill=NAVY)
    # Main title
    add_text(slide, 500000, 1700000, SLIDE_W - 1000000, 600000,
             "チームで集合知化を目指す", size=32, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 2350000, SLIDE_W - 1000000, 400000,
             "Claude Code × Gitea で、個人の知恵をチームの力に変える", size=16,
             color=RGBColor(200, 215, 230), align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 2850000, SLIDE_W - 1000000, 350000,
             "~ 方向性の共有と、一緒に進めたいことの相談 ~", size=12,
             color=RGBColor(180, 195, 215), align=PP_ALIGN.CENTER)
    # Bottom note
    add_text(slide, 500000, SLIDE_H - 500000, SLIDE_W - 1000000, 300000,
             "DX ポータル開発チーム", size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
    return slide


def slide_purpose(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "この資料の目的")
    # 3 purpose cards
    card_w = 2700000
    card_h = 2600000
    gap = 200000
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) // 2
    top = 1000000

    cards = [
        ("1", "現状の共有", "今チームが抱える\n「もったいない」を\n一緒に確認する", NAVY),
        ("2", "方向性の提案", "Gitea + Claude Code で\n「集合知化」する\n具体的なやり方を示す", GREEN),
        ("3", "意見をもらう", "チームとして\n進めるかどうか、\n率直な意見がほしい", ORANGE),
    ]
    for i, (num, title, desc, color) in enumerate(cards):
        left = start_left + i * (card_w + gap)
        # Card body
        add_round_rect(slide, left, top, card_w, card_h,
                       fill=WHITE, border=color, line_pt=2)
        # Number circle
        circle_size = 600000
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         left + (card_w - circle_size) // 2,
                                         top + 200000,
                                         circle_size, circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(slide, left + (card_w - circle_size) // 2, top + 200000,
                 circle_size, circle_size,
                 num, size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(slide, left + 100000, top + 1000000, card_w - 200000, 400000,
                 title, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Desc
        lines = desc.split("\n")
        add_multiline(slide, left + 100000, top + 1450000, card_w - 200000, 1000000,
                      lines, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Bottom message
    add_text(slide, 500000, 4000000, SLIDE_W - 1000000, 400000,
             "結論を押しつけるものではありません。みんなで方向を決めたい。",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_problem_oneliner(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "今、困っていること")

    # Big message
    add_text(slide, 500000, 1100000, SLIDE_W - 1000000, 500000,
             "一言で言うと", size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 1550000, SLIDE_W - 1000000, 700000,
             "優秀な人材がいるのに、", size=24, bold=True, color=DARK_GRAY,
             align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 2150000, SLIDE_W - 1000000, 700000,
             "せっかくのノウハウが消えている。", size=24, bold=True, color=RED,
             align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 2900000, SLIDE_W - 1000000, 400000,
             "もったいない。", size=20, bold=True, color=DARK_GRAY,
             align=PP_ALIGN.CENTER)

    # Supporting box
    box_top = 3700000
    box_h = 900000
    add_round_rect(slide, 1500000, box_top, SLIDE_W - 3000000, box_h,
                   fill=RED_LIGHT, border=RED, line_pt=1.5)
    add_multiline(slide, 1700000, box_top + 100000, SLIDE_W - 3400000, box_h - 200000,
                  ["バイブコーディングはしている。各自がそれなりに使いこなしている。",
                   "でも、やり方がバラバラで、ノウハウが個人の頭の中だけ。",
                   "同じ失敗を繰り返し、同じことを何度も考え直している。"],
                  size=11, color=DARK_GRAY, line_spacing=1.4)
    add_footer(slide, page, total)
    return slide


def slide_example1_attrition(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "具体例 ①：Aさんのノウハウは、Aさんと一緒に消える", "起きていること")

    # Left: Aさんの頭の中
    left_x = 500000
    box_w = 3800000
    box_top = 900000
    box_h = 3700000

    add_round_rect(slide, left_x, box_top, box_w, box_h,
                   fill=VERY_LIGHT_GRAY, border=BORDER_GRAY, line_pt=1)
    add_text(slide, left_x + 100000, box_top + 100000, box_w - 200000, 400000,
             "Aさんの頭の中", size=14, bold=True, color=NAVY)

    items = [
        "DBタイムアウトはコネクションプール",
        "このAPIは夜間バッチで呼ばれる",
        "認証トークンの有効期限は30分",
        "このツールの特殊な使い方",
        "過去に試して失敗した設計3つ",
        "うまくいったプロンプトの書き方",
        "レビューで指摘された大事なこと",
    ]
    add_bullet_list(slide, left_x + 150000, box_top + 600000,
                    box_w - 300000, box_h - 700000,
                    items, size=11, bullet_color=NAVY)

    # Arrow
    arrow_x = left_x + box_w + 100000
    arrow_w = 900000
    add_text(slide, arrow_x, box_top + 1500000, arrow_w, 400000,
             "異動", size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x, box_top + 1900000, arrow_w, 500000,
             "→", size=32, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Right: ?マーク
    right_x = arrow_x + arrow_w + 100000
    add_round_rect(slide, right_x, box_top, box_w, box_h,
                   fill=RED_LIGHT, border=RED, line_pt=2)
    add_text(slide, right_x + 100000, box_top + 100000, box_w - 200000, 400000,
             "チームに残ったもの", size=14, bold=True, color=RED)
    add_text(slide, right_x, box_top + 1200000, box_w, 1200000,
             "???", size=80, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(slide, right_x + 100000, box_top + 2800000, box_w - 200000, 800000,
             "コードは残っている。\nでも「なぜこう作ったか」は\n誰も分からない。",
             size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_example2_repeat(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "具体例 ②：同じ「DBタイムアウト問題」で2人がハマる", "起きていること")

    # Top: Aさん
    add_rect(slide, 500000, 850000, SLIDE_W - 1000000, 1700000,
             fill=WHITE, border=BORDER_GRAY, line_pt=1)
    add_text(slide, 600000, 900000, 1000000, 300000,
             "2月", size=11, bold=True, color=MID_GRAY)
    add_text(slide, 1700000, 900000, 4000000, 300000,
             "Aさんが3時間ハマる", size=13, bold=True, color=RED)
    add_multiline(slide, 600000, 1200000, SLIDE_W - 1200000, 1300000,
                  [("① タイムアウト値を増やす → 効果なし", False),
                   ("② DBドライバを調査 → 関係なし", False),
                   ("③ コネクションプール設定で解決！", False, GREEN, 11),
                   ("", False),
                   ("→ 解決はしたが、どこにも記録しなかった", False, MID_GRAY, 10),
                   ("→ Aさんの頭の中にだけ残った", False, MID_GRAY, 10)],
                  size=11, color=DARK_GRAY, line_spacing=1.3)

    # Arrow down
    add_text(slide, 500000, 2650000, SLIDE_W - 1000000, 250000,
             "↓  3ヶ月後  ↓", size=11, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Bottom: Bさん
    add_rect(slide, 500000, 2950000, SLIDE_W - 1000000, 1700000,
             fill=RED_LIGHT, border=RED, line_pt=1.5)
    add_text(slide, 600000, 3000000, 1000000, 300000,
             "5月", size=11, bold=True, color=MID_GRAY)
    add_text(slide, 1700000, 3000000, 5000000, 300000,
             "Bさんが全く同じ問題で3時間ハマる", size=13, bold=True, color=RED)
    add_multiline(slide, 600000, 3300000, SLIDE_W - 1200000, 1300000,
                  [("① タイムアウト値を増やす → 効果なし", False),
                   ("② DBドライバを調査 → 関係なし", False),
                   ("③ コネクションプール設定で解決！", False, GREEN, 11),
                   ("", False),
                   ("→ Aさんも同じ事を3ヶ月前にやっていたと後で知る", False, RED, 10),
                   ("→ 「教えてくれたら3時間返してほしかった…」", False, RED, 10)],
                  size=11, color=DARK_GRAY, line_spacing=1.3)
    add_footer(slide, page, total)
    return slide


def slide_risk(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "このままだと、何が起きるか")

    items = [
        ("同じ失敗を何度も繰り返す", "Aさんの3時間、Bさんの3時間、Cさんの3時間…。年間で何十時間もの損失", RED),
        ("優秀な人が抜けた瞬間、知識が消える", "「あれどうなってるの？」が誰にも答えられない状態に", RED),
        ("新メンバーが立ち上がるのに時間がかかる", "「先輩に聞いて覚える」しかなく、先輩も新人も疲弊する", ORANGE),
        ("AIを使っているのに、チーム全体は賢くならない", "個々人は熟達するが、チームとしての成長がない", ORANGE),
        ("DXを推進する側なのに、自分たちの開発が旧態依然", "周囲から「DX語る資格あるの？」と見られるリスク", ORANGE),
    ]
    top = 950000
    row_h = 700000
    for i, (title, desc, color) in enumerate(items):
        y = top + i * row_h
        # Number
        num_size = 400000
        add_round_rect(slide, 500000, y + 100000, num_size, num_size,
                       fill=color, border=None)
        add_text(slide, 500000, y + 100000, num_size, num_size,
                 str(i + 1), size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(slide, 1000000, y + 80000, SLIDE_W - 1500000, 300000,
                 title, size=13, bold=True, color=color)
        # Desc
        add_text(slide, 1000000, y + 380000, SLIDE_W - 1500000, 300000,
                 desc, size=10, color=DARK_GRAY)
    add_footer(slide, page, total)
    return slide


def slide_goal(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "目指す姿：集合知化")

    # Definition
    add_text(slide, 500000, 900000, SLIDE_W - 1000000, 400000,
             "一言で言うと", size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 1300000, SLIDE_W - 1000000, 700000,
             "個人の知恵が、消えずに、チームの力になる", size=22, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)

    # Three dimensions
    top = 2300000
    card_w = 2700000
    gap = 150000
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) // 2

    dims = [
        ("残す", "時間軸で", "優秀な人の知恵が\n異動・退職後も残る", NAVY),
        ("揃える", "横軸で", "誰がやっても\n同じレベルで出せる", GREEN),
        ("掛け合わせる", "化学反応で", "Aさんの知恵 + Bさんの知恵\n= 1+1が3になる", ORANGE),
    ]
    for i, (word, sub, desc, color) in enumerate(dims):
        left = start_left + i * (card_w + gap)
        add_round_rect(slide, left, top, card_w, 2000000,
                       fill=WHITE, border=color, line_pt=2)
        add_text(slide, left + 100000, top + 150000, card_w - 200000, 350000,
                 sub, size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, left + 100000, top + 500000, card_w - 200000, 500000,
                 word, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_multiline(slide, left + 100000, top + 1100000, card_w - 200000, 800000,
                      desc.split("\n"), size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_vegapunk(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "たとえ話：ベガパンクの「サテライト」のような仕組み", "イメージをつかむ")

    # Left: explanation
    left_x = 500000
    add_text(slide, left_x, 1000000, 4000000, 400000,
             "ワンピースのベガパンク", size=14, bold=True, color=NAVY)
    add_multiline(slide, left_x, 1400000, 4000000, 2500000,
                  [("1人の天才が、", False),
                   ("6体のサテライト（分身）に分かれている。", False),
                   ("", False),
                   ("それぞれが独立して動くが、", False),
                   ("「パンクレコーズ」という", False),
                   ("共有の知識ベースでつながっている。", False),
                   ("", False),
                   ("だから、1体が学んだことは", False),
                   ("瞬時に全員に共有される。", False)],
                  size=11, color=DARK_GRAY, line_spacing=1.35)

    # Arrow
    add_text(slide, 4400000, 2300000, 400000, 400000,
             "≒", size=32, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Right: 私たちのイメージ
    right_x = 4900000
    add_text(slide, right_x, 1000000, 4000000, 400000,
             "私たちの目指す姿", size=14, bold=True, color=GREEN)
    add_multiline(slide, right_x, 1400000, 4000000, 2500000,
                  [("チーム全員が、", False),
                   ("自分専用の AI 秘書を持つ。", False),
                   ("", False),
                   ("各自が独立して開発するが、", False),
                   ("「team-standards」という", False),
                   ("共有の知識ベースでつながっている。", False),
                   ("", False),
                   ("だから、誰かが学んだことは", False),
                   ("全員の秘書の知識になる。", False)],
                  size=11, color=DARK_GRAY, line_spacing=1.35)

    # Bottom note
    add_rect(slide, 500000, 4200000, SLIDE_W - 1000000, 500000, fill=VERY_LIGHT_GRAY)
    add_text(slide, 500000, 4280000, SLIDE_W - 1000000, 350000,
             "ただしこれは Anthropic が提唱する運用 + コミュニティ実装の組み合わせで、すでに実現可能です",
             size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_why_gitea(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "なぜ GitHub ではなく Gitea か")

    # Three reasons
    top = 950000
    card_h = 3500000

    reasons = [
        ("すでに社内にある",
         "サーバー・アカウントが完備。\n追加投資ほぼゼロ。",
         NAVY),
        ("コードを社外に出せない",
         "セキュリティ・ポリシー上、\nGitHub Cloud は使えない。",
         RED),
        ("コスト", "GitHub Team は $4/user/月。\n5-10人で月数千円でも\n予算化が難しい。",
         ORANGE),
    ]

    card_w = (SLIDE_W - 1000000 - 400000) // 3
    for i, (title, desc, color) in enumerate(reasons):
        x = 500000 + i * (card_w + 200000)
        # Icon bar
        add_rect(slide, x, top, card_w, 100000, fill=color)
        # Card
        add_rect(slide, x, top + 100000, card_w, card_h - 100000,
                 fill=WHITE, border=BORDER_GRAY, line_pt=1)
        # Title
        add_text(slide, x + 150000, top + 300000, card_w - 300000, 400000,
                 title, size=14, bold=True, color=color)
        # Desc
        add_multiline(slide, x + 150000, top + 800000, card_w - 300000, 2000000,
                      desc.split("\n"), size=11, color=DARK_GRAY, line_spacing=1.35)

    add_text(slide, 500000, 4600000, SLIDE_W - 1000000, 350000,
             "→  すでにある Gitea を「ちゃんと使う」だけで、ほぼコストゼロで進められる",
             size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_direction_oneline(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "方向性を、一行で")

    # Big message
    top = 1500000
    add_rect(slide, 500000, top, SLIDE_W - 1000000, 1300000,
             fill=NAVY, border=None)
    add_text(slide, 500000, top + 200000, SLIDE_W - 1000000, 500000,
             "Gitea を", size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 500000, top + 500000, SLIDE_W - 1000000, 500000,
             "「完成コードの墓場」から", size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, 500000, top + 850000, SLIDE_W - 1000000, 500000,
             "「チームの脳と作業場所」に変える",
             size=22, bold=True, color=RGBColor(255, 220, 100),
             align=PP_ALIGN.CENTER)

    # Bottom: what this means
    bot = 3200000
    add_text(slide, 500000, bot, SLIDE_W - 1000000, 400000,
             "つまり", size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_multiline(slide, 500000, bot + 400000, SLIDE_W - 1000000, 1200000,
                  [("書く・考える・決める・学ぶ", True, NAVY, 16),
                   ("その全部がチームの資産として Gitea に溜まる仕組みを作る", False, DARK_GRAY, 13)],
                  size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.5)
    add_footer(slide, page, total)
    return slide


def slide_overview_diagram(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "全体像：仕事場と共同作業場の関係")

    # Left: 職人の作業台 (VS Code + Claude Code)
    left_x = 400000
    box_w = 3500000
    box_h = 2800000
    box_top = 1000000

    add_round_rect(slide, left_x, box_top, box_w, box_h,
                   fill=BLUE_LIGHT, border=BLUE, line_pt=2)
    add_text(slide, left_x + 150000, box_top + 150000, box_w - 300000, 400000,
             "🛠  職人の作業台", size=13, bold=True, color=BLUE)
    add_text(slide, left_x + 150000, box_top + 550000, box_w - 300000, 400000,
             "VS Code + Claude Code", size=14, bold=True, color=DARK_GRAY)
    add_text(slide, left_x + 150000, box_top + 950000, box_w - 300000, 300000,
             "各自のPC", size=10, color=MID_GRAY)

    add_multiline(slide, left_x + 150000, box_top + 1300000, box_w - 300000, 1500000,
                  [("● コードを書く", True, DARK_GRAY, 11),
                   ("● 秘書と壁打ち", True, DARK_GRAY, 11),
                   ("● テストを走らせる", True, DARK_GRAY, 11),
                   ("● デバッグする", True, DARK_GRAY, 11),
                   ("", False),
                   ("→ 今までどおり。むしろ強化", False, MID_GRAY, 10)],
                  size=11, line_spacing=1.4)

    # Middle: arrows
    arrow_x = left_x + box_w + 100000
    arrow_w = 1200000
    add_text(slide, arrow_x, box_top + 900000, arrow_w, 400000,
             "push / PR", size=10, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x, box_top + 1200000, arrow_w, 500000,
             "➡", size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x, box_top + 1750000, arrow_w, 500000,
             "⬅", size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x, box_top + 2250000, arrow_w, 400000,
             "pull / review", size=10, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Right: 工場の共同作業場 (Gitea)
    right_x = arrow_x + arrow_w + 100000
    add_round_rect(slide, right_x, box_top, box_w, box_h,
                   fill=GREEN_LIGHT, border=GREEN, line_pt=2)
    add_text(slide, right_x + 150000, box_top + 150000, box_w - 300000, 400000,
             "🏭  工場の共同作業場", size=13, bold=True, color=GREEN)
    add_text(slide, right_x + 150000, box_top + 550000, box_w - 300000, 400000,
             "Gitea", size=14, bold=True, color=DARK_GRAY)
    add_text(slide, right_x + 150000, box_top + 950000, box_w - 300000, 300000,
             "チームの共有の場所", size=10, color=MID_GRAY)

    add_multiline(slide, right_x + 150000, box_top + 1300000, box_w - 300000, 1500000,
                  [("● Issue（何を・なぜ）", True, DARK_GRAY, 11),
                   ("● 壁打ちログ（経緯）", True, DARK_GRAY, 11),
                   ("● PR / レビュー（判断）", True, DARK_GRAY, 11),
                   ("● CLAUDE.md（チームの脳）", True, DARK_GRAY, 11),
                   ("", False),
                   ("→ 知識が溜まる新しい役割", False, MID_GRAY, 10)],
                  size=11, line_spacing=1.4)

    # Bottom message
    add_text(slide, 500000, 4200000, SLIDE_W - 1000000, 400000,
             "コードを書く場所は変わらない。変えるのは「流れ」と「記憶」の置き場所だけ",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_day_in_life(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "1日の流れ：こう動きます")

    # Horizontal timeline
    top = 1000000
    steps = [
        ("朝", "Issue を\n選ぶ", "Gitea", NAVY),
        ("開始", "壁打ち\nする", "手元", BLUE),
        ("結論", "Issueに\n要約貼る", "Gitea", NAVY),
        ("実装", "コードを\n書く", "手元", BLUE),
        ("途中", "Draft\nPR", "Gitea", NAVY),
        ("確認", "セルフ\nレビュー", "手元", BLUE),
        ("完了", "人間\nレビュー", "Gitea", NAVY),
        ("後", "CLAUDE.md\n更新", "Gitea", GREEN),
    ]
    n = len(steps)
    card_w = (SLIDE_W - 1000000 - (n - 1) * 80000) // n
    card_h = 1500000

    for i, (time, action, place, color) in enumerate(steps):
        x = 500000 + i * (card_w + 80000)
        # Time label
        add_text(slide, x, top, card_w, 300000,
                 time, size=10, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
        # Card
        add_round_rect(slide, x, top + 350000, card_w, card_h,
                       fill=WHITE, border=color, line_pt=1.5)
        # Action
        action_lines = [(line, True, color, 11) for line in action.split("\n")]
        add_multiline(slide, x + 50000, top + 500000, card_w - 100000, 900000,
                      action_lines, size=11, color=color,
                      align=PP_ALIGN.CENTER)
        # Place
        add_text(slide, x, top + 1700000, card_w, 250000,
                 place, size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Arrows between
    for i in range(n - 1):
        x = 500000 + (i + 1) * card_w + i * 80000 + 20000
        add_text(slide, x, top + 900000, 40000, 300000,
                 "→", size=10, color=MID_GRAY)

    # Bottom key
    bot = 3500000
    add_rect(slide, 500000, bot, SLIDE_W - 1000000, 900000, fill=VERY_LIGHT_GRAY)
    add_multiline(slide, 600000, bot + 100000, SLIDE_W - 1200000, 800000,
                  [("● 青い枠＝手元（VS Code + Claude Code）", False, BLUE, 11),
                   ("● 紺色の枠＝Gitea（チームの場所）", False, NAVY, 11),
                   ("● 緑の枠＝学びの定着（CLAUDE.md 更新）", False, GREEN, 11)],
                  size=11, line_spacing=1.5)
    add_footer(slide, page, total)
    return slide


def slide_chat1_morning(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "実際のチャット ①：朝、Issue を選んで壁打ち開始", "具体的な会話イメージ")

    y = 850000
    y = add_chat_user(slide, y, [
        "おはよう。今日は Issue #42「申請フォームのリニューアル」を",
        "担当することになった。まず全体の設計から相談したい。"
    ])

    y = add_chat_ai(slide, y, [
        "おはようございます。Issue #42 を確認しました。",
        "team-standards の CLAUDE.md と合わせて読みました。",
        "",
        "現状の理解：",
        "・入力項目を15個→8個に削減",
        "・紙申請を廃止してオンライン化",
        "・期限は 4/15",
        "",
        "いくつか確認させてください：",
        "① 削減する7項目はもう決まっていますか？",
        "② 既存の申請データの移行は必要ですか？",
        "③ 承認フローは変わりませんか？"
    ])
    add_footer(slide, page, total)
    return slide


def slide_chat2_design(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "実際のチャット ②：設計の方針を決める", "具体的な会話イメージ")

    y = 850000
    y = add_chat_user(slide, y, [
        "① はまだ決まってない。現場の希望を聞きたい。",
        "② 移行は必要。データ構造も変える。",
        "③ 承認フローは変わらない。"
    ])

    y = add_chat_ai(slide, y, [
        "了解です。では進め方として2つの案があります：",
        "",
        "案A：既存画面を段階的に改修",
        "  → リスク小、でも中途半端な期間が発生",
        "",
        "案B：新画面を作って一気に切り替え",
        "  → リスク大、でもクリーンに移行できる",
        "",
        "過去の team-standards の失敗事例を見ると、",
        "『段階的改修は結局2回作業になる』という記録があります。",
        "",
        "今回の期間と規模を考えると、案Bを推奨します。",
        "いかがでしょうか？"
    ])
    add_footer(slide, page, total)
    return slide


def slide_chat3_issue_post(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "実際のチャット ③：壁打ち結論を Issue にポスト", "具体的な会話イメージ")

    y = 850000
    y = add_chat_user(slide, y, [
        "案B でいこう。今の議論を Issue #42 にコメント投稿しておいて。"
    ])

    y = add_chat_ai(slide, y, [
        "4項目テンプレにまとめて、Issue #42 にコメント投稿しました。",
        "",
        "── 投稿内容 ──",
        "【前提】",
        "申請フォームを全面リニューアル。4/15 期限。",
        "",
        "【検討した選択肢】",
        "案A：既存画面の段階的改修",
        "案B：新画面を作って切替",
        "",
        "【結論】",
        "案B を採用。新画面を作ってから一気に切替。",
        "",
        "【却下理由】",
        "案Aは過去事例で「2回作業になった」記録あり（CLAUDE.md）"
    ])
    add_footer(slide, page, total)
    return slide


def slide_chat4_implement(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "実際のチャット ④：実装の開始", "具体的な会話イメージ")

    y = 850000
    y = add_chat_user(slide, y, [
        "では実装に入ろう。まずフォームの画面から。"
    ])

    y = add_chat_ai(slide, y, [
        "了解です。team-standards の CLAUDE.md を確認しました：",
        "・Serendie デザインシステムを使う",
        "・TypeScript で any 禁止",
        "・テストファースト",
        "",
        "この順番で進めます：",
        "1. 型定義（入力項目8個の型）",
        "2. バリデーションのテストを先に書く",
        "3. フォームコンポーネント実装",
        "4. 確認画面実装",
        "",
        "まず型定義から始めます。",
        "src/types/application-form.ts を作成しました..."
    ])
    add_footer(slide, page, total)
    return slide


def slide_chat5_review(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "実際のチャット ⑤：push 前にセルフレビュー", "具体的な会話イメージ")

    y = 850000
    y = add_chat_user(slide, y, [
        "/review-my-changes"
    ])

    y = add_chat_ai(slide, y, [
        "変更を確認しました。",
        "",
        "✓ OK",
        "・Serendie コンポーネントを使っている",
        "・any 型は使われていない",
        "・テストが3件追加されている",
        "",
        "⚠ 気になる点",
        "・validateZipCode 関数にテストがない",
        "・console.log が残っている (src/form.tsx:42)",
        "",
        "❌ CLAUDE.md 違反",
        "・コミットメッセージが prefix 抜け",
        '  （"feat:" が必要）',
        "",
        "修正しましょうか？"
    ])
    add_footer(slide, page, total)
    return slide


def slide_chat6_claudemd_update(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "実際のチャット ⑥：学びを CLAUDE.md に追記提案", "具体的な会話イメージ")

    y = 850000
    y = add_chat_user(slide, y, [
        "今日は郵便番号のバリデーションでハマった。",
        "ハイフンあり/なしの両方対応が大事だと分かった。",
        "CLAUDE.md に追記したい。"
    ])

    y = add_chat_ai(slide, y, [
        "了解です。team-standards の CLAUDE.md への追記 PR を",
        "作成します。以下の内容でよいですか？",
        "",
        "── 追記案 ──",
        "## フォーム入力の過去の失敗あるある",
        "",
        "### 郵便番号のバリデーション",
        "- ハイフンあり（123-4567）/ なし（1234567）",
        "  両方の入力形式を受け付けること",
        "- 内部では常にハイフンなしに正規化してから",
        "  DB に保存すること",
        "",
        "この内容で PR を作成していいですか？"
    ])
    add_footer(slide, page, total)
    return slide


def slide_team_standards_overview(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "team-standards：チームの共有の脳")

    add_text(slide, 500000, 900000, SLIDE_W - 1000000, 400000,
             "Gitea に「team-standards」という1つのリポジトリを作り、チーム全員がローカルに clone して使う",
             size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Center: team-standards box
    box_w = 4200000
    box_h = 2700000
    box_x = (SLIDE_W - box_w) // 2
    box_top = 1500000

    add_round_rect(slide, box_x, box_top, box_w, box_h,
                   fill=NAVY, border=None)
    add_text(slide, box_x, box_top + 150000, box_w, 400000,
             "📦 team-standards", size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, box_x, box_top + 550000, box_w, 300000,
             "（1つのリポジトリ）", size=10, color=RGBColor(200, 215, 230),
             align=PP_ALIGN.CENTER)

    items = [
        ("CLAUDE.md", "常に読まれる基本ルール・規約"),
        ("skills/", "呼び出し式の手順書"),
        ("agents/", "専門の秘書（セキュリティ等）"),
        ("hooks/", "自動化スクリプト"),
        ("mcp/", "外部ツール連携（Gitea 等）"),
    ]
    inner_top = box_top + 900000
    for i, (name, desc) in enumerate(items):
        y = inner_top + i * 340000
        add_text(slide, box_x + 300000, y, 1500000, 280000,
                 name, size=11, bold=True, color=RGBColor(255, 220, 100))
        add_text(slide, box_x + 1700000, y, box_w - 2000000, 280000,
                 desc, size=10, color=WHITE)

    # Bottom
    add_text(slide, 500000, 4500000, SLIDE_W - 1000000, 400000,
             "1箇所で育てて、全プロジェクトに効かせる",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_hierarchy(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "CLAUDE.md の階層構造：共通 + プロジェクト固有")

    # Upper: team common
    upper_top = 900000
    upper_h = 1400000
    add_round_rect(slide, 1500000, upper_top, SLIDE_W - 3000000, upper_h,
                   fill=NAVY, border=None)
    add_text(slide, 1500000, upper_top + 150000, SLIDE_W - 3000000, 350000,
             "共通 CLAUDE.md（team-standards リポジトリ）",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_multiline(slide, 1700000, upper_top + 550000, SLIDE_W - 3400000, 800000,
                  [("● デザインシステムは Serendie を使う", False, WHITE, 11),
                   ("● TypeScript で any 禁止", False, WHITE, 11),
                   ("● 壁打ち→設計→実装の流れ", False, WHITE, 11),
                   ("● コミットメッセージのプレフィックス", False, WHITE, 11)],
                  size=11, line_spacing=1.3)

    # Arrow
    add_text(slide, 500000, 2450000, SLIDE_W - 1000000, 300000,
             "↑  @import で読み込む  ↑",
             size=11, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Lower: project specific (3 projects)
    lower_top = 2850000
    proj_w = 2500000
    proj_h = 1600000
    gap = 300000
    total_w = proj_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) // 2

    projects = [
        ("dx-portal", "React + Python\n申請系ツール", GREEN),
        ("dashboard", "Vue + Go\n分析ダッシュボード", BLUE),
        ("batch-tools", "Python のみ\n夜間バッチ", ORANGE),
    ]

    for i, (name, desc, color) in enumerate(projects):
        x = start_x + i * (proj_w + gap)
        add_round_rect(slide, x, lower_top, proj_w, proj_h,
                       fill=WHITE, border=color, line_pt=1.5)
        add_text(slide, x + 100000, lower_top + 100000, proj_w - 200000, 300000,
                 f"{name} の CLAUDE.md", size=11, bold=True, color=color)
        add_multiline(slide, x + 100000, lower_top + 450000, proj_w - 200000, 1000000,
                      desc.split("\n") + [
                          "",
                          "固有のルール",
                          "・ディレクトリ構造",
                          "・デプロイ手順",
                      ],
                      size=10, color=DARK_GRAY, line_spacing=1.3)
    add_footer(slide, page, total)
    return slide


def slide_claude_md_example(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "CLAUDE.md の具体例（共通ルール部分）", "team-standards リポジトリ")

    code_lines = [
        "# team-standards CLAUDE.md",
        "",
        "## 開発の流れ",
        "1. Issue を立てる（15分以上かかる作業）",
        "2. 壁打ち → 4項目テンプレで Issue にポスト",
        "3. Draft PR を出す（途中確認）",
        "4. セルフレビュー → 人間レビュー",
        "5. マージ後、学びがあれば CLAUDE.md 追記 PR",
        "",
        "## コーディング規約",
        "- TypeScript: any 禁止（unknown で代替）",
        "- デザイン: Serendie Design System を使う",
        "- コミット: 'feat:', 'fix:', 'refactor:' プレフィックス必須",
        "",
        "## 過去の失敗あるある",
        "- DB タイムアウト → コネクションプール設定を確認",
        "- 郵便番号 → ハイフンあり/なし両対応、DB は正規化",
        "- 認証エラー → トークン有効期限を先に確認",
        "",
        "## 壁打ちのポイント",
        "- まず「何を・なぜ」から話す（How より Why）",
        "- 過去の類似事例を秘書に確認してもらう",
    ]
    add_code_block(slide, 500000, 900000, SLIDE_W - 1000000, 3900000,
                   code_lines, size=10)
    add_footer(slide, page, total)
    return slide


def slide_skills_example(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "Skills の具体例：呼び出し式の手順書", "team-standards/skills/implement-feature.md")

    # Explanation
    add_text(slide, 500000, 850000, SLIDE_W - 1000000, 350000,
             "「2回以上やる作業」は Skill にして、/implement-feature と呼ぶだけで使える",
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    code_lines = [
        "---",
        "name: implement-feature",
        "description: 機能追加の標準手順",
        "---",
        "",
        "# 機能追加の標準手順",
        "",
        "1. 要件を Issue で確認する",
        "   - 未決事項があれば質問する",
        "",
        "2. 過去の類似事例を CLAUDE.md で探す",
        "   - あれば方針を参考にする",
        "",
        "3. 4項目テンプレで方針を整理し、Issue にポスト",
        "",
        "4. Draft PR を作成（空のコミットでもOK）",
        "",
        "5. テストを先に書く（RED）",
        "",
        "6. 実装する（GREEN）",
        "",
        "7. /review-my-changes でセルフレビュー",
        "",
        "8. Ready for review にして人間レビューを依頼",
    ]
    add_code_block(slide, 500000, 1300000, SLIDE_W - 1000000, 3500000,
                   code_lines, size=9)
    add_footer(slide, page, total)
    return slide


def slide_subagents_example(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "Subagents の具体例：専門秘書", "team-standards/agents/")

    # Description
    add_text(slide, 500000, 850000, SLIDE_W - 1000000, 400000,
             "特定の領域だけに特化した秘書を用意し、必要なときに呼び出す",
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Three agent cards
    cards = [
        ("🔒", "security-reviewer",
         "脆弱性・認証漏れを\n専門にチェック",
         "例: SQLインジェクション、\n権限チェック漏れ、\nトークン漏洩",
         RED),
        ("🏗", "architecture-reviewer",
         "設計の妥当性を\n専門にレビュー",
         "例: 責務の分離、\n適切な抽象化、\n将来の拡張性",
         NAVY),
        ("🧪", "test-coverage-checker",
         "テスト不足を\n専門にチェック",
         "例: エッジケース、\n異常系、\nカバレッジ",
         GREEN),
    ]
    top = 1350000
    card_w = (SLIDE_W - 1000000 - 300000) // 3
    card_h = 3100000

    for i, (icon, name, role, example, color) in enumerate(cards):
        x = 500000 + i * (card_w + 150000)
        add_round_rect(slide, x, top, card_w, card_h,
                       fill=WHITE, border=color, line_pt=2)
        add_text(slide, x, top + 150000, card_w, 600000,
                 icon, size=32, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, x + 100000, top + 800000, card_w - 200000, 400000,
                 name, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_multiline(slide, x + 100000, top + 1250000, card_w - 200000, 700000,
                      role.split("\n"), size=10, color=DARK_GRAY,
                      align=PP_ALIGN.CENTER, line_spacing=1.3)
        add_rect(slide, x + 150000, top + 2000000, card_w - 300000, 800000,
                 fill=VERY_LIGHT_GRAY)
        add_multiline(slide, x + 200000, top + 2050000, card_w - 400000, 700000,
                      example.split("\n"), size=9, color=MID_GRAY,
                      line_spacing=1.3)

    # Bottom
    add_text(slide, 500000, 4600000, SLIDE_W - 1000000, 350000,
             "「セキュリティレビューして」と言えば、その秘書が呼び出される",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_hooks_mcp(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "Hooks と MCP：自動化と外部連携")

    # Two columns
    col_w = (SLIDE_W - 1200000) // 2
    col_top = 900000
    col_h = 3800000

    # Left: Hooks
    add_round_rect(slide, 500000, col_top, col_w, col_h,
                   fill=WHITE, border=ORANGE, line_pt=2)
    add_rect(slide, 500000, col_top, col_w, 120000, fill=ORANGE)
    add_text(slide, 600000, col_top + 200000, col_w - 200000, 400000,
             "⚡ Hooks", size=16, bold=True, color=ORANGE)
    add_text(slide, 600000, col_top + 650000, col_w - 200000, 300000,
             "特定のイベントで自動実行されるスクリプト",
             size=10, color=MID_GRAY)

    hooks_items = [
        ("コミット前", "lint と型チェックを自動実行"),
        ("セッション終了時", "壁打ちログを自動保存"),
        ("PR作成時", "自動でセルフレビューを起動"),
    ]
    for i, (when, what) in enumerate(hooks_items):
        y = col_top + 1100000 + i * 700000
        add_text(slide, 600000, y, col_w - 200000, 300000,
                 f"● {when}", size=11, bold=True, color=ORANGE)
        add_text(slide, 700000, y + 300000, col_w - 300000, 350000,
                 what, size=10, color=DARK_GRAY)
    add_text(slide, 600000, col_top + 3400000, col_w - 200000, 300000,
             "スクリプトは team-standards/hooks/ で共有",
             size=9, color=MID_GRAY)

    # Right: MCP
    right_x = 500000 + col_w + 200000
    add_round_rect(slide, right_x, col_top, col_w, col_h,
                   fill=WHITE, border=PURPLE, line_pt=2)
    add_rect(slide, right_x, col_top, col_w, 120000, fill=PURPLE)
    add_text(slide, right_x + 100000, col_top + 200000, col_w - 200000, 400000,
             "🔌 MCP", size=16, bold=True, color=PURPLE)
    add_text(slide, right_x + 100000, col_top + 650000, col_w - 200000, 300000,
             "秘書が外部ツールと連携する仕組み",
             size=10, color=MID_GRAY)

    mcp_items = [
        ("Gitea MCP", "Issue/PR/コードを秘書が読み書き"),
        ("Slack MCP", "チャンネルへの投稿を秘書が代行"),
        ("DB MCP", "DB のスキーマを秘書が直接参照"),
    ]
    for i, (name, what) in enumerate(mcp_items):
        y = col_top + 1100000 + i * 700000
        add_text(slide, right_x + 100000, y, col_w - 200000, 300000,
                 f"● {name}", size=11, bold=True, color=PURPLE)
        add_text(slide, right_x + 200000, y + 300000, col_w - 300000, 350000,
                 what, size=10, color=DARK_GRAY)
    add_text(slide, right_x + 100000, col_top + 3400000, col_w - 200000, 300000,
             "Gitea 公式 MCP がすでに存在する",
             size=9, color=MID_GRAY)

    add_footer(slide, page, total)
    return slide


def slide_issue_rule(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "Issue の起票ルール：シンプルに1つだけ")

    # Main rule
    add_rect(slide, 500000, 900000, SLIDE_W - 1000000, 900000,
             fill=GREEN, border=None)
    add_text(slide, 500000, 1050000, SLIDE_W - 1000000, 300000,
             "ルール", size=11, color=RGBColor(200, 230, 210), align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 1320000, SLIDE_W - 1000000, 450000,
             "15分以上かかる作業は、Issue を立ててから始める",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Examples table
    top = 2050000
    add_text(slide, 500000, top, SLIDE_W - 1000000, 350000,
             "判断例", size=12, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    row_top = top + 450000
    row_h = 400000

    # Header
    add_rect(slide, 500000, row_top, 4500000, row_h, fill=VERY_LIGHT_GRAY)
    add_rect(slide, 5000000, row_top, 1800000, row_h, fill=VERY_LIGHT_GRAY)
    add_rect(slide, 6800000, row_top, 1800000, row_h, fill=VERY_LIGHT_GRAY)
    add_text(slide, 600000, row_top + 70000, 4400000, 300000,
             "作業", size=11, bold=True, color=DARK_GRAY)
    add_text(slide, 5000000, row_top + 70000, 1800000, 300000,
             "時間", size=11, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, 6800000, row_top + 70000, 1800000, 300000,
             "Issue？", size=11, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    rows = [
        ("新しい申請フォームを作る", "数日", "必要", GREEN),
        ("API の認証方式を変える", "半日〜1日", "必要", GREEN),
        ("バグ修正：検索が遅い", "半日", "必要", GREEN),
        ("ログ調査：特定ユーザーのエラー原因", "1時間", "必要", GREEN),
        ("タイポの修正", "1分", "不要（直接PR）", MID_GRAY),
        ("変数名のリネーム（1ファイル）", "5分", "不要（直接PR）", MID_GRAY),
    ]
    for i, (work, time, needed, color) in enumerate(rows):
        y = row_top + row_h + i * row_h
        # Alt color
        if i % 2 == 0:
            add_rect(slide, 500000, y, SLIDE_W - 1000000, row_h, fill=WHITE)
        else:
            add_rect(slide, 500000, y, SLIDE_W - 1000000, row_h, fill=VERY_LIGHT_GRAY)
        add_text(slide, 600000, y + 70000, 4400000, 300000, work, size=10, color=DARK_GRAY)
        add_text(slide, 5000000, y + 70000, 1800000, 300000, time, size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, 6800000, y + 70000, 1800000, 300000, needed, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_kabeuchi_template(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "壁打ちログの4項目テンプレ", "Issue コメントに貼る形式")

    # Description
    add_text(slide, 500000, 850000, SLIDE_W - 1000000, 350000,
             "壁打ちログを全部コピペしない。結論と経緯だけ4項目にまとめる",
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Template
    templ_top = 1350000
    templ_h = 3500000
    add_rect(slide, 500000, templ_top, SLIDE_W - 1000000, templ_h,
             fill=VERY_LIGHT_GRAY, border=BORDER_GRAY, line_pt=1)

    items = [
        ("【前提】", "どんな状況で、何を考え始めたか。背景情報。",
         "例：申請フォームを全面リニューアルする。4/15期限。入力項目を15→8に。"),
        ("【検討した選択肢】", "どんな案を出したか。1つだけじゃなく2-3個。",
         "例：A) 段階的改修  B) 新画面を作って切替  C) 現行維持で DB だけ変更"),
        ("【結論】", "どれを採用したか。一言で。",
         "例：案B を採用。新画面を作ってから一気に切替。"),
        ("【却下理由】", "他の案をなぜ選ばなかったか。後で同じ議論を避けるため。",
         "例：A は過去事例で「結局2回作業になった」記録あり（CLAUDE.md）。\n    C は DB 変更の影響範囲が読めずリスク大。"),
    ]
    y = templ_top + 150000
    for title, desc, example in items:
        add_text(slide, 600000, y, SLIDE_W - 1200000, 300000,
                 title, size=12, bold=True, color=NAVY)
        add_text(slide, 800000, y + 300000, SLIDE_W - 1400000, 280000,
                 desc, size=10, color=DARK_GRAY)
        lines = example.split("\n")
        add_multiline(slide, 800000, y + 580000, SLIDE_W - 1400000, 400000,
                      lines, size=9, color=MID_GRAY)
        y += 850000
    add_footer(slide, page, total)
    return slide


def slide_pr_review(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "PR レビュー：秘書と人間の2段階")

    add_text(slide, 500000, 850000, SLIDE_W - 1000000, 400000,
             "秘書が機械的なチェックを済ませた上で、人間は「判断」だけをする",
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Two stages
    stage_top = 1400000
    stage_w = (SLIDE_W - 1400000) // 2
    stage_h = 3100000

    # Stage 1: AI review
    add_round_rect(slide, 500000, stage_top, stage_w, stage_h,
                   fill=WHITE, border=GREEN, line_pt=2)
    add_rect(slide, 500000, stage_top, stage_w, 450000, fill=GREEN)
    add_text(slide, 500000, stage_top + 80000, stage_w, 350000,
             "ステップ1：秘書レビュー", size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    ai_items = [
        ("CLAUDE.md 準拠チェック", "any 使ってないか、命名規約、等"),
        ("テストの存在確認", "新しいコードにテストあるか"),
        ("過去の失敗事例に該当", "類似パターンの警告"),
        ("明らかなバグ・タイポ", "console.log 残りなど"),
        ("コミットメッセージ", "プレフィックス等"),
    ]
    for i, (name, desc) in enumerate(ai_items):
        y = stage_top + 650000 + i * 460000
        add_text(slide, 700000, y, stage_w - 300000, 280000,
                 f"✓ {name}", size=11, bold=True, color=GREEN)
        add_text(slide, 800000, y + 280000, stage_w - 400000, 280000,
                 desc, size=9, color=MID_GRAY)

    # Stage 2: Human review
    right_x = 500000 + stage_w + 400000
    add_round_rect(slide, right_x, stage_top, stage_w, stage_h,
                   fill=WHITE, border=NAVY, line_pt=2)
    add_rect(slide, right_x, stage_top, stage_w, 450000, fill=NAVY)
    add_text(slide, right_x, stage_top + 80000, stage_w, 350000,
             "ステップ2：人間レビュー", size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    human_items = [
        ("アプローチの妥当性", "方向性として合ってるか"),
        ("業務ロジックの意図", "このチームでしか分からない判断"),
        ("もっと良いやり方はないか", "経験に基づく提案"),
        ("影響範囲の見落としがないか", "隣の機能への波及"),
        ("命名・可読性", "3ヶ月後の自分が読めるか"),
    ]
    for i, (name, desc) in enumerate(human_items):
        y = stage_top + 650000 + i * 460000
        add_text(slide, right_x + 200000, y, stage_w - 300000, 280000,
                 f"✓ {name}", size=11, bold=True, color=NAVY)
        add_text(slide, right_x + 300000, y + 280000, stage_w - 400000, 280000,
                 desc, size=9, color=MID_GRAY)
    add_footer(slide, page, total)
    return slide


def slide_claude_md_growth(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "CLAUDE.md の育て方：当番制で継続する")

    # Main flow
    top = 950000
    add_text(slide, 500000, top, SLIDE_W - 1000000, 350000,
             "「気づいた人が書く」だけでは続かない。当番制で仕組み化する",
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Three mechanisms
    mech_top = 1450000
    mech_h = 1000000
    mechanisms = [
        ("1", "週1の振り返り定例（15分）",
         "今週のハマりどころ・うまくいったことを口頭で共有。\nその場で当番が CLAUDE.md 更新 PR を作る",
         GREEN),
        ("2", "その場の気づき PR",
         "「あ、秘書がまた同じミスした」と気づいた人がすぐ 3 行の PR。\n小さくてOK、チームがレビュー",
         NAVY),
        ("3", "月1の刈り込み（Pruning）",
         "古いルールを削除する当番を月替わりで。\n「なくても動く？」で判断。肥大化を防ぐ",
         ORANGE),
    ]
    for i, (num, title, desc, color) in enumerate(mechanisms):
        y = mech_top + i * (mech_h + 50000)
        add_rect(slide, 500000, y, SLIDE_W - 1000000, mech_h,
                 fill=WHITE, border=color, line_pt=1.5)
        # Number
        add_rect(slide, 500000, y, 500000, mech_h, fill=color)
        add_text(slide, 500000, y, 500000, mech_h, num,
                 size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Content
        add_text(slide, 700000, y + 150000, SLIDE_W - 1400000, 350000,
                 title, size=13, bold=True, color=color)
        add_multiline(slide, 700000, y + 500000, SLIDE_W - 1400000, 500000,
                      desc.split("\n"), size=10, color=DARK_GRAY, line_spacing=1.3)
    add_footer(slide, page, total)
    return slide


def slide_parallel_dev(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "2人以上の並列開発：まったく問題ありません")

    # Diagram
    top = 900000
    # A さん
    add_rect(slide, 500000, top, 3800000, 1400000,
             fill=BLUE_LIGHT, border=BLUE, line_pt=1.5)
    add_text(slide, 500000, top + 100000, 3800000, 350000,
             "👤 Aさん", size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_multiline(slide, 600000, top + 500000, 3600000, 800000,
                  [("ブランチ: feature/issue-42", False, DARK_GRAY, 10),
                   ("手元の Claude Code セッション", False, MID_GRAY, 9),
                   ("→ PR #50 を作成", False, DARK_GRAY, 10)],
                  size=10, line_spacing=1.4)

    # B さん
    add_rect(slide, 4800000, top, 3800000, 1400000,
             fill=GREEN_LIGHT, border=GREEN, line_pt=1.5)
    add_text(slide, 4800000, top + 100000, 3800000, 350000,
             "👤 Bさん", size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_multiline(slide, 4900000, top + 500000, 3600000, 800000,
                  [("ブランチ: feature/issue-43", False, DARK_GRAY, 10),
                   ("手元の Claude Code セッション", False, MID_GRAY, 9),
                   ("→ PR #51 を作成", False, DARK_GRAY, 10)],
                  size=10, line_spacing=1.4)

    # Arrow down to Gitea
    add_text(slide, 500000, 2450000, SLIDE_W - 1000000, 350000,
             "↓  それぞれ独立して動く  ↓", size=10, bold=True, color=MID_GRAY,
             align=PP_ALIGN.CENTER)

    # Gitea
    add_round_rect(slide, 500000, 2850000, SLIDE_W - 1000000, 1000000,
                   fill=NAVY, border=None)
    add_text(slide, 500000, 2900000, SLIDE_W - 1000000, 350000,
             "🏭 Gitea", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 3250000, SLIDE_W - 1000000, 300000,
             "PR #50 と PR #51 が並列で存在。お互い干渉しない。",
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 3550000, SLIDE_W - 1000000, 250000,
             "同じファイルに触った場合だけ、マージ時に Git が教えてくれる",
             size=10, color=RGBColor(200, 215, 230), align=PP_ALIGN.CENTER)

    # Bottom note
    add_text(slide, 500000, 4200000, SLIDE_W - 1000000, 400000,
             "これは Git の標準的な使い方。Claude Code が同時に何セッション動いても問題なし",
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_merge_conflict(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "マージ衝突：8割は Claude Code が解決してくれる")

    # 2 columns
    top = 900000

    # Left: できること
    add_round_rect(slide, 500000, top, 4000000, 3800000,
                   fill=GREEN_LIGHT, border=GREEN, line_pt=1.5)
    add_text(slide, 500000, top + 100000, 4000000, 400000,
             "✓ 秘書が得意なこと", size=14, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)

    good_items = [
        ("両方の意図を汲んだマージ", "Aさん=機能追加、Bさん=リファクタ。\n両方取り込める"),
        ("単純な行追加の衝突", "両方の追加を統合する"),
        ("指示しての解決", "『両方の変更の意図を汲んで統合して』と言うと\n説明付きで解決してくれる"),
    ]
    for i, (name, desc) in enumerate(good_items):
        y = top + 600000 + i * 1000000
        add_text(slide, 600000, y, 3800000, 300000,
                 f"● {name}", size=11, bold=True, color=GREEN)
        add_multiline(slide, 700000, y + 300000, 3700000, 700000,
                      desc.split("\n"), size=9, color=DARK_GRAY, line_spacing=1.3)

    # Right: できないこと
    add_round_rect(slide, 4800000, top, 4000000, 3800000,
                   fill=RED_LIGHT, border=RED, line_pt=1.5)
    add_text(slide, 4800000, top + 100000, 4000000, 400000,
             "⚠ 人間の判断が必要なこと", size=14, bold=True, color=RED,
             align=PP_ALIGN.CENTER)

    bad_items = [
        ("本質的に矛盾する設計判断", "『同期 vs 非同期』のような\n二択はチームで話し合う"),
        ("セマンティック衝突", "型が合わないが Git が検知できない\n→ 秘書にレビュー依頼で補完"),
        ("CLAUDE.md の衝突", "チームのルールが対立したら\nチームミーティングで合意を"),
    ]
    for i, (name, desc) in enumerate(bad_items):
        y = top + 600000 + i * 1000000
        add_text(slide, 4900000, y, 3800000, 300000,
                 f"● {name}", size=11, bold=True, color=RED)
        add_multiline(slide, 5000000, y + 300000, 3700000, 700000,
                      desc.split("\n"), size=9, color=DARK_GRAY, line_spacing=1.3)

    add_footer(slide, page, total)
    return slide


def slide_pilot_scope(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "進め方の提案：3名×3ヶ月のパイロット")

    # Main message
    add_rect(slide, 500000, 900000, SLIDE_W - 1000000, 700000,
             fill=ORANGE, border=None)
    add_text(slide, 500000, 1000000, SLIDE_W - 1000000, 500000,
             "いきなり全員ではなく、3名で3ヶ月試してから広げます",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Why
    top = 1800000
    add_text(slide, 500000, top, SLIDE_W - 1000000, 350000,
             "なぜ全員一斉じゃないのか", size=13, bold=True, color=DARK_GRAY,
             align=PP_ALIGN.CENTER)

    reasons = [
        ("失敗コストを小さく", "3名なら影響範囲 1/3。ダメならすぐ戻せる"),
        ("ワークフローを磨く期間", "最初の運用で必ず不具合が出る。小さく始めて調整"),
        ("成功事例を作ってから展開", "3ヶ月後に「こう良かった」と残り 7名に見せられる"),
        ("上司を説得しやすい", "『全面移行』より『試験運用』の方が承認が取りやすい"),
    ]
    for i, (title, desc) in enumerate(reasons):
        y = top + 450000 + i * 600000
        add_text(slide, 600000, y, SLIDE_W - 1200000, 300000,
                 f"{i + 1}. {title}", size=12, bold=True, color=NAVY)
        add_text(slide, 900000, y + 280000, SLIDE_W - 1500000, 280000,
                 desc, size=10, color=DARK_GRAY)
    add_footer(slide, page, total)
    return slide


def slide_pilot_timeline(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "パイロット 3ヶ月の流れ")

    # Timeline
    top = 950000
    phases = [
        ("1ヶ月目", "準備と慣らし",
         ["team-standards リポジトリ作成",
          "CLAUDE.md 初版を3名で書く",
          "1つの小さな Issue で全フロー試す",
          "壁打ちログ4項目テンプレで運用"]),
        ("2ヶ月目", "本格運用",
         ["3名の実作業を完全にこのフローに",
          "Skills・Subagents を追加",
          "週1振り返りで CLAUDE.md 更新",
          "当番制を徹底"]),
        ("3ヶ月目", "チェックと調整",
         ["撤退条件をチェック",
          "推進役2週間休暇テスト",
          "成果と課題を整理",
          "チームへの報告・展開判断"]),
    ]

    phase_w = (SLIDE_W - 1200000) // 3
    for i, (month, title, items) in enumerate(phases):
        x = 500000 + i * (phase_w + 100000)
        color = [GREEN, NAVY, ORANGE][i]
        # Header
        add_rect(slide, x, top, phase_w, 600000, fill=color)
        add_text(slide, x, top + 50000, phase_w, 300000,
                 month, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x, top + 300000, phase_w, 300000,
                 title, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        # Body
        add_rect(slide, x, top + 600000, phase_w, 2800000,
                 fill=WHITE, border=color, line_pt=1.5)
        items_tuples = [("● " + item, False) for item in items]
        add_multiline(slide, x + 150000, top + 700000, phase_w - 300000, 2700000,
                      items_tuples, size=10, color=DARK_GRAY, line_spacing=1.5)
    add_footer(slide, page, total)
    return slide


def slide_checkpoints(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "チェックポイントと撤退条件")

    add_text(slide, 500000, 900000, SLIDE_W - 1000000, 400000,
             "「成果が出るまで続ける」は「永遠に続けるか、感情的にやめるか」になる。先に判断基準を決める",
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Two checkpoints
    top = 1500000
    cp_w = (SLIDE_W - 1200000) // 2
    cp_h = 3200000

    # 3ヶ月
    add_round_rect(slide, 500000, top, cp_w, cp_h,
                   fill=WHITE, border=NAVY, line_pt=2)
    add_rect(slide, 500000, top, cp_w, 450000, fill=NAVY)
    add_text(slide, 500000, top + 80000, cp_w, 350000,
             "3ヶ月目チェック", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    cp3 = [
        "推進役以外の CLAUDE.md 更新 PR が 5件以上",
        "壁打ちログ4項目が週2件以上蓄積",
        "team-standards の CODEOWNERS 2人とも機能",
        "参加3名がワークフローに慣れた実感",
    ]
    for i, item in enumerate(cp3):
        y = top + 600000 + i * 550000
        add_text(slide, 600000, y, 300000, 300000,
                 "□", size=14, bold=True, color=NAVY)
        add_text(slide, 900000, y + 30000, cp_w - 500000, 500000,
                 item, size=10, color=DARK_GRAY)

    # 6ヶ月
    right_x = 500000 + cp_w + 200000
    add_round_rect(slide, right_x, top, cp_w, cp_h,
                   fill=WHITE, border=ORANGE, line_pt=2)
    add_rect(slide, right_x, top, cp_w, 450000, fill=ORANGE)
    add_text(slide, right_x, top + 80000, cp_w, 350000,
             "6ヶ月目チェック", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    cp6 = [
        "推進役2週間休暇テストを実施、回るか",
        "集合知化の実感（アンケート）",
        "定量指標（PR数・レビュー時間等）",
        "残り7名への展開判断",
    ]
    for i, item in enumerate(cp6):
        y = top + 600000 + i * 550000
        add_text(slide, right_x + 100000, y, 300000, 300000,
                 "□", size=14, bold=True, color=ORANGE)
        add_text(slide, right_x + 400000, y + 30000, cp_w - 500000, 500000,
                 item, size=10, color=DARK_GRAY)

    # Bottom
    bot_top = 4800000
    add_rect(slide, 500000, bot_top, SLIDE_W - 1000000, 250000, fill=RED)
    add_text(slide, 500000, bot_top + 20000, SLIDE_W - 1000000, 250000,
             "半数以上の指標が未達 → 設計やり直し or 中止を検討",
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_anti_dependency(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "推進役依存の罠への対策", "一番の失敗パターンを避ける")

    # The trap
    add_rect(slide, 500000, 900000, SLIDE_W - 1000000, 900000,
             fill=RED_LIGHT, border=RED, line_pt=1.5)
    add_text(slide, 600000, 950000, SLIDE_W - 1200000, 350000,
             "⚠ 最大の失敗シナリオ", size=12, bold=True, color=RED)
    add_multiline(slide, 600000, 1300000, SLIDE_W - 1200000, 500000,
                  [("推進役Aさんだけが CLAUDE.md を更新し、", False, DARK_GRAY, 11),
                   ("他メンバーは「Aさんが作った秘書に聞く消費者」になる。", False, DARK_GRAY, 11),
                   ("Aさんが異動した瞬間、全員立ち尽くす。", True, RED, 11)],
                  size=11, line_spacing=1.4)

    # Countermeasures
    top = 2100000
    add_text(slide, 500000, top, SLIDE_W - 1000000, 350000,
             "対策：最初から属人化しない設計に", size=13, bold=True, color=DARK_GRAY,
             align=PP_ALIGN.CENTER)

    measures = [
        ("輪番制の徹底", "CLAUDE.md 更新、振り返り議事、レビューの当番を週替わり"),
        ("CODEOWNERS 2人体制", "team-standards のレビュアーを必ず2人（休暇対策）"),
        ("属人化 KPI", "「推進役以外の人が触った PR の比率」を定期的に確認"),
        ("2週間休暇テスト", "3ヶ月目に推進役を意図的に2週間休ませて、回るか検証"),
    ]
    for i, (title, desc) in enumerate(measures):
        y = top + 500000 + i * 600000
        add_text(slide, 600000, y, 300000, 300000,
                 f"{i + 1}.", size=12, bold=True, color=NAVY)
        add_text(slide, 900000, y, SLIDE_W - 1500000, 300000,
                 title, size=12, bold=True, color=DARK_GRAY)
        add_text(slide, 900000, y + 300000, SLIDE_W - 1500000, 280000,
                 desc, size=10, color=MID_GRAY)
    add_footer(slide, page, total)
    return slide


def slide_first_week(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "もし合意できたら、最初の1週間でやること")

    top = 900000
    items = [
        ("Day 1-2", "準備", [
            "team-standards リポジトリを Gitea に作成",
            "最初の CLAUDE.md を 3名で一緒に書く（200行以内）",
            "各自が ~/dev/team-standards に clone",
        ], NAVY),
        ("Day 3-4", "試運転", [
            "小さな Issue を1つ選んで全フローを試す",
            "Issue → 壁打ち → 4項目テンプレ → Draft PR → レビュー → マージ",
            "つまずいた点を記録する",
        ], BLUE),
        ("Day 5", "振り返り", [
            "3名で振り返り15分",
            "良かった点・悪かった点を書き出す",
            "CLAUDE.md に反映する（早速当番制で）",
        ], GREEN),
    ]

    for i, (days, phase, tasks, color) in enumerate(items):
        y = top + i * 1300000
        # Days label
        add_rect(slide, 500000, y, 1200000, 1100000, fill=color)
        add_text(slide, 500000, y + 200000, 1200000, 300000,
                 days, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, 500000, y + 550000, 1200000, 300000,
                 phase, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Tasks
        add_rect(slide, 1700000, y, SLIDE_W - 2200000, 1100000,
                 fill=WHITE, border=color, line_pt=1.5)
        items_tuples = [("● " + t, False) for t in tasks]
        add_multiline(slide, 1850000, y + 150000, SLIDE_W - 2500000, 900000,
                      items_tuples, size=11, color=DARK_GRAY, line_spacing=1.4)
    add_footer(slide, page, total)
    return slide


def slide_questions_to_team(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "みんなに聞きたいこと")

    add_text(slide, 500000, 900000, SLIDE_W - 1000000, 400000,
             "この提案を進めるかどうかの前に、率直な意見を聞きたい",
             size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

    questions = [
        ("この方向性、ピンと来る？", "「集合知化」という旗に心が動く？ 他に優先したいことがある？"),
        ("最初の3名、誰がやる？", "興味がある人、手を挙げてくれる人。無理強いはしたくない"),
        ("不安・懸念は？", "「これは続かないと思う」「ここが面倒」など、率直に聞きたい"),
        ("もっと良い方法ある？", "Gitea 中心以外のアプローチ、気になる案があれば"),
        ("上司への相談、どう進める？", "タイミング、伝え方、必要な準備。協力してほしい"),
    ]
    top = 1500000
    for i, (q, desc) in enumerate(questions):
        y = top + i * 600000
        add_text(slide, 600000, y, 500000, 300000,
                 "?", size=22, bold=True, color=NAVY)
        add_text(slide, 1000000, y + 40000, SLIDE_W - 1500000, 300000,
                 q, size=13, bold=True, color=DARK_GRAY)
        add_text(slide, 1000000, y + 330000, SLIDE_W - 1500000, 280000,
                 desc, size=10, color=MID_GRAY)
    add_footer(slide, page, total)
    return slide


def slide_summary(prs, page, total):
    slide = new_slide(prs)
    add_title_bar(slide, "まとめ")

    # Main
    add_text(slide, 500000, 900000, SLIDE_W - 1000000, 400000,
             "今日話したこと", size=13, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Key messages
    top = 1350000
    msgs = [
        ("旗", "集合知化 ― 個人の知恵がチームの力になる", NAVY),
        ("手段", "Gitea を『作業場所と脳』に変える（Claude Code と組み合わせて）", GREEN),
        ("場所", "team-standards リポジトリに CLAUDE.md・Skills・Agents を集約", BLUE),
        ("進め方", "3名×3ヶ月のパイロット。撤退条件と属人化対策を最初から", ORANGE),
        ("次の一歩", "みんなの意見を聞く。合意できたら来週から試運転", RED),
    ]
    for i, (label, msg, color) in enumerate(msgs):
        y = top + i * 620000
        # Label
        add_rect(slide, 500000, y, 1200000, 450000, fill=color)
        add_text(slide, 500000, y + 50000, 1200000, 350000,
                 label, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Message
        add_text(slide, 1800000, y + 80000, SLIDE_W - 2300000, 400000,
                 msg, size=12, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom
    add_text(slide, 500000, 4550000, SLIDE_W - 1000000, 400000,
             "完璧な計画より、やってみて直す方が早い。一緒に進めませんか？",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide


def slide_end(prs, page, total):
    slide = new_slide(prs)
    # Simple ending
    add_rect(slide, 0, 1800000, SLIDE_W, 1500000, fill=NAVY)
    add_text(slide, 500000, 2000000, SLIDE_W - 1000000, 500000,
             "ありがとうございました", size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, 500000, 2700000, SLIDE_W - 1000000, 400000,
             "意見・質問、お待ちしています", size=14,
             color=RGBColor(200, 215, 230), align=PP_ALIGN.CENTER)
    return slide


# ==== Main ====

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Define slide order
    builders = [
        slide_title,
        slide_purpose,
        slide_problem_oneliner,
        slide_example1_attrition,
        slide_example2_repeat,
        slide_risk,
        slide_goal,
        slide_vegapunk,
        slide_why_gitea,
        slide_direction_oneline,
        slide_overview_diagram,
        slide_day_in_life,
        slide_chat1_morning,
        slide_chat2_design,
        slide_chat3_issue_post,
        slide_chat4_implement,
        slide_chat5_review,
        slide_chat6_claudemd_update,
        slide_team_standards_overview,
        slide_hierarchy,
        slide_claude_md_example,
        slide_skills_example,
        slide_subagents_example,
        slide_hooks_mcp,
        slide_issue_rule,
        slide_kabeuchi_template,
        slide_pr_review,
        slide_claude_md_growth,
        slide_parallel_dev,
        slide_merge_conflict,
        slide_pilot_scope,
        slide_pilot_timeline,
        slide_checkpoints,
        slide_anti_dependency,
        slide_first_week,
        slide_questions_to_team,
        slide_summary,
        slide_end,
    ]

    total = len(builders)
    for i, build in enumerate(builders, 1):
        build(prs, i, total)
        print(f"  Built slide {i}/{total}: {build.__name__}")

    out_path = Path(__file__).parent / "team-collective-intelligence.pptx"
    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")
    print(f"Total slides: {total}")


if __name__ == "__main__":
    main()
